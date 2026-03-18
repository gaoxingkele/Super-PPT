#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
基于分析结果创建模板 PPTX
"""
import sys
sys.path.insert(0, '.')

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pathlib import Path
import json


# 基于分析结果整理的 4 个模板配置
TEMPLATES = {
    'yili_power': {
        'name': '亿力科技-电力生产安全',
        'source': '国网信通亿力科技有限责任公司-多模态项目汇报',
        'color_scheme': {
            'primary': '#024177',      # 深蓝
            'secondary': '#005D7F',    # 青色
            'accent': '#E8612D',       # 默认橙色强调
            'background': '#FFFFFF',
            'text': '#333333',
            'text_light': '#666666'
        },
        'fonts': {
            'title': '微软雅黑',
            'body': '微软雅黑'
        },
        'style': 'corporate-blue',  # 商务深蓝风格
        'slide_count': 37
    },
    'xmu_graph': {
        'name': '厦门大学-图理论研究',
        'source': '基于图理论的能源企业数字化智能化建设',
        'color_scheme': {
            'primary': '#178F95',      # 青绿
            'secondary': '#4472C4',    # 蓝色
            'accent': '#E8612D',
            'background': '#FFFFFF',
            'text': '#333333',
            'text_light': '#666666'
        },
        'fonts': {
            'title': '微软雅黑',
            'body': '微软雅黑'
        },
        'style': 'academic-teal',  # 学术青绿风格
        'slide_count': 44
    },
    'epri_nature': {
        'name': '电科院-自然基金',
        'source': '电网-重点-中国电力科学研究院',
        'color_scheme': {
            'primary': '#0070C0',      # 标准蓝
            'secondary': '#4C9857',    # 绿色
            'accent': '#E8612D',
            'background': '#FFFFFF',
            'text': '#333333',
            'text_light': '#666666'
        },
        'fonts': {
            'title': '微软雅黑',
            'body': '微软雅黑'
        },
        'style': 'research-dual',  # 研究型蓝绿双色调
        'slide_count': 40
    },
    'zhinang_qa': {
        'name': '输配作业智囊',
        'source': '融合大模型与知识图谱的输配作业智能问答',
        'color_scheme': {
            'primary': '#00479D',      # 深蓝
            'secondary': '#5B9BD5',    # 天蓝
            'accent': '#F2F2F2',       # 浅灰
            'background': '#FFFFFF',
            'text': '#333333',
            'text_light': '#666666'
        },
        'fonts': {
            'title': '微软雅黑',
            'body': '微软雅黑'
        },
        'style': 'tech-gradient',  # 科技渐变风格
        'slide_count': 14
    }
}


def hex_to_rgb(hex_color):
    """十六进制颜色转 RGB"""
    hex_color = hex_color.lstrip('#')
    return RGBColor(
        int(hex_color[0:2], 16),
        int(hex_color[2:4], 16),
        int(hex_color[4:6], 16)
    )


def create_cover_layout(prs, theme):
    """创建封面 layout"""
    blank_layout = prs.slide_layouts[6]  # 空白 layout
    slide = prs.slides.add_slide(blank_layout)
    
    # 添加标题
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(2.5), Inches(12), Inches(1.5)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "标题"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(theme['color_scheme']['primary'])
    p.font.name = theme['fonts']['title']
    
    # 添加副标题
    subtitle_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(4.2), Inches(12), Inches(1)
    )
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.text = "副标题"
    p.font.size = Pt(24)
    p.font.color.rgb = hex_to_rgb(theme['color_scheme']['text'])
    p.font.name = theme['fonts']['body']
    
    # 添加装饰线
    line = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        Inches(0.5), Inches(4.0), Inches(2), Inches(0.05)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = hex_to_rgb(theme['color_scheme']['secondary'])
    line.line.fill.background()
    
    return slide


def create_title_content_layout(prs, theme):
    """创建标题+内容 layout"""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    
    # 标题
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.4), Inches(12), Inches(1)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "标题"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(theme['color_scheme']['primary'])
    p.font.name = theme['fonts']['title']
    
    # 左侧内容
    left_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.6), Inches(6.5), Inches(5)
    )
    tf = left_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "要点内容"
    p.font.size = Pt(18)
    p.font.color.rgb = hex_to_rgb(theme['color_scheme']['text'])
    p.font.name = theme['fonts']['body']
    
    # 添加 bullet
    for i in range(3):
        p = tf.add_paragraph()
        p.text = f"要点 {i+1}"
        p.font.size = Pt(16)
        p.level = 0
    
    # 右侧图片区域（placeholder）
    placeholder = slide.shapes.add_shape(
        1,
        Inches(7.2), Inches(1.6), Inches(5.3), Inches(5)
    )
    placeholder.fill.solid()
    placeholder.fill.fore_color.rgb = hex_to_rgb(theme['color_scheme']['accent'])
    placeholder.line.color.rgb = hex_to_rgb(theme['color_scheme']['secondary'])
    
    # 添加提示文字
    placeholder_text = slide.shapes.add_textbox(
        Inches(7.2), Inches(3.8), Inches(5.3), Inches(0.5)
    )
    tf = placeholder_text.text_frame
    p = tf.paragraphs[0]
    p.text = "[图片/图表区域]"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(14)
    p.font.color.rgb = hex_to_rgb(theme['color_scheme']['text_light'])
    
    return slide


def create_section_break_layout(prs, theme):
    """创建章节过渡页 layout"""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    
    # 背景色块
    bg = slide.shapes.add_shape(
        1,
        Inches(0), Inches(0), Inches(13.333), Inches(7.5)
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = hex_to_rgb(theme['color_scheme']['primary'])
    bg.line.fill.background()
    
    # 章节号
    num_box = slide.shapes.add_textbox(
        Inches(0.8), Inches(2.5), Inches(2), Inches(1)
    )
    tf = num_box.text_frame
    p = tf.paragraphs[0]
    p.text = "01"
    p.font.size = Pt(72)
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb('#FFFFFF')
    
    # 章节标题
    title_box = slide.shapes.add_textbox(
        Inches(0.8), Inches(3.8), Inches(11), Inches(1.2)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "章节标题"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb('#FFFFFF')
    p.font.name = theme['fonts']['title']
    
    return slide


def create_data_chart_layout(prs, theme):
    """创建数据图表页 layout"""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    
    # 标题
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.4), Inches(12), Inches(0.8)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "数据标题"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(theme['color_scheme']['primary'])
    
    # 图表区域
    chart_area = slide.shapes.add_shape(
        1,
        Inches(1), Inches(1.5), Inches(11), Inches(5)
    )
    chart_area.fill.solid()
    chart_area.fill.fore_color.rgb = hex_to_rgb('#F5F5F5')
    chart_area.line.color.rgb = hex_to_rgb(theme['color_scheme']['secondary'])
    
    # 提示文字
    chart_text = slide.shapes.add_textbox(
        Inches(1), Inches(3.8), Inches(11), Inches(0.5)
    )
    tf = chart_text.text_frame
    p = tf.paragraphs[0]
    p.text = "[图表区域 - matplotlib 输出]"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(14)
    p.font.color.rgb = hex_to_rgb(theme['color_scheme']['text_light'])
    
    # Takeaway 文字
    takeaway_box = slide.shapes.add_textbox(
        Inches(1), Inches(6.8), Inches(11), Inches(0.6)
    )
    tf = takeaway_box.text_frame
    p = tf.paragraphs[0]
    p.text = "核心结论：年复合增长率达 60%"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(theme['color_scheme']['secondary'])
    
    return slide


def create_summary_layout(prs, theme):
    """创建总结页 layout"""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    
    # 标题
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.4), Inches(12), Inches(0.8)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "总结与展望"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(theme['color_scheme']['primary'])
    
    # 要点列表
    content_box = slide.shapes.add_textbox(
        Inches(1), Inches(1.8), Inches(11), Inches(4)
    )
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i in range(4):
        p = tf.add_paragraph()
        p.text = f"总结要点 {i+1}"
        p.font.size = Pt(20)
        p.font.color.rgb = hex_to_rgb(theme['color_scheme']['text'])
        p.space_after = Pt(12)
    
    # CTA 区域
    cta_box = slide.shapes.add_shape(
        1,
        Inches(1), Inches(6), Inches(11), Inches(1)
    )
    cta_box.fill.solid()
    cta_box.fill.fore_color.rgb = hex_to_rgb(theme['color_scheme']['secondary'])
    cta_box.line.fill.background()
    
    cta_text = slide.shapes.add_textbox(
        Inches(1), Inches(6.2), Inches(11), Inches(0.6)
    )
    tf = cta_text.text_frame
    p = tf.paragraphs[0]
    p.text = "下一步行动建议"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb('#FFFFFF')
    
    return slide


def create_template(theme_key, theme_config):
    """创建完整模板"""
    print(f"\n创建模板: {theme_config['name']}")
    
    # 创建 presentation
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # 创建 5 个核心 layout
    layouts = [
        ('cover', create_cover_layout),
        ('title_content', create_title_content_layout),
        ('section_break', create_section_break_layout),
        ('data_chart', create_data_chart_layout),
        ('summary', create_summary_layout),
    ]
    
    for name, create_fn in layouts:
        create_fn(prs, theme_config)
        print(f"  [OK] 创建 layout: {name}")
    
    # 保存
    output_path = f"themes/{theme_key}.pptx"
    prs.save(output_path)
    print(f"  [OK] 保存到: {output_path}")
    
    # 生成 style_profile.json
    profile = {
        "theme_name": theme_key,
        "display_name": theme_config['name'],
        "source": theme_config['source'],
        "color_scheme": theme_config['color_scheme'],
        "typography": {
            "title_font": theme_config['fonts']['title'],
            "body_font": theme_config['fonts']['body'],
            "title_size_pt": 32,
            "body_size_pt": 18
        },
        "layout_style": {
            "margin_cm": 2.0,
            "content_alignment": "left",
            "visual_weight": "medium",
            "decoration": "minimal"
        },
        "design_language": theme_config['style'],
        "available_layouts": [name for name, _ in layouts],
        "generation_metadata": {
            "tool": "manual_analysis",
            "source_slides": theme_config['slide_count'],
            "created_from": "pptx_analysis"
        }
    }
    
    profile_path = f"themes/{theme_key}_profile.json"
    with open(profile_path, 'w', encoding='utf-8') as f:
        json.dump(profile, f, ensure_ascii=False, indent=2)
    print(f"  [OK] 保存配置文件: {profile_path}")
    
    return output_path


def main():
    print("="*60)
    print("基于分析结果创建 PPTX 模板")
    print("="*60)
    
    # 确保目录存在
    Path('themes').mkdir(exist_ok=True)
    
    # 创建所有模板
    created = []
    for key, config in TEMPLATES.items():
        path = create_template(key, config)
        created.append(path)
    
    print("\n" + "="*60)
    print("模板创建完成！")
    print("="*60)
    print("\n生成的文件:")
    for f in created:
        print(f"  - {f}")
    print(f"\n共创建 {len(created)} 个模板")


if __name__ == '__main__':
    main()
