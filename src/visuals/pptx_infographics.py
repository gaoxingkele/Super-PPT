# -*- coding: utf-8 -*-
"""
原生 PPTX 信息图渲染器。
使用 python-pptx 原生形状（矩形、圆形、线条、文本框）直接在幻灯片上渲染信息图，
无需生成中间图片。

支持类型：process_flow, stat_display, timeline, hierarchy, comparison,
          matrix, network, pyramid, cycle。
"""
import logging
import math

import src  # noqa: F401

from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

logger = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# 颜色工具
# ---------------------------------------------------------------------------

def _hex_to_rgb(hex_str: str) -> RGBColor:
    """将 '#RRGGBB' 或 'RRGGBB' 格式的十六进制颜色转换为 RGBColor。"""
    h = hex_str.lstrip("#")
    if len(h) != 6:
        h = "333333"
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _lighten_hex(hex_str: str, factor: float = 0.4) -> str:
    """将十六进制颜色变浅，返回新的 hex 字符串。"""
    h = hex_str.lstrip("#")
    if len(h) != 6:
        h = "333333"
    r, g, b = int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
    r = int(r + (255 - r) * factor)
    g = int(g + (255 - g) * factor)
    b = int(b + (255 - b) * factor)
    return f"#{r:02X}{g:02X}{b:02X}"


def _get_palette(color_scheme: dict, n: int) -> list:
    """从配色方案中生成 n 个颜色（hex 字符串列表），循环使用。"""
    base = [
        color_scheme.get("primary", "#1B365D"),
        color_scheme.get("accent", "#E8612D"),
        color_scheme.get("secondary", "#4A90D9"),
        "#2ECC71", "#9B59B6", "#F39C12", "#E74C3C", "#1ABC9C",
    ]
    return (base * ((n // len(base)) + 1))[:n]


# ---------------------------------------------------------------------------
# 形状辅助函数
# ---------------------------------------------------------------------------

def _add_rounded_rect(slide, left, top, width, height, fill_hex,
                      border_hex=None, border_width=Pt(0)):
    """添加圆角矩形形状并返回。"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height,
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = _hex_to_rgb(fill_hex)
    if border_hex:
        shape.line.color.rgb = _hex_to_rgb(border_hex)
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    return shape


def _add_circle(slide, cx, cy, diameter, fill_hex):
    """添加圆形形状，cx/cy 为中心坐标。返回 shape。"""
    r = diameter // 2
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, cx - r, cy - r, diameter, diameter,
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = _hex_to_rgb(fill_hex)
    shape.line.fill.background()
    return shape


def _add_rect(slide, left, top, width, height, fill_hex):
    """添加普通矩形形状。"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height,
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = _hex_to_rgb(fill_hex)
    shape.line.fill.background()
    return shape


def _add_textbox(slide, left, top, width, height, text, font_size=Pt(12),
                 font_color="#FFFFFF", bold=False, alignment=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE):
    """添加文本框并设置文字样式。"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    try:
        tf.vertical_anchor = anchor
    except Exception:
        pass
    p = tf.paragraphs[0]
    p.text = str(text)
    p.font.size = font_size
    p.font.bold = bold
    p.font.color.rgb = _hex_to_rgb(font_color)
    p.alignment = alignment
    return txBox


def _add_arrow_shape(slide, left, top, width, height, fill_hex):
    """添加右箭头形状。"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW, left, top, width, height,
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = _hex_to_rgb(fill_hex)
    shape.line.fill.background()
    return shape


def _add_line(slide, start_x, start_y, end_x, end_y, color_hex,
              width=Pt(2)):
    """添加直线连接器。"""
    connector = slide.shapes.add_connector(
        1,  # MSO_CONNECTOR_TYPE.STRAIGHT = 1
        start_x, start_y, end_x, end_y,
    )
    connector.line.color.rgb = _hex_to_rgb(color_hex)
    connector.line.width = width
    return connector


def _set_shape_text(shape, text, font_size=Pt(12), font_color="#FFFFFF",
                    bold=False, alignment=PP_ALIGN.CENTER):
    """在已有形状上设置文本。"""
    tf = shape.text_frame
    tf.word_wrap = True
    try:
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    except Exception:
        pass
    p = tf.paragraphs[0]
    p.text = str(text)
    p.font.size = font_size
    p.font.bold = bold
    p.font.color.rgb = _hex_to_rgb(font_color)
    p.alignment = alignment


# ---------------------------------------------------------------------------
# 主入口
# ---------------------------------------------------------------------------

def render_native_infographic(slide, visual: dict, color_scheme: dict,
                               left, top, width, height):
    """在 slide 上直接渲染原生 PPTX 信息图。

    Args:
        slide: python-pptx slide object
        visual: slide_plan 中的 visual 字段
        color_scheme: 配色方案
        left, top, width, height: 信息图区域位置和尺寸 (pptx.util 单位)
    """
    infographic_type = visual.get("infographic_type", "process_flow")
    data = visual.get("data", {})
    if not isinstance(data, dict):
        data = {}

    renderer = _NATIVE_RENDERERS.get(infographic_type, _render_native_placeholder)
    try:
        renderer(slide, data, color_scheme, left, top, width, height)
    except Exception as exc:
        logger.warning("原生信息图渲染失败 (%s): %s，使用占位符", infographic_type, exc)
        try:
            _render_native_placeholder(slide, data, color_scheme,
                                       left, top, width, height)
        except Exception:
            pass


# ===========================================================================
#  1. process_flow — 流程图
# ===========================================================================

def _render_process_flow(slide, data, color_scheme, left, top, width, height):
    """横向圆角矩形 + 箭头连接的流程图。"""
    stages = data.get("stages", [])
    if not stages:
        stages = [{"name": "步骤 1", "detail": ""}]

    n = len(stages)
    colors = _get_palette(color_scheme, n)
    accent = color_scheme.get("accent", "#E8612D")

    # 计算每个 box 的尺寸
    arrow_w = Emu(int(width * 0.04)) if n > 1 else 0
    total_arrow_w = arrow_w * max(n - 1, 0)
    box_w = Emu(int((width - total_arrow_w) / n * 0.9))
    gap = Emu(int((width - total_arrow_w - box_w * n) / max(n + 1, 1)))
    box_h = Emu(int(height * 0.65))
    y_box = top + Emu(int(height * 0.18))
    badge_d = Emu(int(min(box_w, box_h) * 0.18))

    x_cursor = left + gap

    for i, stage in enumerate(stages):
        c = colors[i]
        name = stage.get("name", f"步骤 {i+1}") if isinstance(stage, dict) else str(stage)
        detail = stage.get("detail", "") if isinstance(stage, dict) else ""

        # 圆角矩形
        rect = _add_rounded_rect(slide, x_cursor, y_box, box_w, box_h, c)

        # 名称文本（上半部）
        _add_textbox(slide, x_cursor, y_box + Emu(int(box_h * 0.2)),
                     box_w, Emu(int(box_h * 0.35)),
                     name, font_size=Pt(14), font_color="#FFFFFF", bold=True)

        # 详情文本（下半部）
        if detail:
            _add_textbox(slide, x_cursor, y_box + Emu(int(box_h * 0.55)),
                         box_w, Emu(int(box_h * 0.35)),
                         detail, font_size=Pt(10), font_color="#E0E0E0", bold=False)

        # 编号圆圈（左上角）
        badge_x = x_cursor + Emu(int(box_w * 0.06))
        badge_y = y_box + Emu(int(box_h * 0.06))
        _add_circle(slide, badge_x + badge_d // 2, badge_y + badge_d // 2,
                    badge_d, "#FFFFFF")
        _add_textbox(slide, badge_x, badge_y, badge_d, badge_d,
                     str(i + 1), font_size=Pt(10), font_color=c, bold=True)

        x_cursor += box_w

        # 箭头
        if i < n - 1:
            arrow_h = Emu(int(box_h * 0.25))
            arrow_y = y_box + Emu(int(box_h * 0.38))
            _add_arrow_shape(slide, x_cursor, arrow_y, arrow_w, arrow_h, accent)
            x_cursor += arrow_w

        x_cursor += gap


# ===========================================================================
#  2. stat_display — KPI 大数字
# ===========================================================================

def _render_stat_display(slide, data, color_scheme, left, top, width, height):
    """KPI 大数字卡片行。"""
    kpis = data.get("kpis", data.get("items", []))
    if not kpis:
        kpis = [{"label": "数据", "value": "N/A"}]

    n = len(kpis)
    colors = _get_palette(color_scheme, n)
    primary = color_scheme.get("primary", "#1B365D")

    card_gap = Emu(int(width * 0.02))
    card_w = Emu(int((width - card_gap * (n + 1)) / n))
    card_h = Emu(int(height * 0.85))
    y_card = top + Emu(int(height * 0.08))
    bar_h = Emu(int(card_h * 0.06))

    x_cursor = left + card_gap

    for i, kpi in enumerate(kpis):
        c = colors[i]
        value = str(kpi.get("value", "N/A"))
        label = str(kpi.get("label", ""))
        trend = kpi.get("trend", "")

        # 浅色背景卡片
        light_bg = _lighten_hex(c, 0.75)
        _add_rounded_rect(slide, x_cursor, y_card, card_w, card_h,
                          light_bg, border_hex=_lighten_hex(c, 0.5),
                          border_width=Pt(1))

        # 大数字
        value_h = Emu(int(card_h * 0.45))
        _add_textbox(slide, x_cursor, y_card + Emu(int(card_h * 0.1)),
                     card_w, value_h,
                     value, font_size=Pt(44), font_color=c, bold=True)

        # 趋势箭头
        if trend in ("up", "down"):
            arrow_text = "\u25B2" if trend == "up" else "\u25BC"
            arrow_color = "#2ECC71" if trend == "up" else "#E74C3C"
            trend_w = Emu(int(card_w * 0.2))
            _add_textbox(slide, x_cursor + card_w - trend_w - Emu(int(card_w * 0.05)),
                         y_card + Emu(int(card_h * 0.15)),
                         trend_w, Emu(int(card_h * 0.2)),
                         arrow_text, font_size=Pt(20), font_color=arrow_color,
                         bold=True)

        # 标签
        label_h = Emu(int(card_h * 0.2))
        _add_textbox(slide, x_cursor, y_card + Emu(int(card_h * 0.58)),
                     card_w, label_h,
                     label, font_size=Pt(14), font_color="#666666", bold=False)

        # 底部色条
        _add_rect(slide,
                  x_cursor, y_card + card_h - bar_h,
                  card_w, bar_h, c)

        x_cursor += card_w + card_gap


# ===========================================================================
#  3. timeline — 时间线
# ===========================================================================

def _render_timeline(slide, data, color_scheme, left, top, width, height):
    """水平时间线，交替上下节点。"""
    events = data.get("events", data.get("stages", []))
    if not events:
        events = [{"date": "", "title": "事件", "description": ""}]

    n = len(events)
    colors = _get_palette(color_scheme, n)
    primary = color_scheme.get("primary", "#1B365D")

    # 中线
    line_y = top + Emu(int(height * 0.5))
    line_left = left + Emu(int(width * 0.03))
    line_right = left + Emu(int(width * 0.97))
    _add_rect(slide, line_left, line_y - Pt(2), line_right - line_left, Pt(4), primary)

    node_d = Emu(int(min(width / max(n, 1), height) * 0.12))
    node_d = max(node_d, Emu(Inches(0.2)))

    for i, event in enumerate(events):
        title = event.get("title", event.get("name", "")) if isinstance(event, dict) else str(event)
        date = event.get("date", "") if isinstance(event, dict) else ""
        desc = event.get("description", "") if isinstance(event, dict) else ""
        c = colors[i]
        above = (i % 2 == 0)

        # X position along line
        x = line_left + Emu(int((line_right - line_left) * (i + 0.5) / n))

        # 节点圆
        _add_circle(slide, x, line_y, node_d, c)

        # 垂直连接线
        connector_len = Emu(int(height * 0.18))
        if above:
            conn_start_y = line_y - node_d // 2
            conn_end_y = conn_start_y - connector_len
            _add_line(slide, x, conn_start_y, x, conn_end_y, c, Pt(2))
        else:
            conn_start_y = line_y + node_d // 2
            conn_end_y = conn_start_y + connector_len
            _add_line(slide, x, conn_start_y, x, conn_end_y, c, Pt(2))

        # 文本区域
        text_w = Emu(int(width / n * 0.85))
        text_h = Emu(int(height * 0.22))
        text_x = x - text_w // 2

        if above:
            text_y = conn_end_y - text_h
        else:
            text_y = conn_end_y

        # 日期
        if date:
            date_h = Emu(int(text_h * 0.35))
            if above:
                _add_textbox(slide, text_x, text_y + text_h - date_h,
                             text_w, date_h,
                             str(date), font_size=Pt(9), font_color="#888888",
                             bold=False)
                title_y = text_y
            else:
                _add_textbox(slide, text_x, text_y,
                             text_w, date_h,
                             str(date), font_size=Pt(9), font_color="#888888",
                             bold=False)
                title_y = text_y + date_h
        else:
            title_y = text_y

        # 标题
        title_h = Emu(int(text_h * 0.65))
        _add_textbox(slide, text_x, title_y, text_w, title_h,
                     title, font_size=Pt(11), font_color=c, bold=True)


# ===========================================================================
#  4. hierarchy — 层级图
# ===========================================================================

def _render_hierarchy(slide, data, color_scheme, left, top, width, height):
    """自上而下的树状层级图。"""
    nodes = data.get("nodes", [])
    if not nodes:
        nodes = [{"name": "根节点", "children": []}]

    colors = _get_palette(color_scheme, 8)
    primary = color_scheme.get("primary", "#1B365D")

    # -- Level 0: Root --
    root = nodes[0] if nodes else {"name": "根节点"}
    root_name = root.get("name", "根节点") if isinstance(root, dict) else str(root)
    root_w = Emu(int(width * 0.25))
    root_h = Emu(int(height * 0.14))
    root_x = left + (width - root_w) // 2
    root_y = top + Emu(int(height * 0.02))

    root_shape = _add_rounded_rect(slide, root_x, root_y, root_w, root_h, primary)
    _set_shape_text(root_shape, root_name, font_size=Pt(16), font_color="#FFFFFF",
                    bold=True)

    root_cx = root_x + root_w // 2
    root_bottom = root_y + root_h

    # -- Level 1: Children --
    children = []
    if isinstance(root, dict):
        children = root.get("children", [])
    if not children and len(nodes) > 1:
        children = nodes[1:]

    if not children:
        return

    nc = len(children)
    child_w = Emu(int(min(width * 0.18, width / nc * 0.85)))
    child_h = Emu(int(height * 0.12))
    child_y = top + Emu(int(height * 0.28))
    child_gap = Emu(int((width - child_w * nc) / (nc + 1)))

    child_centers = []
    for i, child in enumerate(children):
        child_name = child.get("name", str(child)) if isinstance(child, dict) else str(child)
        c = colors[i % len(colors)]
        cx = left + child_gap * (i + 1) + child_w * i
        child_shape = _add_rounded_rect(slide, cx, child_y, child_w, child_h, c)
        _set_shape_text(child_shape, child_name, font_size=Pt(12),
                        font_color="#FFFFFF", bold=True)

        child_cx = cx + child_w // 2
        child_centers.append((child_cx, c))

        # 连线 root -> child
        _add_line(slide, root_cx, root_bottom, child_cx, child_y, "#CCCCCC", Pt(2))

        # -- Level 2: Grandchildren --
        grandchildren = child.get("children", []) if isinstance(child, dict) else []
        if grandchildren:
            gc_y = top + Emu(int(height * 0.55))
            ngc = len(grandchildren)
            gc_w = Emu(int(min(child_w * 0.85, width / max(nc * ngc, 1) * 0.8)))
            gc_h = Emu(int(height * 0.10))

            # 在 child 下方均匀分布
            gc_total_w = gc_w * ngc + Emu(int(gc_w * 0.15)) * max(ngc - 1, 0)
            gc_start_x = child_cx - gc_total_w // 2

            child_bottom = child_y + child_h

            for j, gc in enumerate(grandchildren[:5]):
                gc_name = gc.get("name", str(gc)) if isinstance(gc, dict) else str(gc)
                gc_color = _lighten_hex(c, 0.25)
                gx = gc_start_x + Emu(int((gc_w + Emu(int(gc_w * 0.15))) * j))
                gc_shape = _add_rounded_rect(slide, gx, gc_y, gc_w, gc_h, gc_color)
                _set_shape_text(gc_shape, gc_name, font_size=Pt(10),
                                font_color="#FFFFFF", bold=False)

                gc_cx = gx + gc_w // 2
                _add_line(slide, child_cx, child_bottom, gc_cx, gc_y,
                          "#DDDDDD", Pt(1))


# ===========================================================================
#  5. comparison — 对比图
# ===========================================================================

def _render_comparison(slide, data, color_scheme, left, top, width, height):
    """双栏对比图，中间 VS 徽章。"""
    items = data.get("items", [])
    if not items or len(items) < 2:
        items = [{"name": "A", "metrics": {}}, {"name": "B", "metrics": {}}]

    primary = color_scheme.get("primary", "#1B365D")
    accent = color_scheme.get("accent", "#E8612D")

    col_w = Emu(int(width * 0.42))
    vs_w = Emu(int(width * 0.10))
    col_gap = Emu(int(width * 0.03))

    left_col_x = left
    right_col_x = left + col_w + vs_w + col_gap * 2
    vs_x = left + col_w + col_gap

    header_h = Emu(int(height * 0.16))
    header_y = top

    # -- Left header --
    left_item = items[0]
    left_name = left_item.get("name", "A")
    _add_rounded_rect(slide, left_col_x, header_y, col_w, header_h, primary)
    _add_textbox(slide, left_col_x, header_y, col_w, header_h,
                 left_name, font_size=Pt(20), font_color="#FFFFFF", bold=True)

    # -- Right header --
    right_item = items[1]
    right_name = right_item.get("name", "B")
    _add_rounded_rect(slide, right_col_x, header_y, col_w, header_h, accent)
    _add_textbox(slide, right_col_x, header_y, col_w, header_h,
                 right_name, font_size=Pt(20), font_color="#FFFFFF", bold=True)

    # -- VS badge --
    vs_d = Emu(int(min(vs_w, header_h) * 0.9))
    vs_cx = vs_x + vs_w // 2
    vs_cy = header_y + header_h // 2
    _add_circle(slide, vs_cx, vs_cy, vs_d, "#F0F0F0")
    _add_textbox(slide, vs_cx - vs_d // 2, vs_cy - vs_d // 2, vs_d, vs_d,
                 "VS", font_size=Pt(14), font_color="#999999", bold=True)

    # -- Metrics --
    left_metrics = left_item.get("metrics", {})
    right_metrics = right_item.get("metrics", {})

    if isinstance(left_metrics, dict):
        keys = list(left_metrics.keys())[:6]
    elif isinstance(left_metrics, list):
        keys = [str(m) for m in left_metrics[:6]]
    else:
        keys = []

    if not keys:
        return

    metric_area_top = header_y + header_h + Emu(int(height * 0.05))
    row_h = Emu(int((height - (metric_area_top - top)) / max(len(keys), 1) * 0.9))
    row_gap = Emu(int(row_h * 0.1))
    bar_h = Emu(int(row_h * 0.35))

    for idx, key in enumerate(keys):
        row_y = metric_area_top + Emu(int((row_h + row_gap) * idx))

        # 中间标签
        _add_textbox(slide, vs_x, row_y, vs_w, row_h,
                     str(key)[:10], font_size=Pt(10), font_color="#888888",
                     bold=False)

        # 获取数值
        lv = left_metrics.get(key, 0) if isinstance(left_metrics, dict) else 0
        rv = right_metrics.get(key, 0) if isinstance(right_metrics, dict) else 0

        try:
            lv_f = float(lv)
            rv_f = float(rv)
        except (ValueError, TypeError):
            lv_f, rv_f = 1.0, 1.0

        max_v = max(abs(lv_f), abs(rv_f), 0.001)

        # 左侧条（从右向左）
        l_bar_w = Emu(int(col_w * 0.85 * abs(lv_f) / max_v))
        l_bar_x = left_col_x + col_w - l_bar_w
        bar_y = row_y + (row_h - bar_h) // 2
        _add_rounded_rect(slide, l_bar_x, bar_y, l_bar_w, bar_h, primary)

        # 左数值文本
        _add_textbox(slide, left_col_x, row_y, col_w - l_bar_w - Emu(int(col_w * 0.02)),
                     row_h, str(lv), font_size=Pt(10), font_color=primary,
                     bold=True, alignment=PP_ALIGN.RIGHT)

        # 右侧条
        r_bar_w = Emu(int(col_w * 0.85 * abs(rv_f) / max_v))
        r_bar_x = right_col_x
        _add_rounded_rect(slide, r_bar_x, bar_y, r_bar_w, bar_h, accent)

        # 右数值文本
        _add_textbox(slide, right_col_x + r_bar_w + Emu(int(col_w * 0.02)), row_y,
                     col_w - r_bar_w - Emu(int(col_w * 0.02)), row_h,
                     str(rv), font_size=Pt(10), font_color=accent,
                     bold=True, alignment=PP_ALIGN.LEFT)


# ===========================================================================
#  6. matrix — 矩阵卡片
# ===========================================================================

def _render_matrix(slide, data, color_scheme, left, top, width, height):
    """网格卡片布局。"""
    items = data.get("items", data.get("cells", data.get("nodes", [])))
    if not items:
        items = [{"name": "项目", "detail": ""}]

    n = len(items)
    cols = min(3, n)
    rows = (n + cols - 1) // cols
    colors = _get_palette(color_scheme, n)

    card_gap_x = Emu(int(width * 0.03))
    card_gap_y = Emu(int(height * 0.04))
    card_w = Emu(int((width - card_gap_x * (cols + 1)) / cols))
    card_h = Emu(int((height - card_gap_y * (rows + 1)) / rows))
    accent_bar_h = Emu(int(card_h * 0.06))

    for idx, item in enumerate(items):
        row = idx // cols
        col = idx % cols

        name = item.get("name", str(item)) if isinstance(item, dict) else str(item)
        detail = item.get("detail", item.get("description", "")) if isinstance(item, dict) else ""
        c = colors[idx]

        cx = left + card_gap_x * (col + 1) + card_w * col
        cy = top + card_gap_y * (row + 1) + card_h * row

        # 浅色背景卡片
        light_bg = _lighten_hex(c, 0.75)
        _add_rounded_rect(slide, cx, cy, card_w, card_h, light_bg,
                          border_hex=_lighten_hex(c, 0.5), border_width=Pt(1))

        # 顶部彩色条
        _add_rect(slide, cx, cy, card_w, accent_bar_h, c)

        # 名称
        name_h = Emu(int(card_h * 0.35))
        _add_textbox(slide, cx, cy + accent_bar_h + Emu(int(card_h * 0.05)),
                     card_w, name_h,
                     name, font_size=Pt(13), font_color=c, bold=True)

        # 详情
        if detail:
            detail_h = Emu(int(card_h * 0.40))
            _add_textbox(slide, cx + Emu(int(card_w * 0.05)),
                         cy + accent_bar_h + name_h + Emu(int(card_h * 0.05)),
                         Emu(int(card_w * 0.9)), detail_h,
                         detail, font_size=Pt(10), font_color="#666666", bold=False,
                         anchor=MSO_ANCHOR.TOP)


# ===========================================================================
#  7. network — 关系网络
# ===========================================================================

def _render_network(slide, data, color_scheme, left, top, width, height):
    """中心辐射式关系网络图。"""
    nodes = data.get("nodes", data.get("items", []))
    if not nodes:
        nodes = ["节点"]

    center_label = data.get("center", data.get("title", "核心"))
    if isinstance(center_label, dict):
        center_label = center_label.get("name", "核心")

    colors = _get_palette(color_scheme, len(nodes) + 1)
    primary = color_scheme.get("primary", "#1B365D")

    # 区域中心
    area_cx = left + width // 2
    area_cy = top + height // 2
    min_dim = min(width, height)

    # 中心节点
    center_d = Emu(int(min_dim * 0.22))
    _add_circle(slide, area_cx, area_cy, center_d, primary)
    _add_textbox(slide, area_cx - center_d // 2, area_cy - center_d // 2,
                 center_d, center_d,
                 str(center_label)[:8], font_size=Pt(14), font_color="#FFFFFF",
                 bold=True)

    # 外围节点
    n = len(nodes)
    radius = Emu(int(min_dim * 0.38))
    node_d = Emu(int(min_dim * 0.15))

    for i, node in enumerate(nodes):
        name = node.get("name", str(node)) if isinstance(node, dict) else str(node)
        c = colors[(i + 1) % len(colors)]

        angle = 2 * math.pi * i / n - math.pi / 2
        nx = area_cx + Emu(int(radius * math.cos(angle)))
        ny = area_cy + Emu(int(radius * math.sin(angle)))

        # 连线（先画，在底层）
        _add_line(slide, area_cx, area_cy, nx, ny, "#CCCCCC", Pt(1.5))

        # 节点圆
        _add_circle(slide, nx, ny, node_d, c)
        _add_textbox(slide, nx - node_d // 2, ny - node_d // 2,
                     node_d, node_d,
                     str(name)[:8], font_size=Pt(10), font_color="#FFFFFF",
                     bold=True)


# ===========================================================================
#  8. pyramid — 金字塔
# ===========================================================================

def _render_pyramid(slide, data, color_scheme, left, top, width, height):
    """堆叠梯形金字塔图，窄在上宽在下。

    使用矩形近似梯形效果（python-pptx 不直接支持自由梯形），
    每层宽度递增，居中排列。
    """
    levels = data.get("levels", data.get("items", data.get("stages", [])))
    if not levels:
        levels = ["顶层"]

    n = len(levels)
    colors = _get_palette(color_scheme, n)

    layer_gap = Emu(int(height * 0.02))
    total_gap = layer_gap * (n - 1)
    layer_h = Emu(int((height - total_gap) / n))

    # 最窄宽度（顶）和最宽宽度（底）
    min_w = Emu(int(width * 0.25))
    max_w = Emu(int(width * 0.95))

    for i in range(n):
        name = levels[i] if isinstance(levels[i], str) else levels[i].get("name", str(levels[i]))
        c = colors[i]

        # 宽度从窄到宽
        if n > 1:
            frac = i / (n - 1)
        else:
            frac = 1.0
        layer_w = Emu(int(min_w + (max_w - min_w) * frac))
        layer_x = left + (width - layer_w) // 2
        layer_y = top + Emu(int((layer_h + layer_gap) * i))

        # 使用 CHEVRON 或 PENTAGON 形状来模拟梯形
        # 改用六边形（HEXAGON）更好看，但最安全的是圆角矩形
        shape = _add_rounded_rect(slide, layer_x, layer_y, layer_w, layer_h, c)
        _set_shape_text(shape, name, font_size=Pt(14), font_color="#FFFFFF",
                        bold=True)


# ===========================================================================
#  9. cycle — 循环图
# ===========================================================================

def _render_cycle(slide, data, color_scheme, left, top, width, height):
    """环形循环图，节点沿圆排列并用箭头连接。"""
    stages = data.get("stages", data.get("items", []))
    if not stages:
        stages = [{"name": "阶段"}]

    n = len(stages)
    colors = _get_palette(color_scheme, n)

    area_cx = left + width // 2
    area_cy = top + height // 2
    min_dim = min(width, height)

    radius = Emu(int(min_dim * 0.35))
    node_d = Emu(int(min_dim * 0.16))

    # 计算每个节点位置
    positions = []
    for i in range(n):
        angle = 2 * math.pi * i / n - math.pi / 2
        nx = area_cx + Emu(int(radius * math.cos(angle)))
        ny = area_cy + Emu(int(radius * math.sin(angle)))
        positions.append((nx, ny))

    # 先画箭头连接器（顺时针方向）
    if n > 1:
        accent = color_scheme.get("accent", "#E8612D")
        for i in range(n):
            j = (i + 1) % n
            sx, sy = positions[i]
            ex, ey = positions[j]

            # 计算缩短的起止点，避免箭头覆盖节点
            angle = math.atan2(ey - sy, ex - sx)
            offset = node_d // 2 + Emu(int(min_dim * 0.02))
            sx_adj = sx + Emu(int(offset * math.cos(angle)))
            sy_adj = sy + Emu(int(offset * math.sin(angle)))
            ex_adj = ex - Emu(int(offset * math.cos(angle)))
            ey_adj = ey - Emu(int(offset * math.sin(angle)))

            # 使用小箭头形状替代连接线
            # 计算箭头长度和位置
            dx = ex_adj - sx_adj
            dy = ey_adj - sy_adj
            arr_len = Emu(int(math.sqrt(dx * dx + dy * dy)))

            # 简单连接线
            _add_line(slide, sx_adj, sy_adj, ex_adj, ey_adj, "#CCCCCC", Pt(2))

            # 在终点附近添加小三角箭头
            arrow_size = Emu(int(min_dim * 0.03))
            mid_x = ex_adj
            mid_y = ey_adj
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.ISOSCELES_TRIANGLE,
                mid_x - arrow_size // 2,
                mid_y - arrow_size // 2,
                arrow_size, arrow_size,
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = _hex_to_rgb("#CCCCCC")
            arrow.line.fill.background()
            # 旋转箭头指向正确方向
            angle_deg = math.degrees(angle) + 90  # 三角形默认朝上
            arrow.rotation = angle_deg

    # 画节点
    for i, stage in enumerate(stages):
        name = stage.get("name", str(stage)) if isinstance(stage, dict) else str(stage)
        c = colors[i]
        nx, ny = positions[i]

        _add_circle(slide, nx, ny, node_d, c)
        _add_textbox(slide, nx - node_d // 2, ny - node_d // 2,
                     node_d, node_d,
                     str(name)[:8], font_size=Pt(11), font_color="#FFFFFF",
                     bold=True)


# ===========================================================================
#  占位符（未知类型回退）
# ===========================================================================

def _render_native_placeholder(slide, data, color_scheme, left, top, width, height):
    """未知信息图类型的通用占位。"""
    primary = color_scheme.get("primary", "#1B365D")
    light_bg = _lighten_hex(primary, 0.85)

    _add_rounded_rect(slide, left, top, width, height, light_bg,
                      border_hex=primary, border_width=Pt(2))

    # 顶部装饰条
    bar_h = Emu(int(height * 0.06))
    _add_rect(slide, left, top, width, bar_h, primary)

    # 文字
    desc = data.get("description", "信息图")
    if not desc:
        desc = "信息图"
    _add_textbox(slide, left + Emu(int(width * 0.05)),
                 top + bar_h + Emu(int(height * 0.1)),
                 Emu(int(width * 0.9)), Emu(int(height * 0.7)),
                 str(desc), font_size=Pt(14), font_color=primary,
                 bold=False, anchor=MSO_ANCHOR.TOP)


# ===========================================================================
#  渲染器注册表
# ===========================================================================

_NATIVE_RENDERERS = {
    "process_flow": _render_process_flow,
    "stat_display": _render_stat_display,
    "timeline": _render_timeline,
    "hierarchy": _render_hierarchy,
    "comparison": _render_comparison,
    "matrix": _render_matrix,
    "network": _render_network,
    "pyramid": _render_pyramid,
    "cycle": _render_cycle,
}
