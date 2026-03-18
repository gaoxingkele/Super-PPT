# -*- coding: utf-8 -*-
"""
原生 python-pptx 图表渲染器。
在 slide 上直接使用 PPTX 原生图表 API，无需 matplotlib 图片。
"""
import src  # noqa: F401

from pptx.util import Inches, Pt, Emu
from pptx.chart.data import CategoryChartData, XyChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor


# ============ 颜色工具 ============

def _get_colors(color_scheme: dict, n: int = 6) -> list:
    """从配色方案生成 N 个颜色字符串（不含 #）。"""
    base_colors = [
        color_scheme.get("primary", "#1B365D"),
        color_scheme.get("accent", "#E8612D"),
        color_scheme.get("secondary", "#4A90D9"),
        "#2ECC71", "#9B59B6", "#F39C12", "#E74C3C", "#1ABC9C",
    ]
    colors = (base_colors * ((n // len(base_colors)) + 1))[:n]
    return [c.lstrip("#") for c in colors]


def _rgb(hex_color: str) -> RGBColor:
    """将 hex 颜色字符串转为 RGBColor，兼容带/不带 # 前缀。"""
    return RGBColor.from_string(hex_color.lstrip("#"))


# ============ 字体配置 ============

_ZH_FONT = "Microsoft YaHei"


def _set_chart_font(chart, font_name: str = _ZH_FONT, size: Pt = Pt(10)):
    """统一设置图表字体，支持中文。"""
    try:
        chart.font.name = font_name
        chart.font.size = size
    except Exception:
        pass


def _style_chart_common(chart, title: str = ""):
    """通用图表样式：设置标题、字体、去除多余网格线。"""
    # 标题
    if title:
        chart.has_title = True
        chart.chart_title.text_frame.text = title
        for para in chart.chart_title.text_frame.paragraphs:
            for run in para.runs:
                run.font.name = _ZH_FONT
                run.font.size = Pt(14)
                run.font.bold = True
    else:
        chart.has_title = False

    # 字体
    _set_chart_font(chart)

    # 图例
    if chart.has_legend:
        chart.legend.font.name = _ZH_FONT
        chart.legend.font.size = Pt(9)
        chart.legend.include_in_layout = False


def _apply_series_color(series, hex_no_hash: str):
    """为图表系列应用填充颜色。"""
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = RGBColor.from_string(hex_no_hash)


def _add_data_labels(series, font_size: Pt = Pt(10), number_format: str = "General",
                     position=None):
    """为系列添加数据标签。"""
    series.has_data_labels = True
    data_labels = series.data_labels
    data_labels.font.size = font_size
    data_labels.font.name = _ZH_FONT
    data_labels.number_format = number_format
    data_labels.show_value = True
    data_labels.show_category_name = False
    data_labels.show_series_name = False
    if position is not None:
        data_labels.label_position = position


def _remove_gridlines(chart):
    """移除多余的网格线，保持简洁。"""
    try:
        value_axis = chart.value_axis
        value_axis.has_major_gridlines = False
        value_axis.has_minor_gridlines = False
    except Exception:
        pass
    try:
        category_axis = chart.category_axis
        category_axis.has_major_gridlines = False
        category_axis.has_minor_gridlines = False
    except Exception:
        pass


def _style_axis_fonts(chart):
    """设置坐标轴字体为中文字体。"""
    for axis_attr in ("value_axis", "category_axis"):
        try:
            axis = getattr(chart, axis_attr)
            axis.tick_labels.font.name = _ZH_FONT
            axis.tick_labels.font.size = Pt(9)
        except Exception:
            pass


# ============ 统一入口 ============

def render_native_chart(slide, visual: dict, color_scheme: dict,
                        left, top, width, height):
    """在 slide 上直接渲染原生 PPTX 图表。

    Args:
        slide: python-pptx slide object
        visual: slide_plan 中的 visual 字段，包含 chart 和 data
        color_scheme: 配色方案 {"primary": "#xxx", "accent": "#xxx", ...}
        left, top, width, height: 图表位置和尺寸 (pptx.util 单位)
    """
    chart_type = visual.get("chart") or visual.get("chart_type", "bar")
    data = visual.get("data", {})

    if not data:
        return None

    # === 数据规范化：兼容 LLM 各种输出格式 ===
    data = _normalize_chart_data(data, chart_type)

    renderer = _CHART_RENDERERS.get(chart_type)
    if not renderer:
        # 回退到 bar 图
        renderer = _render_bar

    return renderer(slide, data, color_scheme, left, top, width, height)


def _normalize_chart_data(data: dict, chart_type: str) -> dict:
    """规范化 LLM 输出的数据格式，统一为渲染器期望的格式。"""
    data = dict(data)  # shallow copy

    # 1) labels 别名适配：categories / companies / items -> labels
    if not data.get("labels"):
        for alias in ("categories", "companies", "items", "names"):
            if data.get(alias):
                data["labels"] = data.pop(alias)
                break

    # 2) series: list-of-dict -> dict 格式
    #    LLM 常输出 [{"name":"A","values":[1,2]}, ...] 而渲染器期望 {"A":[1,2], ...}
    series = data.get("series")
    if isinstance(series, list) and series and isinstance(series[0], dict):
        new_series = {}
        for item in series:
            name = item.get("name", item.get("label", f"series_{len(new_series)}"))
            vals = item.get("values", item.get("data", []))
            new_series[name] = vals
        data["series"] = new_series

    # 3) 自动从多个同名 *_values 字段合成 series（如 gaas_values / gan_values）
    if not data.get("series") and not data.get("values"):
        val_keys = [k for k in data if k.endswith("_values") or k.endswith("_data")]
        if len(val_keys) >= 2 and data.get("labels"):
            new_series = {}
            for k in val_keys:
                name = k.replace("_values", "").replace("_data", "").upper()
                new_series[name] = data[k]
            data["series"] = new_series

    # 4) 自动从多个同名数值列合成 series（如 sales_2022 / sales_2024）
    if not data.get("series") and not data.get("values") and data.get("labels"):
        labels = data["labels"]
        num_keys = [k for k in data
                    if k != "labels" and isinstance(data[k], list)
                    and len(data[k]) == len(labels)
                    and all(isinstance(v, (int, float)) for v in data[k])]
        if len(num_keys) >= 2:
            new_series = {}
            for k in num_keys:
                new_series[k] = data[k]
            data["series"] = new_series
        elif len(num_keys) == 1:
            data["values"] = data[num_keys[0]]

    # 5) radar 图用 categories 而非 labels
    if chart_type == "radar" and not data.get("categories") and data.get("labels"):
        data["categories"] = data["labels"]

    return data


# ============ 各类图表渲染器 ============

def _render_bar(slide, data: dict, color_scheme: dict,
                left, top, width, height):
    """柱状图 — COLUMN_CLUSTERED。"""
    # 自动委托：有 series 无 values 时转为 grouped_bar
    if data.get("series") and not data.get("values"):
        return _render_grouped_bar(slide, data, color_scheme,
                                   left, top, width, height)

    labels = data.get("labels", [])
    values = data.get("values", [])
    if not labels or not values:
        return None

    chart_data = CategoryChartData()
    chart_data.categories = labels
    chart_data.add_series("数据", values)

    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, left, top, width, height, chart_data
    )
    chart = chart_frame.chart
    colors = _get_colors(color_scheme, 1)

    # 样式
    _style_chart_common(chart, data.get("title", ""))
    _remove_gridlines(chart)
    _style_axis_fonts(chart)

    # 系列颜色和数据标签
    series = chart.series[0]
    _apply_series_color(series, colors[0])
    _add_data_labels(series, position=XL_LABEL_POSITION.OUTSIDE_END)

    # 隐藏图例（单系列不需要）
    chart.has_legend = False

    return chart_frame


def _render_grouped_bar(slide, data: dict, color_scheme: dict,
                        left, top, width, height):
    """分组柱状图 — COLUMN_CLUSTERED（多系列）。"""
    labels = data.get("labels", [])
    series_data = data.get("series", {})
    if not labels or not series_data:
        return None

    chart_data = CategoryChartData()
    chart_data.categories = labels
    for name, values in series_data.items():
        chart_data.add_series(name, values)

    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, left, top, width, height, chart_data
    )
    chart = chart_frame.chart
    colors = _get_colors(color_scheme, len(series_data))

    _style_chart_common(chart, data.get("title", ""))
    _remove_gridlines(chart)
    _style_axis_fonts(chart)

    # 图例
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False

    # 各系列颜色和数据标签
    for i, series in enumerate(chart.series):
        _apply_series_color(series, colors[i])
        _add_data_labels(series, position=XL_LABEL_POSITION.OUTSIDE_END)

    return chart_frame


def _render_line(slide, data: dict, color_scheme: dict,
                 left, top, width, height):
    """折线图 — LINE_MARKERS。"""
    labels = data.get("labels", [])
    series_data = data.get("series", {})
    if not series_data:
        values = data.get("values", [])
        if values:
            series_data = {"数据": values}

    if not labels or not series_data:
        return None

    chart_data = CategoryChartData()
    chart_data.categories = labels
    for name, values in series_data.items():
        chart_data.add_series(name, values)

    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE_MARKERS, left, top, width, height, chart_data
    )
    chart = chart_frame.chart
    colors = _get_colors(color_scheme, len(series_data))

    _style_chart_common(chart, data.get("title", ""))
    _remove_gridlines(chart)
    _style_axis_fonts(chart)

    # 图例
    if len(series_data) > 1:
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
    else:
        chart.has_legend = False

    # 各系列颜色、线条样式和数据标签
    for i, series in enumerate(chart.series):
        _apply_series_color(series, colors[i])
        series.format.line.color.rgb = RGBColor.from_string(colors[i])
        series.format.line.width = Pt(2.5)
        series.smooth = False
        _add_data_labels(series, position=XL_LABEL_POSITION.ABOVE)

    return chart_frame


def _render_pie(slide, data: dict, color_scheme: dict,
                left, top, width, height):
    """饼图 — PIE。"""
    labels = data.get("labels", [])
    values = data.get("values", [])
    if not labels or not values:
        return None

    chart_data = CategoryChartData()
    chart_data.categories = labels
    chart_data.add_series("数据", values)

    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.PIE, left, top, width, height, chart_data
    )
    chart = chart_frame.chart
    colors = _get_colors(color_scheme, len(labels))

    _style_chart_common(chart, data.get("title", ""))

    # 图例
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False

    # 各扇区颜色
    series = chart.series[0]
    for i in range(len(labels)):
        point = series.points[i]
        point.format.fill.solid()
        point.format.fill.fore_color.rgb = RGBColor.from_string(colors[i])

    # 数据标签：显示百分比
    series.has_data_labels = True
    data_labels = series.data_labels
    data_labels.font.size = Pt(10)
    data_labels.font.name = _ZH_FONT
    data_labels.number_format = "0.0%"
    data_labels.show_percentage = True
    data_labels.show_value = False
    data_labels.show_category_name = True
    data_labels.show_series_name = False
    data_labels.label_position = XL_LABEL_POSITION.OUTSIDE_END

    return chart_frame


def _render_donut(slide, data: dict, color_scheme: dict,
                  left, top, width, height):
    """环形图 — DOUGHNUT。"""
    labels = data.get("labels", [])
    values = data.get("values", [])
    if not labels or not values:
        return None

    chart_data = CategoryChartData()
    chart_data.categories = labels
    chart_data.add_series("数据", values)

    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.DOUGHNUT, left, top, width, height, chart_data
    )
    chart = chart_frame.chart
    colors = _get_colors(color_scheme, len(labels))

    _style_chart_common(chart, data.get("title", ""))

    # 图例
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False

    # 各扇区颜色
    series = chart.series[0]
    for i in range(len(labels)):
        point = series.points[i]
        point.format.fill.solid()
        point.format.fill.fore_color.rgb = RGBColor.from_string(colors[i])

    # 数据标签
    series.has_data_labels = True
    data_labels = series.data_labels
    data_labels.font.size = Pt(10)
    data_labels.font.name = _ZH_FONT
    data_labels.number_format = "0.0%"
    data_labels.show_percentage = True
    data_labels.show_value = False
    data_labels.show_category_name = True
    data_labels.show_series_name = False

    return chart_frame


def _render_radar(slide, data: dict, color_scheme: dict,
                  left, top, width, height):
    """雷达图 — RADAR_FILLED。"""
    categories = data.get("categories", [])
    series_data = data.get("series", {})
    if not categories or not series_data:
        return None

    chart_data = CategoryChartData()
    chart_data.categories = categories
    for name, values in series_data.items():
        chart_data.add_series(name, values)

    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.RADAR_FILLED, left, top, width, height, chart_data
    )
    chart = chart_frame.chart
    colors = _get_colors(color_scheme, len(series_data))

    _style_chart_common(chart, data.get("title", ""))

    # 图例
    if len(series_data) > 1:
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
    else:
        chart.has_legend = False

    # 各系列颜色（半透明填充）
    for i, series in enumerate(chart.series):
        _apply_series_color(series, colors[i])
        series.format.fill.solid()
        series.format.fill.fore_color.rgb = RGBColor.from_string(colors[i])
        # 线条颜色
        series.format.line.color.rgb = RGBColor.from_string(colors[i])
        series.format.line.width = Pt(1.5)

    return chart_frame


def _render_scatter(slide, data: dict, color_scheme: dict,
                    left, top, width, height):
    """散点图 — XY_SCATTER。"""
    x_values = data.get("x", [])
    y_values = data.get("y", [])
    if not x_values or not y_values:
        return None

    chart_data = XyChartData()
    series_obj = chart_data.add_series("数据")
    for x, y in zip(x_values, y_values):
        series_obj.add_data_point(x, y)

    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.XY_SCATTER, left, top, width, height, chart_data
    )
    chart = chart_frame.chart
    colors = _get_colors(color_scheme, 1)

    _style_chart_common(chart, data.get("title", ""))
    _style_axis_fonts(chart)
    chart.has_legend = False

    # 系列样式 — 散点无线条
    series = chart.series[0]
    series.format.line.fill.background()
    _apply_series_color(series, colors[0])

    # 数据标签（如果有 labels）
    point_labels = data.get("labels", [])
    if point_labels:
        series.has_data_labels = True
        data_labels = series.data_labels
        data_labels.font.size = Pt(8)
        data_labels.font.name = _ZH_FONT
        data_labels.show_value = False
        data_labels.show_category_name = True
        data_labels.show_series_name = False

    return chart_frame


# ============ 用形状模拟的图表 ============

def _render_waterfall(slide, data: dict, color_scheme: dict,
                      left, top, width, height):
    """瀑布图 — 清晰的增减可视化，带连接线和图例。"""
    labels = data.get("labels", [])
    values = data.get("values", [])
    if not labels or not values:
        return None

    n = len(labels)
    primary = color_scheme.get("primary", "#1B365D")
    accent = color_scheme.get("accent", "#E8612D")
    secondary = color_scheme.get("secondary", "#4A90D9")
    positive_color = _rgb("#2ECC71")   # 绿色 = 增长
    negative_color = _rgb(accent)       # 橙红 = 下降
    total_color = _rgb(primary)         # 深蓝 = 基准/总计

    # 计算累积值
    cumulative = [0]
    for v in values[:-1]:
        cumulative.append(cumulative[-1] + v)
    cumulative.append(0)  # 最后一项从 0 起

    # 找出绘图范围
    all_vals = []
    for i, v in enumerate(values):
        base = cumulative[i]
        is_last = (i == n - 1)
        if is_last:
            all_vals.extend([0, v])
        else:
            all_vals.extend([base, base + v])

    max_val = max(all_vals) if all_vals else 1
    min_val = min(min(all_vals), 0)
    padding = (max_val - min_val) * 0.1
    max_val += padding
    val_range = max_val - min_val if max_val != min_val else 1

    left_emu = int(left)
    top_emu = int(top)
    width_emu = int(width)
    height_emu = int(height)

    # 布局：上方标题，中间图表，下方标签 + 图例
    title = data.get("title", "")
    title_h = Pt(28) if title else 0
    label_h = Pt(40)
    legend_h = Pt(24)
    chart_top = top_emu + int(title_h)
    chart_height = height_emu - int(title_h) - int(label_h) - int(legend_h)
    label_top = chart_top + chart_height + Pt(4)

    # 柱子尺寸
    margin_left = width_emu // 10
    chart_w = width_emu - margin_left * 2
    bar_width = chart_w // n
    gap = bar_width // 4
    actual_bar_width = bar_width - gap

    # 标题
    if title:
        txbox = slide.shapes.add_textbox(left_emu, top_emu, width_emu, int(title_h))
        tf = txbox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.name = _ZH_FONT
        p.alignment = PP_ALIGN.CENTER

    shapes_added = []

    def val_to_y(v):
        return chart_top + int(chart_height * (1 - (v - min_val) / val_range))

    for i, (label, val) in enumerate(zip(labels, values)):
        base = cumulative[i]
        is_first = (i == 0)
        is_last = (i == n - 1)

        # 确定颜色和柱子范围
        if is_first or is_last:
            color = total_color
            bar_bottom_val = 0
            bar_top_val = val
        elif val >= 0:
            color = positive_color
            bar_bottom_val = base
            bar_top_val = base + val
        else:
            color = negative_color
            bar_bottom_val = base + val
            bar_top_val = base

        bar_y_top = val_to_y(bar_top_val)
        bar_y_bottom = val_to_y(bar_bottom_val)
        bar_h = max(bar_y_bottom - bar_y_top, Pt(4))

        bar_x = left_emu + margin_left + i * bar_width + gap // 2

        # 绘制圆角矩形
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, bar_x, bar_y_top, actual_bar_width, bar_h
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        shapes_added.append(shape)

        # 数值标签（柱顶上方）
        val_text = f"{val:+d}" if not (is_first or is_last) else str(val)
        val_box = slide.shapes.add_textbox(
            bar_x - gap, bar_y_top - Pt(22), actual_bar_width + gap * 2, Pt(20)
        )
        vp = val_box.text_frame.paragraphs[0]
        vp.text = val_text
        vp.font.size = Pt(13)
        vp.font.bold = True
        vp.font.name = _ZH_FONT
        vp.font.color.rgb = color
        vp.alignment = PP_ALIGN.CENTER

        # 底部标签（多行支持）
        lbl_box = slide.shapes.add_textbox(
            bar_x - gap, label_top, actual_bar_width + gap * 2, int(label_h)
        )
        lbl_tf = lbl_box.text_frame
        lbl_tf.word_wrap = True
        lp = lbl_tf.paragraphs[0]
        lp.text = label
        lp.font.size = Pt(11)
        lp.font.name = _ZH_FONT
        lp.font.color.rgb = RGBColor(0x55, 0x55, 0x55)
        lp.alignment = PP_ALIGN.CENTER

        # 连接虚线（非首尾柱）
        if not is_last and i < n - 1:
            next_base = cumulative[i + 1] if i + 1 < len(cumulative) else 0
            connect_val = base + val if not (is_first or is_last) else val
            connect_y = val_to_y(connect_val)
            next_bar_x = left_emu + margin_left + (i + 1) * bar_width + gap // 2
            line = slide.shapes.add_connector(
                1,  # straight connector
                bar_x + actual_bar_width, connect_y,
                next_bar_x, connect_y,
            )
            line.line.color.rgb = RGBColor(0xBB, 0xBB, 0xBB)
            line.line.dash_style = 2  # dash
            line.line.width = Pt(1)

    # 图例
    legend_y = label_top + int(label_h) + Pt(2)
    legend_items = [
        (total_color, "基准/结果"),
        (positive_color, "增长"),
        (negative_color, "下降"),
    ]
    legend_x = left_emu + width_emu // 2 - Pt(120)
    for j, (lcolor, ltext) in enumerate(legend_items):
        lx = legend_x + j * Pt(90)
        dot = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     lx, legend_y, Pt(12), Pt(12))
        dot.fill.solid()
        dot.fill.fore_color.rgb = lcolor
        dot.line.fill.background()
        ltb = slide.shapes.add_textbox(lx + Pt(16), legend_y - Pt(2), Pt(70), Pt(16))
        ltp = ltb.text_frame.paragraphs[0]
        ltp.text = ltext
        ltp.font.size = Pt(10)
        ltp.font.name = _ZH_FONT
        ltp.font.color.rgb = RGBColor(0x66, 0x66, 0x66)

    return shapes_added


def _render_funnel(slide, data: dict, color_scheme: dict,
                   left, top, width, height):
    """漏斗图 — 使用水平矩形模拟，宽度递减。"""
    labels = data.get("labels", [])
    values = data.get("values", [])
    if not labels or not values:
        return None

    n = len(labels)
    colors = _get_colors(color_scheme, n)
    max_val = max(values) if values else 1

    left_emu = int(left)
    top_emu = int(top)
    width_emu = int(width)
    height_emu = int(height)

    # 标题
    title = data.get("title", "")
    title_height = Pt(24) if title else 0
    if title:
        txbox = slide.shapes.add_textbox(left_emu, top_emu, width_emu, int(title_height))
        tf = txbox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.name = _ZH_FONT
        p.alignment = PP_ALIGN.CENTER

    chart_top = top_emu + int(title_height)
    chart_height = height_emu - int(title_height)

    row_height = chart_height // n
    gap = row_height // 8

    shapes_added = []
    for i, (label, val) in enumerate(zip(labels, values)):
        ratio = val / max_val if max_val else 0
        bar_w = int(width_emu * ratio * 0.9)
        bar_x = left_emu + (width_emu - bar_w) // 2
        bar_y = chart_top + i * row_height + gap
        bar_h = row_height - gap * 2

        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, bar_x, bar_y, bar_w, bar_h
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor.from_string(colors[i])
        shape.line.fill.background()

        # 在矩形内显示标签和数值
        shape.text_frame.word_wrap = True
        p = shape.text_frame.paragraphs[0]
        p.text = f"{label}  {val}"
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.name = _ZH_FONT
        p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        p.alignment = PP_ALIGN.CENTER

        shapes_added.append(shape)

    return shapes_added


def _render_heatmap(slide, data: dict, color_scheme: dict,
                    left, top, width, height):
    """热力图 — 使用彩色矩形网格模拟。"""
    matrix = data.get("matrix", [])
    x_labels = data.get("x_labels", [])
    y_labels = data.get("y_labels", [])
    if not matrix:
        return None

    rows = len(matrix)
    cols = len(matrix[0]) if matrix else 0
    if rows == 0 or cols == 0:
        return None

    left_emu = int(left)
    top_emu = int(top)
    width_emu = int(width)
    height_emu = int(height)

    # 标题
    title = data.get("title", "")
    title_height = Pt(24) if title else 0
    if title:
        txbox = slide.shapes.add_textbox(left_emu, top_emu, width_emu, int(title_height))
        tf = txbox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.name = _ZH_FONT
        p.alignment = PP_ALIGN.CENTER

    chart_top = top_emu + int(title_height)
    chart_height = height_emu - int(title_height)

    # 预留行/列标签空间
    label_w = width_emu // (cols + 2)  # Y 标签列宽
    label_h = chart_height // (rows + 2)  # X 标签行高

    cell_w = (width_emu - label_w) // max(cols, 1)
    cell_h = (chart_height - label_h) // max(rows, 1)
    grid_left = left_emu + label_w
    grid_top = chart_top + label_h

    # 找出值范围，用于颜色映射
    flat_vals = [v for row in matrix for v in row]
    min_v = min(flat_vals) if flat_vals else 0
    max_v = max(flat_vals) if flat_vals else 1
    val_range = max_v - min_v if max_v != min_v else 1

    # 基础颜色（从浅到深）
    base_rgb = _rgb(color_scheme.get("primary", "#1B365D"))
    base_r, base_g, base_b = base_rgb[0], base_rgb[1], base_rgb[2]

    shapes_added = []

    # X 标签（顶部）
    for j, xlabel in enumerate(x_labels[:cols]):
        bx = grid_left + j * cell_w
        txbox = slide.shapes.add_textbox(bx, chart_top, cell_w, label_h)
        p = txbox.text_frame.paragraphs[0]
        p.text = str(xlabel)
        p.font.size = Pt(8)
        p.font.name = _ZH_FONT
        p.alignment = PP_ALIGN.CENTER

    # Y 标签（左侧）
    for i, ylabel in enumerate(y_labels[:rows]):
        by = grid_top + i * cell_h
        txbox = slide.shapes.add_textbox(left_emu, by, label_w, cell_h)
        p = txbox.text_frame.paragraphs[0]
        p.text = str(ylabel)
        p.font.size = Pt(8)
        p.font.name = _ZH_FONT
        p.alignment = PP_ALIGN.RIGHT

    # 绘制单元格
    for i, row in enumerate(matrix):
        for j, val in enumerate(row):
            cx = grid_left + j * cell_w
            cy = grid_top + i * cell_h

            # 颜色插值：从白色到基础颜色
            t = (val - min_v) / val_range
            r = int(255 - (255 - base_r) * t)
            g = int(255 - (255 - base_g) * t)
            b = int(255 - (255 - base_b) * t)

            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, cx, cy, cell_w, cell_h
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(r, g, b)
            shape.line.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            shape.line.width = Pt(0.5)

            # 单元格内显示数值
            shape.text_frame.word_wrap = False
            p = shape.text_frame.paragraphs[0]
            p.text = f"{val:.1f}" if isinstance(val, float) else str(val)
            p.font.size = Pt(8)
            p.font.name = _ZH_FONT
            p.alignment = PP_ALIGN.CENTER
            # 深色背景用白字，浅色用黑字
            brightness = (r * 299 + g * 587 + b * 114) / 1000
            p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF) if brightness < 128 else RGBColor(0, 0, 0)

            shapes_added.append(shape)

    return shapes_added


# ============ 渲染器注册表 ============

_CHART_RENDERERS = {
    "bar": _render_bar,
    "grouped_bar": _render_grouped_bar,
    "line": _render_line,
    "pie": _render_pie,
    "donut": _render_donut,
    "radar": _render_radar,
    "scatter": _render_scatter,
    "waterfall": _render_waterfall,
    "funnel": _render_funnel,
    "heatmap": _render_heatmap,
}
