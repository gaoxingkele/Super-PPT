# -*- coding: utf-8 -*-
"""
参考模板风格提取器。
从 PPTX/PDF/图片中提取配色方案、字体、设计语言等风格信息。
"""
import json
from pathlib import Path

import src  # noqa: F401


def extract_style(template_path: Path) -> dict:
    """
    从参考模板提取风格信息。

    Args:
        template_path: 模板文件路径（.pptx / .pdf / .png / .jpg）

    Returns:
        style_profile dict
    """
    suffix = template_path.suffix.lower()

    if suffix == ".pptx":
        return _extract_from_pptx(template_path)
    elif suffix == ".pdf":
        return _extract_from_pdf(template_path)
    elif suffix in (".png", ".jpg", ".jpeg", ".webp"):
        return _extract_from_image(template_path)
    else:
        raise ValueError(f"不支持的模板格式: {suffix}")


def _extract_from_pptx(path: Path) -> dict:
    """从 PPTX 模板提取风格。"""
    from pptx import Presentation
    from pptx.util import Pt

    prs = Presentation(str(path))

    # 提取配色
    color_scheme = {}
    try:
        theme = prs.slide_masters[0].slide_layouts[0]
        # 从母版的颜色主题提取
    except Exception:
        pass

    # 提取字体和布局信息
    fonts = set()
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    if para.font and para.font.name:
                        fonts.add(para.font.name)

    # 提取布局类型
    layout_names = []
    for layout in prs.slide_layouts:
        layout_names.append(layout.name)

    profile = {
        "source": str(path),
        "color_scheme": color_scheme or {
            "primary": "#1B365D",
            "secondary": "#4A90D9",
            "accent": "#E8612D",
            "background": "#FFFFFF",
            "text": "#333333",
        },
        "typography": {
            "fonts_detected": list(fonts),
            "title_font": list(fonts)[0] if fonts else "微软雅黑",
            "body_font": list(fonts)[-1] if fonts else "微软雅黑",
        },
        "layout_style": {
            "visual_weight": "moderate",
            "decoration": "minimal",
        },
        "extracted_layouts": layout_names,
        "slide_count": len(prs.slides),
    }
    return profile


def _extract_from_pdf(path: Path) -> dict:
    """从 PDF 模板提取风格（使用 LLM 视觉分析）。"""
    # Phase 1: 基础占位
    # TODO Phase 3: pdf2image 截图 + chat_vision 分析
    return {
        "source": str(path),
        "color_scheme": {
            "primary": "#1B365D",
            "secondary": "#4A90D9",
            "accent": "#E8612D",
            "background": "#FFFFFF",
            "text": "#333333",
        },
        "design_language": "auto-detected-from-pdf",
        "note": "PDF 风格提取需要 Phase 3 实现",
    }


def _extract_from_image(path: Path) -> dict:
    """从图片模板提取风格（使用 LLM 视觉分析）。"""
    # Phase 1: 基础占位
    # TODO Phase 3: chat_vision 直接分析图片
    return {
        "source": str(path),
        "color_scheme": {
            "primary": "#1B365D",
            "secondary": "#4A90D9",
            "accent": "#E8612D",
            "background": "#FFFFFF",
            "text": "#333333",
        },
        "design_language": "auto-detected-from-image",
        "note": "图片风格提取需要 Phase 3 实现",
    }
