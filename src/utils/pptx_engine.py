# -*- coding: utf-8 -*-
"""
python-pptx 装配引擎。
将幻灯片大纲 + 视觉资产 + 主题模板装配为最终 PPTX。
"""
from pathlib import Path
from typing import Optional

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE


CJK_FONT = "微软雅黑"
EN_FONT = "Times New Roman"

# ── CJK 文字长度自适应（借鉴 PPTAgent length_factor） ──
# 不同布局区域的字符容量上限（按区域宽度和字号估算）
# 格式: (宽度inches, 字号pt) → 每行约容纳的中文字数
# 中文字符宽度 ≈ 字号pt × 0.0139 英寸，英文/数字约一半
_LAYOUT_CHAR_LIMITS = {
    # (区域宽度, 字号) → 单行字数, 最大行数 → 每条bullet最大字数
    "title_content":    {"max_chars_per_bullet": 45, "max_bullets": 6},
    "data_chart":       {"max_chars_per_bullet": 35, "max_bullets": 5},
    "infographic":      {"max_chars_per_bullet": 35, "max_bullets": 5},
    "two_column":       {"max_chars_per_bullet": 30, "max_bullets": 5},
    "key_insight":      {"max_chars_per_bullet": 40, "max_bullets": 5},
    "methodology":      {"max_chars_per_bullet": 30, "max_bullets": 6},
    "architecture":     {"max_chars_per_bullet": 30, "max_bullets": 6},
    "summary":          {"max_chars_per_bullet": 40, "max_bullets": 6},
    "table":            {"max_chars_per_bullet": 45, "max_bullets": 5},
    "timeline":         {"max_chars_per_bullet": 35, "max_bullets": 5},
    "agenda":           {"max_chars_per_bullet": 50, "max_bullets": 10},
}


def _adapt_bullet_text(text: str, layout: str, is_title: bool = False) -> str:
    """
    根据布局类型自适应截断 bullet 文字，防止溢出文本框。

    借鉴 PPTAgent 的 length_factor 思路：
    - CJK 字符按 1.0 计宽
    - 拉丁字母/数字按 0.5 计宽
    - 超出容量时智能截断到最近的标点/空格处
    """
    if is_title:
        # 标题最大约25个中文字宽度
        max_width = 30
    else:
        limits = _LAYOUT_CHAR_LIMITS.get(layout, {"max_chars_per_bullet": 45})
        max_width = limits["max_chars_per_bullet"]

    # 计算等效中文字宽度
    width = 0.0
    for ch in text:
        if '\u4e00' <= ch <= '\u9fff' or '\u3000' <= ch <= '\u303f' or '\uff00' <= ch <= '\uffef':
            width += 1.0  # CJK 全角
        else:
            width += 0.55  # 拉丁/数字/标点（半角）

    if width <= max_width:
        return text

    # 超长，需要截断
    # 找到 max_width 对应的字符位置
    acc = 0.0
    cut_pos = len(text)
    for i, ch in enumerate(text):
        if '\u4e00' <= ch <= '\u9fff' or '\u3000' <= ch <= '\u303f' or '\uff00' <= ch <= '\uffef':
            acc += 1.0
        else:
            acc += 0.55
        if acc > max_width - 2:  # 留2字宽给省略号
            cut_pos = i
            break

    # 回退到最近的标点/空格处（更自然的断句）
    truncated = text[:cut_pos]
    for sep in ('，', '；', '、', '。', ',', ';', ' ', '）', ')'):
        last_sep = truncated.rfind(sep)
        if last_sep > cut_pos * 0.6:  # 不要回退太多
            truncated = truncated[:last_sep + 1]
            break

    return truncated.rstrip() + "…"


def _adapt_bullets_for_layout(bullets: list, layout: str) -> list:
    """对整个 bullet 列表做布局自适应：截断过长条目，裁剪过多条目。"""
    limits = _LAYOUT_CHAR_LIMITS.get(layout, {"max_chars_per_bullet": 45, "max_bullets": 6})
    max_bullets = limits["max_bullets"]

    # 裁剪过多的 bullets
    if len(bullets) > max_bullets:
        bullets = bullets[:max_bullets]

    # 逐条截断
    return [_adapt_bullet_text(str(b), layout) for b in bullets]


def _set_font_name(run, text: str = ""):
    """为 run 设置字体名称：中文用微软雅黑，英文/数字用 Times New Roman。"""
    run.font.name = CJK_FONT
    try:
        from pptx.oxml.ns import qn
        r_elem = run._r
        rPr = r_elem.find(qn("a:rPr"))
        if rPr is None:
            from lxml import etree
            rPr = etree.SubElement(r_elem, qn("a:rPr"))
        # 设置东亚字体
        ea = rPr.find(qn("a:ea"))
        if ea is None:
            from lxml import etree
            ea = etree.SubElement(rPr, qn("a:ea"))
        ea.set("typeface", CJK_FONT)
        # 设置拉丁字体
        latin = rPr.find(qn("a:latin"))
        if latin is None:
            from lxml import etree
            latin = etree.SubElement(rPr, qn("a:latin"))
        latin.set("typeface", EN_FONT)
    except Exception:
        pass


class PPTXBuilder:
    """PPTX 装配器。"""

    def __init__(self, template_path: Optional[Path] = None, color_scheme: dict = None,
                 assets_dir: Optional[Path] = None, disable_auto_split: bool = False):
        if template_path and template_path.is_file():
            self.prs = Presentation(str(template_path))
        else:
            self.prs = Presentation()
            self.prs.slide_width = Inches(13.333)
            self.prs.slide_height = Inches(7.5)

        self.colors = color_scheme or {}
        self.assets_dir = assets_dir
        self._slide_count = 0
        self._disable_auto_split = disable_auto_split

    def add_slide(self, slide_spec: dict, asset_path: Optional[Path] = None):
        """添加一张幻灯片。"""
        layout_type = slide_spec.get("layout", "title_content")

        # CJK 文字长度自适应：渲染前截断过长/过多的 bullets
        if "bullets" in slide_spec and isinstance(slide_spec["bullets"], list):
            slide_spec["bullets"] = _adapt_bullets_for_layout(
                slide_spec["bullets"], layout_type
            )

        handler = _LAYOUT_HANDLERS.get(layout_type, _add_title_content)
        handler(self, slide_spec, asset_path)
        self._slide_count += 1

    def save(self, output_path: Path):
        output_path.parent.mkdir(parents=True, exist_ok=True)
        self.prs.save(str(output_path))

    def _get_color(self, key: str, default: str = "#333333") -> RGBColor:
        hex_color = self.colors.get(key, default).lstrip("#")
        return RGBColor.from_string(hex_color)

    def _add_blank_slide(self):
        layout = self.prs.slide_layouts[6]  # 空白布局
        return self.prs.slides.add_slide(layout)


# ============ 辅助函数 ============

def _get_image_size(image_path: str):
    """获取图片原始尺寸，返回 (width, height) 像素；失败返回 None。"""
    try:
        from PIL import Image as PILImage
        img = PILImage.open(image_path)
        size = img.size
        img.close()
        return size
    except Exception:
        return None


def _image_fill_ratio(image_path: str, area_width, area_height) -> float:
    """计算图片在指定区域内等比缩放后的面积利用率 (0~1)。"""
    size = _get_image_size(image_path)
    if not size:
        return 1.0  # 无法判断时假设填满
    img_w, img_h = size
    mw, mh = int(area_width), int(area_height)
    scale = min(mw / img_w, mh / img_h)
    return (img_w * scale * img_h * scale) / (mw * mh)


def _add_picture_fit(slide, image_path: str, left, top, max_width, max_height):
    """在指定区域内等比缩放插入图片，保持原始比例不变形，居中放置。"""
    size = _get_image_size(image_path)
    if not size:
        return slide.shapes.add_picture(str(image_path), left, top, max_width, max_height)

    img_w, img_h = size
    mw = int(max_width)
    mh = int(max_height)
    scale = min(mw / img_w, mh / img_h)

    final_w = int(img_w * scale)
    final_h = int(img_h * scale)

    # 居中偏移
    offset_x = (mw - final_w) // 2
    offset_y = (mh - final_h) // 2

    return slide.shapes.add_picture(
        str(image_path),
        int(left) + offset_x,
        int(top) + offset_y,
        final_w,
        final_h,
    )


def _add_split_image_text(builder, spec: dict, asset_path: Path,
                           image_area_width=None, image_area_height=None):
    """
    拆页辅助：当图文共存但图片面积利用率过低时，拆为两页。
    - Page 1: 标题 + 全幅图片（大面积展示）
    - Page 2: 标题 + bullets + takeaway（文字详解）
    返回 True 表示已拆页处理，False 表示不需要拆页。
    """
    if not asset_path or not asset_path.is_file():
        return False
    bullets = spec.get("bullets", [])
    if len(bullets) < 2:
        return False

    # 计算如果共存时图片的面积利用率
    aw = image_area_width or Inches(7.5)
    ah = image_area_height or Inches(5)
    fill = _image_fill_ratio(str(asset_path), aw, ah)

    # 面积利用率 > 45% 表示图片与区域比例相近，无需拆分
    if fill > 0.45:
        return False

    # ---- Page 1: 全幅图片页 ----
    slide1 = builder._add_blank_slide()
    _add_title_textbox(slide1, spec.get("title", ""), builder)
    _add_bottom_bar(slide1, builder)

    # 图片占据主要内容区域（全宽）
    _add_picture_fit(slide1, str(asset_path),
                     Inches(0.5), Inches(1.6), Inches(12.3), Inches(5.3))

    # 底部 takeaway（如果有）
    takeaway = spec.get("takeaway", "")
    if takeaway:
        txBox = slide1.shapes.add_textbox(Inches(0.8), Inches(6.9), Inches(11), Inches(0.5))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = f"\u25b6 {takeaway}"
        run.font.size = Pt(14)
        run.font.bold = True
        run.font.color.rgb = builder._get_color("accent", "#C00000")
        _set_font_name(run)

    if spec.get("notes"):
        slide1.notes_slide.notes_text_frame.text = spec["notes"]
    _add_page_number(slide1, builder, builder._slide_count + 1)
    builder._slide_count += 1

    # ---- Page 2: 文字详解页 ----
    slide2 = builder._add_blank_slide()
    detail_title = spec.get("title", "") + " — 要点解读"
    _add_title_textbox(slide2, detail_title, builder)
    _add_bottom_bar(slide2, builder)

    # 副标题
    subtitle = spec.get("subtitle", "")
    if subtitle:
        txBox = slide2.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(12), Inches(0.4))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(15)
        p.font.color.rgb = RGBColor(0x88, 0x88, 0x88)

    # Bullets（大面积展示，双栏排列）
    n = len(bullets)
    if n <= 4:
        # 单栏大字
        txBox = slide2.shapes.add_textbox(Inches(0.8), Inches(1.9), Inches(11.5), Inches(4.5))
        tf = txBox.text_frame
        tf.word_wrap = True
        for i, bullet in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            _add_rich_text(p, f"\u2022 {bullet}", builder)
            p.font.size = Pt(18)
            p.space_after = Pt(14)
            p.line_spacing = Pt(30)
    else:
        # 双栏排列
        mid = (n + 1) // 2
        for col, (col_bullets, x_start) in enumerate([
            (bullets[:mid], Inches(0.8)),
            (bullets[mid:], Inches(7.0)),
        ]):
            txBox = slide2.shapes.add_textbox(x_start, Inches(1.9), Inches(5.8), Inches(4.5))
            tf = txBox.text_frame
            tf.word_wrap = True
            for i, bullet in enumerate(col_bullets):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                _add_rich_text(p, f"\u2022 {bullet}", builder)
                p.font.size = Pt(16)
                p.space_after = Pt(10)
                p.line_spacing = Pt(26)

    # Takeaway
    if takeaway:
        txBox = slide2.shapes.add_textbox(Inches(0.8), Inches(6.5), Inches(11), Inches(0.7))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = f"\u25b6 {takeaway}"
        run.font.size = Pt(14)
        run.font.bold = True
        run.font.color.rgb = builder._get_color("accent", "#C00000")
        _set_font_name(run)

    _add_page_number(slide2, builder, builder._slide_count + 1)
    # 注意：外层 add_slide 会再 +1，所以这里不额外加

    return True


def _set_shape_transparency(shape, alpha_val: int):
    """设置形状填充的透明度。alpha_val: 0-100000 (0=全透明, 100000=不透明)。"""
    try:
        from pptx.oxml.ns import qn
        from lxml import etree
        # 访问形状的 spPr/a:solidFill/a:srgbClr 节点
        sp = shape._element
        spPr = sp.find(qn("p:spPr"))
        if spPr is None:
            return
        solidFill = spPr.find(qn("a:solidFill"))
        if solidFill is None:
            return
        color_elem = solidFill[0] if len(solidFill) else None
        if color_elem is not None:
            # 移除已有 alpha
            for old_alpha in color_elem.findall(qn("a:alpha")):
                color_elem.remove(old_alpha)
            alpha = etree.SubElement(color_elem, qn("a:alpha"))
            alpha.set("val", str(alpha_val))
    except Exception:
        pass  # 静默失败，不影响主流程


def _add_title_textbox(slide, title: str, builder: PPTXBuilder):
    """在幻灯片顶部添加标题 + 底部装饰线。"""
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = builder._get_color("primary", "#002060")
    _set_font_name(run)

    # 标题下方装饰线
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.5), Inches(1.25), Inches(2), Pt(4),
    )
    line.fill.solid()
    line.fill.fore_color.rgb = builder._get_color("accent", "#C00000")
    line.line.fill.background()


def _add_page_number(slide, builder: PPTXBuilder, page_num: int):
    """在右下角添加页码。"""
    txBox = slide.shapes.add_textbox(
        Inches(12.2), Inches(7.0), Inches(1), Inches(0.4),
    )
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = str(page_num)
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(0x99, 0x99, 0x99)
    p.alignment = PP_ALIGN.RIGHT


def _add_bottom_bar(slide, builder: PPTXBuilder):
    """底部品牌色条。"""
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(7.35), builder.prs.slide_width, Pt(8),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = builder._get_color("primary", "#002060")
    bar.line.fill.background()


def _add_rich_text(paragraph, text: str, builder: PPTXBuilder):
    """解析 **bold** 标记，将加粗部分渲染为深红色粗体。"""
    import re
    # Clear default text
    paragraph.clear()
    parts = re.split(r'(\*\*.*?\*\*)', text)
    for part in parts:
        if part.startswith('**') and part.endswith('**'):
            run = paragraph.add_run()
            run.text = part[2:-2]
            run.font.bold = True
            run.font.color.rgb = builder._get_color("accent", "#C00000")
            run.font.size = Pt(17)
            _set_font_name(run)
        else:
            if part:
                run = paragraph.add_run()
                run.text = part
                run.font.color.rgb = builder._get_color("text", "#333333")
                run.font.size = Pt(17)
                _set_font_name(run)


# ============ 纯文字页分块卡片 ============

# 交替浅色背景色板
_BLOCK_COLORS = [
    RGBColor(0xE8, 0xF0, 0xFE),  # 浅蓝 #E8F0FE
    RGBColor(0xFF, 0xF3, 0xE0),  # 浅橙 #FFF3E0
    RGBColor(0xE8, 0xF5, 0xE9),  # 浅绿 #E8F5E9
    RGBColor(0xFC, 0xE4, 0xEC),  # 浅粉 #FCE4EC
    RGBColor(0xF3, 0xE5, 0xF5),  # 浅紫 #F3E5F5
    RGBColor(0xE0, 0xF7, 0xFA),  # 浅青 #E0F7FA
]


def _add_text_block_bullets(slide, bullets: list, builder):
    """纯文字页：每个要点用不同浅色背景卡片呈现，提升阅读体验。"""
    n = len(bullets)
    if n == 0:
        return

    # 布局参数
    area_top = Inches(1.7)
    area_left = Inches(0.5)
    area_width = Inches(12.3)
    area_height = Inches(4.8)
    gap = Inches(0.12)

    if n <= 3:
        # 单列：每个卡片全宽
        card_h = (area_height - gap * (n - 1)) / n
        for i, bullet in enumerate(bullets):
            y = area_top + i * (card_h + gap)
            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                area_left, y, area_width, card_h,
            )
            card.fill.solid()
            card.fill.fore_color.rgb = _BLOCK_COLORS[i % len(_BLOCK_COLORS)]
            card.line.fill.background()
            # 左侧色条
            bar = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                area_left, y, Inches(0.08), card_h,
            )
            bar.fill.solid()
            bar.fill.fore_color.rgb = builder._get_color("primary", "#002060")
            bar.line.fill.background()
            # 文字
            txBox = slide.shapes.add_textbox(
                area_left + Inches(0.25), y + Inches(0.1),
                area_width - Inches(0.4), card_h - Inches(0.2),
            )
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            _add_rich_text(p, str(bullet), builder)
            p.font.size = Pt(16)
            p.line_spacing = Pt(24)
    else:
        # 双列：2xN 网格
        cols = 2
        rows = (n + 1) // 2
        col_width = (area_width - gap) / 2
        card_h = (area_height - gap * (rows - 1)) / rows
        for i, bullet in enumerate(bullets):
            col = i % cols
            row = i // cols
            x = area_left + col * (col_width + gap)
            y = area_top + row * (card_h + gap)
            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                x, y, col_width, card_h,
            )
            card.fill.solid()
            card.fill.fore_color.rgb = _BLOCK_COLORS[i % len(_BLOCK_COLORS)]
            card.line.fill.background()
            # 顶部色条
            bar = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                x, y, col_width, Pt(4),
            )
            bar.fill.solid()
            bar.fill.fore_color.rgb = builder._get_color("secondary", "#0060A8")
            bar.line.fill.background()
            # 文字
            txBox = slide.shapes.add_textbox(
                x + Inches(0.15), y + Inches(0.12),
                col_width - Inches(0.3), card_h - Inches(0.2),
            )
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            _add_rich_text(p, str(bullet), builder)
            p.font.size = Pt(15)
            p.line_spacing = Pt(22)


# ============ 布局处理器 ============

def _add_cover(builder: PPTXBuilder, spec: dict, asset_path: Optional[Path]):
    """封面页：背景图 + 半透明遮罩 + 标题。"""
    slide = builder._add_blank_slide()

    # 背景图
    if asset_path and asset_path.is_file():
        slide.shapes.add_picture(
            str(asset_path), Inches(0), Inches(0),
            builder.prs.slide_width, builder.prs.slide_height,
        )

    # 半透明遮罩（深色矩形增加文字可读性）
    overlay = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), builder.prs.slide_width, builder.prs.slide_height,
    )
    overlay.fill.solid()
    overlay.fill.fore_color.rgb = builder._get_color("primary", "#002060")
    _set_shape_transparency(overlay, 40000)  # 40% 不透明
    overlay.line.fill.background()

    # 左侧装饰线
    accent_line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.8), Inches(2.2), Pt(6), Inches(2.8),
    )
    accent_line.fill.solid()
    accent_line.fill.fore_color.rgb = builder._get_color("accent", "#C00000")
    accent_line.line.fill.background()

    # 标题
    title = spec.get("title", "")
    txBox = slide.shapes.add_textbox(Inches(1.2), Inches(2.3), Inches(10.5), Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(44)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    _set_font_name(run)
    p.alignment = PP_ALIGN.LEFT

    # 副标题
    subtitle = spec.get("subtitle", "")
    if subtitle:
        p2 = tf.add_paragraph()
        run2 = p2.add_run()
        run2.text = subtitle
        run2.font.size = Pt(20)
        run2.font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
        _set_font_name(run2)
        p2.alignment = PP_ALIGN.LEFT
        p2.space_before = Pt(16)


def _add_agenda(builder: PPTXBuilder, spec: dict, asset_path: Optional[Path]):
    """目录页：图形化大纲列表。"""
    slide = builder._add_blank_slide()
    _add_title_textbox(slide, spec.get("title", "目录"), builder)
    _add_bottom_bar(slide, builder)

    # 副标题
    subtitle = spec.get("subtitle", "")
    if subtitle:
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(12), Inches(0.5))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(0x88, 0x88, 0x88)
        p.font.italic = True

    items = spec.get("bullets", spec.get("items", []))
    n = len(items)
    start_y = 2.0
    item_height = min(1.0, 4.5 / max(n, 1))

    for i, item in enumerate(items):
        y = Inches(start_y) + i * Inches(item_height)

        # 背景卡片
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(1.2), y - Inches(0.05), Inches(10.5), Inches(item_height - 0.1),
        )
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(0xF5, 0xF7, 0xFA)
        card.line.fill.background()

        # 左侧色条
        color_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(1.2), y - Inches(0.05), Pt(5), Inches(item_height - 0.1),
        )
        color_bar.fill.solid()
        color_bar.fill.fore_color.rgb = builder._get_color("accent", "#C00000")
        color_bar.line.fill.background()

        # 编号圆圈
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(1.6), y + Inches(0.05), Inches(0.45), Inches(0.45),
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = builder._get_color("primary", "#002060")
        circle.line.fill.background()
        txc = circle.text_frame
        txc.paragraphs[0].text = f"0{i + 1}" if i < 9 else str(i + 1)
        txc.paragraphs[0].font.size = Pt(14)
        txc.paragraphs[0].font.bold = True
        txc.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        txc.paragraphs[0].alignment = PP_ALIGN.CENTER
        txc.word_wrap = False

        # 项目文字
        txBox = slide.shapes.add_textbox(Inches(2.3), y + Inches(0.05), Inches(9), Inches(0.5))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = str(item)
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = builder._get_color("text", "#333333")

    if spec.get("notes"):
        slide.notes_slide.notes_text_frame.text = spec["notes"]
    _add_page_number(slide, builder, builder._slide_count + 1)


def _add_section_break(builder: PPTXBuilder, spec: dict, asset_path: Optional[Path]):
    """章节过渡页：背景图/色块 + 编号徽章 + 大标题。"""
    slide = builder._add_blank_slide()

    # 背景图（如果有）
    if asset_path and asset_path.is_file():
        slide.shapes.add_picture(
            str(asset_path), Inches(0), Inches(0),
            builder.prs.slide_width, builder.prs.slide_height,
        )
        # 加遮罩
        overlay = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0), builder.prs.slide_width, builder.prs.slide_height,
        )
        overlay.fill.solid()
        overlay.fill.fore_color.rgb = builder._get_color("primary", "#002060")
        _set_shape_transparency(overlay, 55000)
        overlay.line.fill.background()
    else:
        # 纯色背景
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = builder._get_color("primary", "#002060")

    # 编号徽章（大圆形）
    subtitle = spec.get("subtitle", "")
    # 从 subtitle 提取编号（如 "PART 01", "第一章" 等）
    import re
    badge_text = ""
    if subtitle:
        m = re.search(r'(\d+)', subtitle)
        if m:
            badge_text = f"{int(m.group(1)):02d}"
    if not badge_text:
        # 尝试从已有 slide 数计算
        section_num = sum(1 for _ in range(builder._slide_count)
                         if False) or ""

    if badge_text:
        # 大圆形编号徽章
        badge = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(1.5), Inches(1.8), Inches(1.6), Inches(1.6),
        )
        badge.fill.solid()
        badge.fill.fore_color.rgb = builder._get_color("accent", "#C00000")
        badge.line.fill.background()

        txb = badge.text_frame
        txb.paragraphs[0].alignment = PP_ALIGN.CENTER
        run = txb.paragraphs[0].add_run()
        run.text = badge_text
        run.font.size = Pt(56)
        run.font.bold = True
        run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        _set_font_name(run)
        txb.vertical_anchor = MSO_ANCHOR.MIDDLE

    # 左侧装饰线
    accent_line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(1.5), Inches(3.6), Inches(2), Pt(5),
    )
    accent_line.fill.solid()
    accent_line.fill.fore_color.rgb = builder._get_color("accent", "#C00000")
    accent_line.line.fill.background()

    # 大标题
    txBox = slide.shapes.add_textbox(Inches(1.5), Inches(3.8), Inches(10), Inches(2))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = spec.get("title", "")
    run.font.size = Pt(42)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    _set_font_name(run)
    p.alignment = PP_ALIGN.LEFT

    # 副标题
    if subtitle:
        p2 = tf.add_paragraph()
        run2 = p2.add_run()
        run2.text = subtitle
        run2.font.size = Pt(20)
        run2.font.color.rgb = RGBColor(0xBB, 0xBB, 0xBB)
        _set_font_name(run2)
        p2.space_before = Pt(12)


def _add_title_content(builder: PPTXBuilder, spec: dict, asset_path: Optional[Path]):
    """标题 + 副标题 + 要点 + 右侧图片。图文比例不合时自动拆页。"""
    if asset_path and spec.get("bullets") and not builder._disable_auto_split:
        if _add_split_image_text(builder, spec, asset_path,
                                  Inches(4.8), Inches(5)):
            return

    slide = builder._add_blank_slide()
    _add_title_textbox(slide, spec.get("title", ""), builder)
    _add_bottom_bar(slide, builder)

    # 副标题
    subtitle = spec.get("subtitle", "")
    if subtitle:
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(12), Inches(0.4))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(15)
        p.font.color.rgb = RGBColor(0x88, 0x88, 0x88)

    bullets = spec.get("bullets", [])
    has_image = asset_path and asset_path.is_file()

    # 原生图表渲染（title_content 也可能有 matplotlib）
    visual = spec.get("visual", {})
    chart_rendered = False
    if visual and isinstance(visual, dict) and visual.get("type") == "matplotlib":
        try:
            from src.visuals.pptx_charts import render_native_chart
            if bullets:
                result = render_native_chart(slide, visual, builder.colors,
                                    Inches(0.3), Inches(1.8), Inches(7.2), Inches(4.5))
            else:
                result = render_native_chart(slide, visual, builder.colors,
                                    Inches(1), Inches(1.6), Inches(11), Inches(5))
            chart_rendered = result is not None
        except Exception as e:
            print(f"  [WARN] title_content 原生图表渲染失败 ({spec.get('id','')}): {e}", flush=True)

    if bullets and not has_image and not chart_rendered:
        # ---- 纯文字页：分块背景卡片布局 ----
        _add_text_block_bullets(slide, bullets, builder)
    elif bullets:
        # ---- 图文页/图表页：左文右图 ----
        left_width = Inches(6.8)
        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.9), left_width, Inches(4.3))
        tf = txBox.text_frame
        tf.word_wrap = True
        for i, bullet in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            _add_rich_text(p, f"\u2022 {bullet}", builder)
            p.font.size = Pt(17)
            p.space_after = Pt(10)
            p.line_spacing = Pt(28)

    if has_image:
        _add_picture_fit(slide, str(asset_path), Inches(8), Inches(1.6), Inches(4.8), Inches(5))

    takeaway = spec.get("takeaway", "")
    if takeaway:
        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(6.5), Inches(11), Inches(0.7))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = f"\u25b6 {takeaway}"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = builder._get_color("accent", "#C00000")

    if spec.get("notes"):
        slide.notes_slide.notes_text_frame.text = spec["notes"]
    _add_page_number(slide, builder, builder._slide_count + 1)


def _add_data_chart(builder: PPTXBuilder, spec: dict, asset_path: Optional[Path]):
    """数据图表页：左侧图表 + 右侧要点说明。支持透明产品背景图增强视觉。"""
    slide = builder._add_blank_slide()

    # 产品背景图：局部放置在右下角，不干扰前景内容
    slide_id = spec.get("id", "")
    has_bg = False
    if builder.assets_dir and slide_id:
        bg_candidate = builder.assets_dir / f"{slide_id}_bg.png"
        if bg_candidate.is_file():
            has_bg = True
            # 右下角局部放置产品图（不遮挡图表和文字区域）
            slide.shapes.add_picture(
                str(bg_candidate), Inches(8.5), Inches(3.5),
                Inches(5), Inches(4),
            )
            # 在产品图上方叠一层渐变白色遮罩，使其更柔和
            overlay = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(8.5), Inches(3.5), Inches(5), Inches(4),
            )
            overlay.fill.solid()
            overlay.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            _set_shape_transparency(overlay, 20000)  # 20% 不透明 = 80% 透明
            overlay.line.fill.background()

    _add_title_textbox(slide, spec.get("title", ""), builder)
    _add_bottom_bar(slide, builder)

    # 副标题
    subtitle = spec.get("subtitle", "")
    if subtitle:
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(12), Inches(0.4))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(0x88, 0x88, 0x88)

    bullets = spec.get("bullets", [])
    has_bullets = bool(bullets)

    # 渲染图表：优先原生渲染，回退到图片
    visual = spec.get("visual", {})
    chart_rendered = False
    if visual and isinstance(visual, dict) and visual.get("type") == "matplotlib":
        try:
            from src.visuals.pptx_charts import render_native_chart
            if has_bullets:
                result = render_native_chart(slide, visual, builder.colors,
                                    Inches(0.3), Inches(1.8), Inches(8.2), Inches(4.8))
            else:
                result = render_native_chart(slide, visual, builder.colors,
                                    Inches(1), Inches(1.6), Inches(11), Inches(5))
            chart_rendered = result is not None
        except Exception as e:
            print(f"  [WARN] data_chart 原生图表渲染失败 ({spec.get('id','')}): {e}", flush=True)

    if not chart_rendered and asset_path and asset_path.is_file():
        if has_bullets:
            _add_picture_fit(slide, str(asset_path), Inches(0.3), Inches(1.8), Inches(8.2), Inches(4.8))
        else:
            _add_picture_fit(slide, str(asset_path), Inches(1), Inches(1.6), Inches(11), Inches(5))

    if has_bullets:
        # 右侧要点说明区（带白色背景卡片，确保文字清晰可读）
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(8.8), Inches(1.8), Inches(4.2), Inches(4.8),
        )
        card.fill.solid()
        if has_bg:
            # 有背景图时用纯白底，确保文字完全清晰
            card.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            _set_shape_transparency(card, 90000)  # 90% 不透明
        else:
            card.fill.fore_color.rgb = RGBColor(0xF5, 0xF7, 0xFA)
        card.line.fill.background()

        # 要点标题
        txTitle = slide.shapes.add_textbox(Inches(9.0), Inches(1.9), Inches(3.8), Inches(0.4))
        tf_t = txTitle.text_frame
        p_t = tf_t.paragraphs[0]
        p_t.text = "关键解读"
        p_t.font.size = Pt(16)
        p_t.font.bold = True
        p_t.font.color.rgb = builder._get_color("primary", "#002060")

        # 要点列表（支持 **bold** 标记）
        txBox = slide.shapes.add_textbox(Inches(9.0), Inches(2.5), Inches(3.8), Inches(3.8))
        tf = txBox.text_frame
        tf.word_wrap = True
        for i, bullet in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            _add_rich_text(p, f"\u25b8 {bullet}", builder)
            p.font.size = Pt(14)
            p.space_after = Pt(6)
            p.line_spacing = Pt(20)

    takeaway = spec.get("takeaway", "")
    if takeaway:
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(6.7), Inches(12), Inches(0.6))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = f"\u25b6 {takeaway}"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = builder._get_color("accent", "#C00000")

    if spec.get("notes"):
        slide.notes_slide.notes_text_frame.text = spec["notes"]
    _add_page_number(slide, builder, builder._slide_count + 1)


def _add_infographic(builder: PPTXBuilder, spec: dict, asset_path: Optional[Path]):
    """信息图页：图 + 要点说明。图文比例不合时自动拆页。"""
    # 尝试智能拆页：当图片在左侧区域(7.5"×5")面积利用率低时拆分
    if asset_path and spec.get("bullets") and not builder._disable_auto_split:
        if _add_split_image_text(builder, spec, asset_path,
                                  Inches(7.5), Inches(5)):
            return  # 已拆页处理

    slide = builder._add_blank_slide()
    _add_title_textbox(slide, spec.get("title", ""), builder)
    _add_bottom_bar(slide, builder)

    # 副标题
    subtitle = spec.get("subtitle", "")
    if subtitle:
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(12), Inches(0.4))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(0x88, 0x88, 0x88)

    bullets = spec.get("bullets", [])
    has_bullets = bool(bullets)

    # 渲染信息图：使用 AI 生成的图片
    if asset_path and asset_path.is_file():
        if has_bullets:
            _add_picture_fit(slide, str(asset_path), Inches(0.3), Inches(1.8), Inches(7.5), Inches(5))
        else:
            _add_picture_fit(slide, str(asset_path), Inches(1), Inches(1.6), Inches(11), Inches(5.2))

    if has_bullets:
        # 右侧要点说明（紧凑布局）
        n_bullets = len(bullets)
        bullet_spacing = min(1.15, 4.5 / max(n_bullets, 1))
        for i, bullet in enumerate(bullets):
            y = Inches(2.0) + i * Inches(bullet_spacing)
            # 小色块图标
            icon = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(8.2), y, Inches(0.25), Inches(0.25),
            )
            icon.fill.solid()
            icon.fill.fore_color.rgb = builder._get_color("accent", "#C00000")
            icon.line.fill.background()

            # 文字（支持 **bold** 标记）
            txBox = slide.shapes.add_textbox(Inches(8.6), y - Inches(0.05), Inches(4.2), Inches(bullet_spacing))
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            _add_rich_text(p, str(bullet), builder)
            p.font.size = Pt(14)
            p.line_spacing = Pt(20)

    takeaway = spec.get("takeaway", "")
    if takeaway:
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(6.7), Inches(12), Inches(0.6))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = f"\u25b6 {takeaway}"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = builder._get_color("accent", "#C00000")

    if spec.get("notes"):
        slide.notes_slide.notes_text_frame.text = spec["notes"]
    _add_page_number(slide, builder, builder._slide_count + 1)


def _add_two_column(builder: PPTXBuilder, spec: dict, asset_path: Optional[Path]):
    """双栏对比页。支持 left/right 字段 或 bullets 自动分栏。图文比例不合时自动拆页。"""
    if asset_path and spec.get("bullets") and not builder._disable_auto_split:
        if _add_split_image_text(builder, spec, asset_path,
                                  Inches(12), Inches(1.8)):
            return

    slide = builder._add_blank_slide()
    _add_title_textbox(slide, spec.get("title", ""), builder)
    _add_bottom_bar(slide, builder)

    left = spec.get("left", {})
    right = spec.get("right", {})
    bullets = spec.get("bullets", [])

    # 如果没有 left/right 但有 bullets，自动分为两栏
    if not left.get("bullets") and not right.get("bullets") and bullets:
        mid = (len(bullets) + 1) // 2
        left = {"bullets": bullets[:mid]}
        right = {"bullets": bullets[mid:]}
        # 如果标题有对比关键词（vs / VS），尝试拆分为两个标题
        title = spec.get("title", "")
        import re
        vs_match = re.split(r'\s+vs\.?\s+|\s+VS\.?\s+|\s*[：:]\s*', title, maxsplit=1)
        if len(vs_match) == 2:
            left["heading"] = vs_match[0].strip()
            right["heading"] = vs_match[1].strip()

    has_visual = asset_path and asset_path.is_file()
    # 原生图表渲染（two_column 也可能有 matplotlib）
    visual = spec.get("visual", {})
    chart_rendered = False
    if visual and isinstance(visual, dict) and visual.get("type") == "matplotlib":
        try:
            from src.visuals.pptx_charts import render_native_chart
            chart_y = Inches(5.3)
            chart_h = Inches(1.8)
            if not left.get("bullets") and not right.get("bullets"):
                chart_y = Inches(1.8)
                chart_h = Inches(5.0)
            result = render_native_chart(slide, visual, builder.colors,
                                Inches(0.5), chart_y, Inches(12), chart_h)
            chart_rendered = result is not None
        except Exception as e:
            print(f"  [WARN] two_column 原生图表渲染失败 ({spec.get('id','')}): {e}", flush=True)

    content_height = Inches(3.2) if (has_visual or chart_rendered) else Inches(4.8)

    # 中间分隔线
    divider = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(6.5), Inches(1.8), Pt(2), content_height,
    )
    divider.fill.solid()
    divider.fill.fore_color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    divider.line.fill.background()

    # 左栏
    txBox_l = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(5.8), content_height)
    tf_l = txBox_l.text_frame
    tf_l.word_wrap = True
    first_para = True
    if left.get("heading"):
        p = tf_l.paragraphs[0]
        run = p.add_run()
        run.text = left["heading"]
        run.font.size = Pt(20)
        run.font.bold = True
        run.font.color.rgb = builder._get_color("primary")
        _set_font_name(run)
        p.space_after = Pt(8)
        first_para = False
    for bullet in left.get("bullets", []):
        p = tf_l.paragraphs[0] if first_para else tf_l.add_paragraph()
        first_para = False
        _add_rich_text(p, f"• {bullet}", builder)
        p.font.size = Pt(15)
        p.space_after = Pt(6)
        p.line_spacing = Pt(22)

    # 右栏
    txBox_r = slide.shapes.add_textbox(Inches(7), Inches(1.8), Inches(5.8), content_height)
    tf_r = txBox_r.text_frame
    tf_r.word_wrap = True
    first_para = True
    if right.get("heading"):
        p = tf_r.paragraphs[0]
        run = p.add_run()
        run.text = right["heading"]
        run.font.size = Pt(20)
        run.font.bold = True
        run.font.color.rgb = builder._get_color("accent", "#C00000")
        _set_font_name(run)
        p.space_after = Pt(8)
        first_para = False
    for bullet in right.get("bullets", []):
        p = tf_r.paragraphs[0] if first_para else tf_r.add_paragraph()
        first_para = False
        _add_rich_text(p, f"• {bullet}", builder)
        p.font.size = Pt(15)
        p.space_after = Pt(6)
        p.line_spacing = Pt(22)

    # 底部图表/图片（未原生渲染时回退到图片）
    if not chart_rendered and has_visual:
        _add_picture_fit(slide, str(asset_path), Inches(0.5), Inches(5.3), Inches(12), Inches(1.8))

    if spec.get("notes"):
        slide.notes_slide.notes_text_frame.text = spec["notes"]
    _add_page_number(slide, builder, builder._slide_count + 1)


def _add_key_insight(builder: PPTXBuilder, spec: dict, asset_path: Optional[Path]):
    """核心发现页：标题 + 要点列表 + KPI 图。图文比例不合时自动拆页。"""
    if asset_path and spec.get("bullets") and not builder._disable_auto_split:
        if _add_split_image_text(builder, spec, asset_path,
                                  Inches(5), Inches(4.5)):
            return

    slide = builder._add_blank_slide()
    _add_title_textbox(slide, spec.get("title", ""), builder)
    _add_bottom_bar(slide, builder)

    # 顶部强调色条
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), builder.prs.slide_width, Inches(0.12),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = builder._get_color("accent", "#C00000")
    bar.line.fill.background()

    bullets = spec.get("bullets", [])
    visual = spec.get("visual", {})

    has_visual = asset_path and asset_path.is_file()

    if bullets:
        # 左侧要点列表
        text_width = Inches(6.5) if has_visual else Inches(11.5)
        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), text_width, Inches(4.5))
        tf = txBox.text_frame
        tf.word_wrap = True
        for i, bullet in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            _add_rich_text(p, f"\u25b8 {bullet}", builder)
            p.font.size = Pt(16)
            p.space_after = Pt(10)
            p.line_spacing = Pt(26)

    # 插入信息图图片
    if asset_path and asset_path.is_file():
        if bullets:
            _add_picture_fit(slide, str(asset_path), Inches(7.8), Inches(1.8), Inches(5), Inches(4.5))
        else:
            _add_picture_fit(slide, str(asset_path), Inches(1.5), Inches(1.8), Inches(10), Inches(4.8))

    takeaway = spec.get("takeaway", "")
    if takeaway:
        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(6.5), Inches(11), Inches(0.7))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = f"\u25b6 {takeaway}"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = builder._get_color("accent", "#C00000")

    if spec.get("notes"):
        slide.notes_slide.notes_text_frame.text = spec["notes"]
    _add_page_number(slide, builder, builder._slide_count + 1)


def _parse_bullets_as_table(bullets: list):
    """
    Try to extract table structure from bullet strings.
    Supports delimiters: "|" (pipe) and ":" (colon, key-value pairs).
    Returns (headers, rows) or ([], []) if parsing fails.
    """
    if not bullets or len(bullets) < 2:
        return [], []

    # ── Strategy 1: pipe-delimited bullets ──
    if any("|" in str(b) for b in bullets):
        parsed = [str(b).split("|") for b in bullets]
        parsed = [[c.strip() for c in row if c.strip()] for row in parsed]
        col_counts = [len(r) for r in parsed if len(r) >= 2]
        if col_counts:
            # Use first row as header if all rows have similar column counts
            mode_cols = max(set(col_counts), key=col_counts.count)
            consistent = [r for r in parsed if len(r) == mode_cols]
            if len(consistent) >= 2:
                return consistent[0], consistent[1:]

    # ── Strategy 2: colon-delimited key-value pairs ──
    kv_pairs = []
    for b in bullets:
        b_str = str(b).strip()
        # Remove leading bullet markers like "•", "-", "·"
        b_str = b_str.lstrip("•·-– ").strip()
        if ":" in b_str or "：" in b_str:
            # Split on first colon (either half-width or full-width)
            for sep in ["：", ":"]:
                if sep in b_str:
                    key, val = b_str.split(sep, 1)
                    kv_pairs.append((key.strip(), val.strip()))
                    break
        else:
            kv_pairs = []
            break

    if len(kv_pairs) >= 2:
        headers = ["项目", "内容"]
        rows = [[k, v] for k, v in kv_pairs]
        return headers, rows

    return [], []


def _add_table(builder: PPTXBuilder, spec: dict, asset_path: Optional[Path]):
    """表格页。支持从 bullets 自动解析表格结构，无法解析时回退到 title_content。"""
    table_data = spec.get("table", {})
    headers = table_data.get("headers", [])
    rows = table_data.get("rows", [])

    if not headers:
        # ── Fallback: try to parse bullets into a table structure ──
        bullets = spec.get("bullets", [])
        if bullets:
            parsed_headers, parsed_rows = _parse_bullets_as_table(bullets)
            if parsed_headers:
                headers = parsed_headers
                rows = parsed_rows
            else:
                # Cannot parse as table → delegate to title_content layout
                _add_title_content(builder, spec, asset_path)
                return
        if not headers:
            # No table data and no bullets → still create slide with title only
            slide = builder._add_blank_slide()
            _add_title_textbox(slide, spec.get("title", ""), builder)
            _add_bottom_bar(slide, builder)
            return

    slide = builder._add_blank_slide()
    _add_title_textbox(slide, spec.get("title", ""), builder)
    _add_bottom_bar(slide, builder)

    total_rows = len(rows) + 1
    total_cols = len(headers)
    table_shape = slide.shapes.add_table(
        total_rows, total_cols,
        Inches(0.8), Inches(1.8), Inches(11.5), Inches(4.8),
    )
    table = table_shape.table

    # 表头样式
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = str(h)
        cell.fill.solid()
        cell.fill.fore_color.rgb = builder._get_color("primary", "#002060")
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(14)
            paragraph.font.bold = True
            paragraph.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # 数据行（交替背景色）
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            if j < total_cols:
                cell = table.cell(i + 1, j)
                cell.text = str(val)
                if i % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(0xF5, 0xF7, 0xFA)
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.size = Pt(11)

    if spec.get("notes"):
        slide.notes_slide.notes_text_frame.text = spec["notes"]
    _add_page_number(slide, builder, builder._slide_count + 1)


def _add_quote(builder: PPTXBuilder, spec: dict, asset_path: Optional[Path]):
    """引用页：居中大字。"""
    slide = builder._add_blank_slide()

    # 浅色背景
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = RGBColor(0xF8, 0xFA, 0xFC)

    # 引号装饰
    txq = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(2), Inches(1.5))
    tf_q = txq.text_frame
    p_q = tf_q.paragraphs[0]
    p_q.text = "\u201C"  # 左双引号
    p_q.font.size = Pt(120)
    p_q.font.color.rgb = builder._get_color("accent", "#C00000")
    p_q.font.bold = True

    quote = spec.get("quote", spec.get("title", ""))
    txBox = slide.shapes.add_textbox(Inches(2), Inches(2.5), Inches(9), Inches(2.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = quote
    p.font.size = Pt(32)
    p.font.italic = True
    p.font.color.rgb = builder._get_color("primary")
    p.alignment = PP_ALIGN.CENTER

    source = spec.get("source", "")
    if source:
        p2 = tf.add_paragraph()
        p2.text = f"— {source}"
        p2.font.size = Pt(16)
        p2.font.color.rgb = RGBColor(0x99, 0x99, 0x99)
        p2.alignment = PP_ALIGN.RIGHT
        p2.space_before = Pt(20)

    _add_bottom_bar(slide, builder)


def _add_summary(builder: PPTXBuilder, spec: dict, asset_path: Optional[Path]):
    """总结页：要点回顾 + 底部 KPI + 行动号召。图文比例不合时自动拆页。"""
    if asset_path and spec.get("bullets") and not builder._disable_auto_split:
        if _add_split_image_text(builder, spec, asset_path,
                                  Inches(12), Inches(2)):
            return

    slide = builder._add_blank_slide()

    _add_title_textbox(slide, spec.get("title", "总结"), builder)
    _add_bottom_bar(slide, builder)

    # 内容卡片
    content_card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(1.6), Inches(12), Inches(3.2),
    )
    content_card.fill.solid()
    content_card.fill.fore_color.rgb = RGBColor(0xF5, 0xF7, 0xFA)
    content_card.line.fill.background()

    bullets = spec.get("bullets", [])
    txBox = slide.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(11), Inches(3.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        _add_rich_text(p, f"\u2713 {bullet}", builder)
        p.font.size = Pt(16)
        p.space_after = Pt(6)
        p.line_spacing = Pt(24)

    # 底部 KPI 信息图（仅此一处使用图片，不做全页背景叠加）
    if asset_path and asset_path.is_file():
        _add_picture_fit(slide, str(asset_path), Inches(0.5), Inches(5.0), Inches(12), Inches(2.0))

    cta = spec.get("call_to_action", "")
    if cta:
        # 行动号召背景框
        cta_bg = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(2), Inches(6.2), Inches(9), Inches(1.0),
        )
        cta_bg.fill.solid()
        cta_bg.fill.fore_color.rgb = builder._get_color("accent", "#C00000")
        cta_bg.line.fill.background()

        txBox2 = slide.shapes.add_textbox(Inches(2.5), Inches(6.3), Inches(8), Inches(0.8))
        tf2 = txBox2.text_frame
        p = tf2.paragraphs[0]
        p.text = cta
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        p.alignment = PP_ALIGN.CENTER

    if spec.get("notes"):
        slide.notes_slide.notes_text_frame.text = spec["notes"]
    _add_page_number(slide, builder, builder._slide_count + 1)


def _add_image_full(builder: PPTXBuilder, spec: dict, asset_path: Optional[Path]):
    """全屏图片页。"""
    slide = builder._add_blank_slide()
    if asset_path and asset_path.is_file():
        slide.shapes.add_picture(
            str(asset_path), Inches(0), Inches(0),
            builder.prs.slide_width, builder.prs.slide_height,
        )
    if spec.get("notes"):
        slide.notes_slide.notes_text_frame.text = spec["notes"]


def _add_timeline(builder: PPTXBuilder, spec: dict, asset_path: Optional[Path]):
    """时间线页（同 infographic）。"""
    _add_infographic(builder, spec, asset_path)


def _add_methodology(builder: PPTXBuilder, spec: dict, asset_path: Optional[Path]):
    """研究方法/技术路线页：编号分步 + 技术路线图。图文比例不合时自动拆页。"""
    if asset_path and spec.get("bullets") and not builder._disable_auto_split:
        if _add_split_image_text(builder, spec, asset_path,
                                  Inches(11), Inches(3)):
            return

    slide = builder._add_blank_slide()
    _add_title_textbox(slide, spec.get("title", ""), builder)
    _add_bottom_bar(slide, builder)

    bullets = spec.get("bullets", [])
    has_visual = asset_path and asset_path.is_file()

    # 编号分步展示（顶部水平排列）
    if bullets:
        n = len(bullets)
        step_width = min(3.0, 11.5 / max(n, 1))
        start_x = (13.333 - step_width * n) / 2

        for i, bullet in enumerate(bullets):
            x = Inches(start_x + i * step_width)
            y_top = Inches(1.8) if has_visual else Inches(2.0)

            # 编号圆形徽章
            badge = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                x + Inches(step_width / 2 - 0.25), y_top, Inches(0.5), Inches(0.5),
            )
            badge.fill.solid()
            badge.fill.fore_color.rgb = builder._get_color("secondary", "#0060A8")
            badge.line.fill.background()
            txb = badge.text_frame
            txb.paragraphs[0].alignment = PP_ALIGN.CENTER
            run = txb.paragraphs[0].add_run()
            run.text = str(i + 1)
            run.font.size = Pt(18)
            run.font.bold = True
            run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            _set_font_name(run)
            txb.vertical_anchor = MSO_ANCHOR.MIDDLE

            # 步骤文字
            txBox = slide.shapes.add_textbox(
                x, y_top + Inches(0.6), Inches(step_width - 0.1), Inches(1.2),
            )
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            _add_rich_text(p, str(bullet), builder)
            p.font.size = Pt(14)
            p.alignment = PP_ALIGN.CENTER

            # 箭头连接（除最后一个）
            if i < n - 1:
                arrow = slide.shapes.add_shape(
                    MSO_SHAPE.RIGHT_ARROW,
                    x + Inches(step_width - 0.3), y_top + Inches(0.1),
                    Inches(0.3), Inches(0.3),
                )
                arrow.fill.solid()
                arrow.fill.fore_color.rgb = builder._get_color("accent", "#C00000")
                arrow.line.fill.background()

    # 下方技术路线图/流程图
    if has_visual:
        img_top = Inches(4.0) if bullets else Inches(1.8)
        img_height = Inches(3.0) if bullets else Inches(5.0)
        _add_picture_fit(slide, str(asset_path), Inches(1), img_top, Inches(11), img_height)

    takeaway = spec.get("takeaway", "")
    if takeaway:
        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(6.7), Inches(11), Inches(0.6))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = f"\u25b6 {takeaway}"
        run.font.size = Pt(14)
        run.font.bold = True
        run.font.color.rgb = builder._get_color("accent", "#C00000")
        _set_font_name(run)

    if spec.get("notes"):
        slide.notes_slide.notes_text_frame.text = spec["notes"]
    _add_page_number(slide, builder, builder._slide_count + 1)


def _add_architecture(builder: PPTXBuilder, spec: dict, asset_path: Optional[Path]):
    """系统架构页：全宽架构图 + 底部要点。"""
    slide = builder._add_blank_slide()
    _add_title_textbox(slide, spec.get("title", ""), builder)
    _add_bottom_bar(slide, builder)

    bullets = spec.get("bullets", [])
    has_visual = asset_path and asset_path.is_file()

    # 架构图（全宽展示）
    if has_visual:
        img_height = Inches(4.0) if bullets else Inches(5.2)
        _add_picture_fit(slide, str(asset_path), Inches(0.8), Inches(1.6), Inches(11.5), img_height)

    # 底部要点（水平卡片排列）
    if bullets:
        n = len(bullets)
        card_width = min(3.5, 11.5 / max(n, 1))
        start_x = (13.333 - card_width * n) / 2
        y = Inches(5.8)

        for i, bullet in enumerate(bullets):
            x = Inches(start_x + i * card_width)

            # 卡片背景
            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                x, y, Inches(card_width - 0.15), Inches(1.2),
            )
            card.fill.solid()
            card.fill.fore_color.rgb = RGBColor(0xF0, 0xF4, 0xFA)
            card.line.fill.background()

            # 顶部色条
            bar = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                x, y, Inches(card_width - 0.15), Pt(4),
            )
            bar.fill.solid()
            bar.fill.fore_color.rgb = builder._get_color("secondary", "#0060A8")
            bar.line.fill.background()

            # 文字
            txBox = slide.shapes.add_textbox(
                x + Inches(0.1), y + Inches(0.15), Inches(card_width - 0.35), Inches(0.95),
            )
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            _add_rich_text(p, str(bullet), builder)
            p.font.size = Pt(14)

    if spec.get("notes"):
        slide.notes_slide.notes_text_frame.text = spec["notes"]
    _add_page_number(slide, builder, builder._slide_count + 1)


def _add_end(builder: PPTXBuilder, spec: dict, asset_path: Optional[Path]):
    """结束页：致谢大字 + 联系信息，与封面风格呼应。"""
    slide = builder._add_blank_slide()

    # 背景图
    if asset_path and asset_path.is_file():
        slide.shapes.add_picture(
            str(asset_path), Inches(0), Inches(0),
            builder.prs.slide_width, builder.prs.slide_height,
        )

    # 深色遮罩
    overlay = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), builder.prs.slide_width, builder.prs.slide_height,
    )
    overlay.fill.solid()
    overlay.fill.fore_color.rgb = builder._get_color("primary", "#002060")
    _set_shape_transparency(overlay, 45000)
    overlay.line.fill.background()

    # 装饰线
    accent_line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(5.5), Inches(2.5), Inches(2.3), Pt(5),
    )
    accent_line.fill.solid()
    accent_line.fill.fore_color.rgb = builder._get_color("accent", "#C00000")
    accent_line.line.fill.background()

    # 大字致谢
    title = spec.get("title", "感谢聆听")
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.8), Inches(11), Inches(2))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(56)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    _set_font_name(run)
    p.alignment = PP_ALIGN.CENTER

    # 副标题/联系信息
    subtitle = spec.get("subtitle", "")
    if subtitle:
        p2 = tf.add_paragraph()
        run2 = p2.add_run()
        run2.text = subtitle
        run2.font.size = Pt(18)
        run2.font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
        _set_font_name(run2)
        p2.alignment = PP_ALIGN.CENTER
        p2.space_before = Pt(20)

    # 底部色条
    _add_bottom_bar(slide, builder)


# ============ 布局处理器注册 ============
_LAYOUT_HANDLERS = {
    "cover": _add_cover,
    "agenda": _add_agenda,
    "section_break": _add_section_break,
    "title_content": _add_title_content,
    "data_chart": _add_data_chart,
    "infographic": _add_infographic,
    "two_column": _add_two_column,
    "key_insight": _add_key_insight,
    "table": _add_table,
    "image_full": _add_image_full,
    "quote": _add_quote,
    "timeline": _add_timeline,
    "summary": _add_summary,
    "methodology": _add_methodology,
    "architecture": _add_architecture,
    "end": _add_end,
}
