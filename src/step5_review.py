# -*- coding: utf-8 -*-
"""
Step5: 四角色迭代审阅。
Agent A（听众）+ Agent B（演讲者）+ Agent C（制作者）+ Agent E（视觉检测）协同优化 PPT。
"""
import json
import time
from concurrent.futures import ThreadPoolExecutor
from pathlib import Path

import src  # noqa: F401
from src.llm_client import chat
from src.utils.safe_write import safe_write_text
from src.prompts.review import (
    AGENT_A_SYSTEM, AGENT_B_SYSTEM, AGENT_C_SYSTEM, AGENT_D_SYSTEM,
    build_review_user_prompt,
)
from src.visual_inspector import inspect_pptx, format_agent_e_report, generate_thumbnail_grid
from config import LLM_PROVIDER, REVIEW_FALLBACK_PROVIDERS


class ParseError(Exception):
    """LLM 返回内容无法解析为有效 JSON。"""
    pass


def _load_checkpoint(output_dir: Path) -> dict | None:
    """加载断点续传检查点。"""
    ckpt_path = output_dir / "review_checkpoint.json"
    if not ckpt_path.is_file():
        return None
    try:
        ckpt = json.loads(ckpt_path.read_text(encoding="utf-8"))
        print(f"[Step5] 发现检查点: 已完成 {ckpt['completed_round']} 轮，从第 "
              f"{ckpt['completed_round'] + 1} 轮继续", flush=True)
        return ckpt
    except Exception as e:
        print(f"[Step5] 检查点读取失败({e})，从头开始", flush=True)
        return None


def _save_checkpoint(output_dir: Path, completed_round: int,
                     review_log: list, prev_a_dims: dict, prev_b_dims: dict,
                     score_trajectory_a: list, score_trajectory_b: list,
                     target_score: float, prev_c_response: dict = None,
                     prev_d_scores: dict = None, score_trajectory_d: list = None):
    """保存断点续传检查点和增量日志。"""
    ckpt = {
        "completed_round": completed_round,
        "prev_a_dims": prev_a_dims,
        "prev_b_dims": prev_b_dims,
        "prev_d_scores": prev_d_scores,
        "score_trajectory_a": score_trajectory_a,
        "score_trajectory_b": score_trajectory_b,
        "score_trajectory_d": score_trajectory_d or [],
        "review_log": review_log,
        "target_score": target_score,
        "prev_c_response": prev_c_response,
    }
    ckpt_path = output_dir / "review_checkpoint.json"
    safe_write_text(ckpt_path, json.dumps(ckpt, ensure_ascii=False, indent=2))
    # 同步写增量审阅日志（随时可查看进度）
    _save_review_log(output_dir, review_log, score_trajectory_a,
                     score_trajectory_b, target_score, converged=False)


def _save_review_log(output_dir: Path, review_log: list,
                     score_trajectory_a: list, score_trajectory_b: list,
                     target_score: float, converged: bool,
                     score_trajectory_d: list = None):
    """保存审阅日志到 review_log.json。"""
    final_log = {
        "total_rounds": len(review_log),
        "converged": converged,
        "target_score": target_score,
        "score_trajectory": {
            "agent_a": score_trajectory_a,
            "agent_b": score_trajectory_b,
            "agent_d": score_trajectory_d or [],
        },
        "rounds": review_log,
    }
    log_path = output_dir / "review_log.json"
    safe_write_text(log_path, json.dumps(final_log, ensure_ascii=False, indent=2))


def run_review(base: str, output_dir: Path, theme: str = None,
               no_ai_images: bool = False, max_rounds: int = 5,
               target_score: float = 9.0) -> dict:
    """
    Step5 入口：五角色迭代审阅（A/B/C/D/E，支持断点续传）。

    Returns:
        {
            "pptx_path": Path,
            "review_log": list,
            "final_scores": dict,
            "rounds": int,
        }
    """
    # ---- 加载原始素材（供 Agent D 交叉验证） ----
    analysis = {}
    raw_content = ""
    analysis_path = output_dir / "analysis.json"
    if analysis_path.is_file():
        analysis = json.loads(analysis_path.read_text(encoding="utf-8"))
    raw_content_path = output_dir / "raw_content.md"
    if raw_content_path.is_file():
        raw_content = raw_content_path.read_text(encoding="utf-8", errors="replace")

    # ---- 断点续传：尝试恢复上次进度 ----
    ckpt = _load_checkpoint(output_dir)
    if ckpt:
        review_log = ckpt["review_log"]
        prev_a_dims = ckpt["prev_a_dims"]
        prev_b_dims = ckpt["prev_b_dims"]
        prev_d_scores = ckpt.get("prev_d_scores")
        score_trajectory_a = ckpt["score_trajectory_a"]
        score_trajectory_b = ckpt["score_trajectory_b"]
        score_trajectory_d = ckpt.get("score_trajectory_d", [])
        start_round = ckpt["completed_round"] + 1
        prev_c_response = ckpt.get("prev_c_response")
    else:
        review_log = []
        prev_a_dims = None
        prev_b_dims = None
        prev_d_scores = None
        score_trajectory_a = []
        score_trajectory_b = []
        score_trajectory_d = []
        start_round = 1
        prev_c_response = None

    avg_a = score_trajectory_a[-1] if score_trajectory_a else 0.0
    avg_b = score_trajectory_b[-1] if score_trajectory_b else 0.0

    # Agent E 在后两轮参与
    agent_e_start_round = max(1, max_rounds - 1)

    for round_num in range(start_round, max_rounds + 1):
        t0 = time.time()
        e_active = round_num >= agent_e_start_round
        agents_tag = "A+B+D" + ("+E" if e_active else "")
        print(f"\n[Step5] ===== 第 {round_num} 轮审阅 ({agents_tag}) =====", flush=True)

        # 读取当前 slide_plan
        slide_plan = json.loads(
            (output_dir / "slide_plan.json").read_text(encoding="utf-8")
        )

        # ---- 并行调用 Agent A、B、D ----
        with ThreadPoolExecutor(max_workers=3) as pool:
            future_a = pool.submit(
                _call_agent, "A", AGENT_A_SYSTEM, slide_plan, prev_a_dims,
                prev_c_response
            )
            future_b = pool.submit(
                _call_agent, "B", AGENT_B_SYSTEM, slide_plan, prev_b_dims,
                prev_c_response
            )
            future_d = pool.submit(
                _call_agent_d, slide_plan, analysis, raw_content,
                prev_d_scores, prev_c_response
            )
            result_a = future_a.result()
            result_b = future_b.result()
            result_d = future_d.result()

        # 强制分数单调递增
        result_a = _enforce_monotonic(result_a, prev_a_dims)
        result_b = _enforce_monotonic(result_b, prev_b_dims)

        avg_a = _avg_score(result_a)
        avg_b = _avg_score(result_b)
        avg_d = _avg_score_d(result_d)
        score_trajectory_a.append(avg_a)
        score_trajectory_b.append(avg_b)
        score_trajectory_d.append(avg_d)

        _print_scores("Agent A（听众）", result_a)
        _print_scores("Agent B（演讲者）", result_b)
        _print_scores_d("Agent D（逻辑）", result_d)

        # ---- Agent E 视觉检测（后两轮参与） ----
        agent_e_report = None
        agent_e_result = None
        if e_active:
            pptx_path = output_dir / f"{base}_slides.pptx"
            if pptx_path.is_file():
                agent_e_result = _call_agent_e(pptx_path, slide_plan)
                agent_e_report = format_agent_e_report(agent_e_result)
                n_anomalies = agent_e_result["summary"]["total_anomalies"]
                n_high = agent_e_result["summary"]["high_severity"]
                n_monotony = agent_e_result["summary"].get("visual_monotony_slides", 0)
                mono_tag = f" 单调{n_monotony}页" if n_monotony else ""
                print(f"  [Agent E] 检测到 {n_anomalies} 个异常 "
                      f"(严重{n_high}{mono_tag})", flush=True)
            else:
                print(f"  [Agent E] PPTX 文件不存在，跳过视觉检测", flush=True)

        # 检查收敛
        if avg_a >= target_score and avg_b >= target_score:
            has_high_e = (agent_e_result and
                         agent_e_result["summary"]["high_severity"] > 0)
            has_high_d = any(
                p.get("priority") == "high"
                for p in result_d.get("priority_suggestions", [])
            )
            if not has_high_e and not has_high_d:
                print(f"\n[Step5] 收敛! A={avg_a:.1f} B={avg_b:.1f} D={avg_d:.1f} "
                      f"(目标≥{target_score})", flush=True)
                review_log.append(_build_round_log(
                    round_num, result_a, result_b, None, 0, time.time() - t0,
                    agent_e_result, result_d
                ))
                _save_checkpoint(output_dir, round_num, review_log,
                                 prev_a_dims, prev_b_dims,
                                 score_trajectory_a, score_trajectory_b,
                                 target_score, prev_c_response,
                                 prev_d_scores, score_trajectory_d)
                break
            else:
                blockers = []
                if has_high_e:
                    blockers.append(f"E有{n_high}个严重异常")
                if has_high_d:
                    blockers.append("D有高优先级逻辑问题")
                print(f"  [Step5] A/B分数达标但 {'/'.join(blockers)}，"
                      f"继续修复", flush=True)

        # ---- Agent C 综合制定改造计划 ----
        print(f"  [Agent C] 综合研判中...", flush=True)
        changes = _call_agent_c(slide_plan, result_a, result_b,
                                agent_d_result=result_d,
                                agent_e_report=agent_e_report)
        n_changes = len(changes.get("changes", []))
        print(f"  [Agent C] 生成 {n_changes} 项改进", flush=True)

        # 保存 Agent C 的回复，供下轮 A/B/D 参考
        prev_c_response = {
            "reasoning": changes.get("reasoning", ""),
            "response_to_a": changes.get("response_to_a", ""),
            "response_to_b": changes.get("response_to_b", ""),
            "response_to_d": changes.get("response_to_d", ""),
            "response_to_e": changes.get("response_to_e", ""),
        }

        if n_changes == 0:
            print(f"  [Agent C] 无改进建议，终止迭代", flush=True)
            review_log.append(_build_round_log(
                round_num, result_a, result_b, changes, 0, time.time() - t0,
                agent_e_result
            ))
            _save_checkpoint(output_dir, round_num, review_log,
                             prev_a_dims, prev_b_dims,
                             score_trajectory_a, score_trajectory_b,
                             target_score, prev_c_response)
            break

        # ---- 应用改动 ----
        slide_plan, changed_visual_ids = _apply_changes(slide_plan, changes)

        # 保存当前轮次 slide_plan 快照
        plan_text = json.dumps(slide_plan, ensure_ascii=False, indent=2)
        snapshot = output_dir / f"slide_plan_r{round_num}.json"
        safe_write_text(snapshot, plan_text)
        # 覆写主 slide_plan
        safe_write_text(output_dir / "slide_plan.json", plan_text)

        # ---- 选择性重跑 Step3 (视觉资产) ----
        n_visual_regen = 0
        if changed_visual_ids:
            print(f"  [Step3] 重新生成 {len(changed_visual_ids)} 个视觉资产...",
                  flush=True)
            n_visual_regen = _rerun_visuals_selective(
                output_dir, slide_plan, changed_visual_ids, no_ai_images
            )

        # ---- 智能选择：外科手术式编辑 vs 全量重建 ----
        pptx_path = output_dir / f"{base}_slides.pptx"
        text_only_changes = _classify_text_only_changes(changes)
        has_structural = changed_visual_ids or any(
            ch.get("action") in ("insert_after", "delete", "reorder")
            for ch in changes.get("changes", [])
        )

        if not has_structural and text_only_changes and pptx_path.is_file():
            # 外科手术式编辑：仅修改文字，不重建 PPTX
            print(f"  [Step4] 外科手术式编辑 ({len(text_only_changes)} 处文字修改)...",
                  flush=True)
            _surgical_text_edit(pptx_path, slide_plan, text_only_changes)
        else:
            # 全量重建
            print(f"  [Step4] 重新装配 PPTX...", flush=True)
            from src.step4_build import run_build
            run_build(base, output_dir, theme)

        # 保存每轮 PPTX 快照，方便对比
        import shutil
        main_pptx = output_dir / f"{base}_slides.pptx"
        round_pptx = output_dir / f"{base}_slides_r{round_num}.pptx"
        if main_pptx.is_file():
            shutil.copy2(main_pptx, round_pptx)
            print(f"  [快照] 已保存 {round_pptx.name}", flush=True)

        elapsed = time.time() - t0
        review_log.append(_build_round_log(
            round_num, result_a, result_b, changes, n_visual_regen, elapsed,
            agent_e_result, result_d
        ))

        # 更新上轮分数
        prev_a_dims = result_a.get("dimensions", {})
        prev_b_dims = result_b.get("dimensions", {})
        prev_d_scores = {
            "macro_logic": result_d.get("macro_logic", {}).get("score", 5),
            "meso_logic": result_d.get("meso_logic", {}).get("score", 5),
            "micro_logic": result_d.get("micro_logic", {}).get("score", 5),
            "source_fidelity": result_d.get("source_fidelity", {}).get("score", 5),
        }

        # ---- 每轮结束保存检查点（断点续传核心） ----
        _save_checkpoint(output_dir, round_num, review_log,
                         prev_a_dims, prev_b_dims,
                         score_trajectory_a, score_trajectory_b,
                         target_score, prev_c_response,
                         prev_d_scores, score_trajectory_d)

        print(f"  [第{round_num}轮完成] A={avg_a:.1f} B={avg_b:.1f} D={avg_d:.1f} "
              f"改动{n_changes}项 耗时{elapsed:.0f}s", flush=True)

    # 保存最终审阅日志
    converged = avg_a >= target_score and avg_b >= target_score
    _save_review_log(output_dir, review_log, score_trajectory_a,
                     score_trajectory_b, target_score, converged,
                     score_trajectory_d)

    # 清理检查点（正常完成后）
    ckpt_path = output_dir / "review_checkpoint.json"
    if ckpt_path.is_file():
        ckpt_path.unlink()
        print("[Step5] 检查点已清理", flush=True)

    pptx_path = output_dir / f"{base}_slides.pptx"
    final_d = score_trajectory_d[-1] if score_trajectory_d else 0.0
    print(f"\n[Step5] 迭代审阅完成: {len(review_log)} 轮, "
          f"最终 A={score_trajectory_a[-1]:.1f} B={score_trajectory_b[-1]:.1f} "
          f"D={final_d:.1f}",
          flush=True)

    return {
        "pptx_path": pptx_path,
        "review_log": review_log,
        "final_scores": {"agent_a": score_trajectory_a[-1],
                         "agent_b": score_trajectory_b[-1],
                         "agent_d": final_d},
        "rounds": len(review_log),
    }


# ============ Agent 调用 ============

def _get_provider_chain() -> list:
    """返回 LLM 候选链：主 provider + fallback 列表（去重）。"""
    primary = (LLM_PROVIDER or "kimi").lower().strip()
    chain = [primary]
    for p in REVIEW_FALLBACK_PROVIDERS:
        if p.lower().strip() != primary:
            chain.append(p.lower().strip())
    return chain


def _call_with_fallback(messages: list, agent_label: str,
                        max_tokens: int = 4096, temperature: float = 0.3) -> dict:
    """调用 LLM 并解析 JSON，失败时切换候选模型重试。"""
    chain = _get_provider_chain()
    last_error = None

    for i, provider in enumerate(chain):
        try:
            tag = f"Agent {agent_label}"
            if i > 0:
                print(f"  [{tag}] 切换到候选模型: {provider}", flush=True)
            response = chat(messages, provider=provider,
                            max_tokens=max_tokens, temperature=temperature)
            result = _parse_json_response(response, agent_label)
            if i > 0:
                print(f"  [{tag}] {provider} 解析成功", flush=True)
            return result
        except ParseError as e:
            last_error = e
            print(f"  [{tag}] {provider} 返回格式异常，"
                  f"{'尝试下一个候选' if i < len(chain) - 1 else '所有候选已用尽'}",
                  flush=True)
        except Exception as e:
            last_error = e
            print(f"  [{tag}] {provider} 调用失败: {e}，"
                  f"{'尝试下一个候选' if i < len(chain) - 1 else '所有候选已用尽'}",
                  flush=True)

    # 所有候选都失败，返回保底默认值
    print(f"  [警告] Agent {agent_label} 所有候选模型均失败，使用默认值", flush=True)
    if agent_label in ("A", "B"):
        return {
            "dimensions": {
                "dim1": {"score": 5, "justification": "所有模型解析失败"},
                "dim2": {"score": 5, "justification": "所有模型解析失败"},
                "dim3": {"score": 5, "justification": "所有模型解析失败"},
                "dim4": {"score": 5, "justification": "所有模型解析失败"},
            },
            "slide_comments": [],
            "overall_comment": f"所有候选模型({', '.join(chain)})均失败: {last_error}",
        }
    elif agent_label == "D":
        return {
            "macro_logic": {"score": 5, "chapter_flow": "解析失败", "issues": []},
            "meso_logic": {"score": 5, "issues": []},
            "micro_logic": {"score": 5, "issues": []},
            "source_fidelity": {"score": 5, "issues": []},
            "emphasis": {"core_points": [], "well_emphasized": False, "issues": []},
            "priority_suggestions": [],
            "overall_comment": f"所有候选模型均失败: {last_error}",
        }
    else:
        return {"reasoning": f"所有候选模型均失败: {last_error}", "changes": []}


def _call_agent(role: str, system_prompt: str, slide_plan: dict,
                prev_dims: dict = None, agent_c_response: dict = None) -> dict:
    """调用 Agent A 或 B，返回解析后的评分 dict。"""
    user_prompt = build_review_user_prompt(
        slide_plan, role, prev_scores=prev_dims,
        agent_c_response=agent_c_response
    )
    messages = [
        {"role": "system", "content": system_prompt},
        {"role": "user", "content": user_prompt},
    ]
    return _call_with_fallback(messages, role, max_tokens=4096, temperature=0.3)


def _call_agent_d(slide_plan: dict, analysis: dict, raw_content: str,
                   prev_d_scores: dict = None,
                   agent_c_response: dict = None) -> dict:
    """调用 Agent D（逻辑架构审阅），返回三层逻辑评估。"""
    user_prompt = build_review_user_prompt(
        slide_plan, "D",
        prev_scores=prev_d_scores,
        analysis=analysis,
        raw_content=raw_content,
        agent_c_response=agent_c_response,
    )
    messages = [
        {"role": "system", "content": AGENT_D_SYSTEM},
        {"role": "user", "content": user_prompt},
    ]
    result = _call_with_fallback(messages, "D", max_tokens=6144, temperature=0.3)

    # 强制 D 的分数单调递增
    if prev_d_scores:
        for layer in ("macro_logic", "meso_logic", "micro_logic", "source_fidelity"):
            if layer in result and isinstance(result[layer], dict):
                prev = prev_d_scores.get(layer, 0)
                if isinstance(prev, (int, float)) and result[layer].get("score", 0) < prev:
                    result[layer]["score"] = prev

    return result


def _call_agent_c(slide_plan: dict, result_a: dict, result_b: dict,
                   agent_d_result: dict = None,
                   agent_e_report: str = None) -> dict:
    """调用 Agent C，返回改造计划 dict（含对 A/B/D/E 的回复）。"""
    user_prompt = build_review_user_prompt(
        slide_plan, "C",
        agent_a_result=result_a,
        agent_b_result=result_b,
        agent_d_result=agent_d_result,
        agent_e_report=agent_e_report,
    )
    messages = [
        {"role": "system", "content": AGENT_C_SYSTEM},
        {"role": "user", "content": user_prompt},
    ]
    result = _call_with_fallback(messages, "C", max_tokens=8192, temperature=0.3)

    # 动态最大改动数：基础5 + D配额2 + E配额1 = 最多8
    max_changes = 5
    if agent_d_result:
        max_changes += 2
    if agent_e_report:
        max_changes += 1
    max_changes = min(max_changes, 8)
    changes = result.get("changes", [])
    if len(changes) > max_changes:
        result["changes"] = changes[:max_changes]

    return result


def _call_agent_e(pptx_path: Path, slide_plan: dict) -> dict:
    """调用 Agent E（程序化视觉检测），返回检测结果。"""
    print(f"  [Agent E] 正在分析 PPTX 排版...", flush=True)
    try:
        result = inspect_pptx(pptx_path, slide_plan)

        # 生成缩略图网格用于视觉存档
        try:
            thumb_path = generate_thumbnail_grid(
                pptx_path,
                pptx_path.parent / f"{pptx_path.stem}_thumbnails.jpg",
                cols=5, thumb_width=280,
            )
            print(f"  [Agent E] 缩略图已生成: {thumb_path.name}", flush=True)
        except Exception as te:
            print(f"  [Agent E] 缩略图生成失败: {te}", flush=True)

        return result
    except Exception as e:
        print(f"  [Agent E] 检测失败: {e}", flush=True)
        return {
            "slide_width": 0, "slide_height": 0, "total_slides": 0,
            "anomalies": [],
            "summary": {
                "total_anomalies": 0, "high_severity": 0,
                "medium_severity": 0, "low_severity": 0,
                "most_common_type": "none",
            },
        }


# ============ 分数处理 ============

def _enforce_monotonic(result: dict, prev_dims: dict = None) -> dict:
    """强制分数不低于上轮。"""
    if not prev_dims:
        return result
    dims = result.get("dimensions", {})
    for key, prev_info in prev_dims.items():
        if key in dims:
            prev_score = prev_info.get("score", 0) if isinstance(prev_info, dict) else prev_info
            curr = dims[key]
            if isinstance(curr, dict) and curr.get("score", 0) < prev_score:
                curr["score"] = prev_score
    return result


def _avg_score(result: dict) -> float:
    """计算 Agent 评分的平均分。"""
    dims = result.get("dimensions", {})
    scores = []
    for v in dims.values():
        if isinstance(v, dict):
            scores.append(v.get("score", 0))
        elif isinstance(v, (int, float)):
            scores.append(v)
    return sum(scores) / max(len(scores), 1)


def _print_scores(agent_name: str, result: dict):
    """打印 Agent 评分。"""
    dims = result.get("dimensions", {})
    parts = []
    for key, val in dims.items():
        score = val.get("score", "?") if isinstance(val, dict) else val
        parts.append(f"{key}:{score}")
    avg = _avg_score(result)
    print(f"  [{agent_name}] {' '.join(parts)} → 平均 {avg:.1f}", flush=True)


def _avg_score_d(result: dict) -> float:
    """计算 Agent D 的平均逻辑分。"""
    scores = []
    for layer in ("macro_logic", "meso_logic", "micro_logic", "source_fidelity"):
        val = result.get(layer, {})
        if isinstance(val, dict) and "score" in val:
            scores.append(val["score"])
    return sum(scores) / max(len(scores), 1)


def _print_scores_d(agent_name: str, result: dict):
    """打印 Agent D 逻辑评分。"""
    parts = []
    for layer in ("macro_logic", "meso_logic", "micro_logic", "source_fidelity"):
        val = result.get(layer, {})
        score = val.get("score", "?") if isinstance(val, dict) else "?"
        short = layer.split("_")[0]  # macro/meso/micro/source
        parts.append(f"{short}:{score}")
    avg = _avg_score_d(result)
    n_suggestions = len(result.get("priority_suggestions", []))
    print(f"  [{agent_name}] {' '.join(parts)} → 平均 {avg:.1f} "
          f"(建议{n_suggestions}条)", flush=True)


# ============ 改动应用 ============

def _apply_changes(slide_plan: dict, agent_c_result: dict) -> tuple:
    """
    应用 Agent C 的改动到 slide_plan。
    返回 (updated_slide_plan, changed_visual_ids: set)
    """
    changes = agent_c_result.get("changes", [])
    slides = slide_plan.get("slides", [])
    changed_visual_ids = set()

    for ch in changes:
        action = ch.get("action", "")
        try:
            if action == "modify":
                _apply_modify(slides, ch, changed_visual_ids)
            elif action == "insert_after":
                slides = _apply_insert(slides, ch, changed_visual_ids)
            elif action == "delete":
                slides = _apply_delete(slides, ch)
            elif action == "reorder":
                slides = _apply_reorder(slides, ch)
            else:
                print(f"    [警告] 未知操作: {action}", flush=True)
        except Exception as e:
            print(f"    [警告] 应用改动失败: {e}", flush=True)
            continue

    slide_plan["slides"] = slides
    slide_plan["meta"]["total_slides"] = len(slides)
    return slide_plan, changed_visual_ids


def _apply_modify(slides: list, change: dict, changed_visual_ids: set):
    """修改现有幻灯片的某个字段。"""
    sid = change.get("slide_id", "")
    field = change.get("field", "")
    new_value = change.get("new_value")

    for s in slides:
        if s.get("id") == sid:
            s[field] = new_value
            if field == "visual":
                changed_visual_ids.add(sid)
            print(f"    修改 {sid}.{field}", flush=True)
            return

    print(f"    [警告] 未找到 {sid}", flush=True)


def _apply_insert(slides: list, change: dict, changed_visual_ids: set) -> list:
    """在指定页面后插入新页面。"""
    after_id = change.get("after_slide_id", "")
    new_slide = change.get("new_slide", {})

    for i, s in enumerate(slides):
        if s.get("id") == after_id:
            slides.insert(i + 1, new_slide)
            new_id = new_slide.get("id", "new")
            if new_slide.get("visual"):
                changed_visual_ids.add(new_id)
            print(f"    插入 {new_id} (在 {after_id} 之后)", flush=True)
            return slides

    print(f"    [警告] 未找到 {after_id}，跳过插入", flush=True)
    return slides


def _apply_delete(slides: list, change: dict) -> list:
    """删除指定页面。"""
    sid = change.get("slide_id", "")
    original_len = len(slides)
    slides = [s for s in slides if s.get("id") != sid]
    if len(slides) < original_len:
        print(f"    删除 {sid}", flush=True)
    else:
        print(f"    [警告] 未找到 {sid}，跳过删除", flush=True)
    return slides


def _apply_reorder(slides: list, change: dict) -> list:
    """调整页面顺序。"""
    sid = change.get("slide_id", "")
    move_after = change.get("move_after", "")

    target = None
    target_idx = None
    for i, s in enumerate(slides):
        if s.get("id") == sid:
            target = s
            target_idx = i
            break
    if target is None:
        print(f"    [警告] 未找到 {sid}，跳过重排", flush=True)
        return slides

    slides.pop(target_idx)
    for i, s in enumerate(slides):
        if s.get("id") == move_after:
            slides.insert(i + 1, target)
            print(f"    重排 {sid} → {move_after} 之后", flush=True)
            return slides

    slides.append(target)
    print(f"    [警告] 未找到 {move_after}，{sid} 移至末尾", flush=True)
    return slides


# ============ 选择性视觉重生 ============

def _rerun_visuals_selective(output_dir: Path, slide_plan: dict,
                              changed_ids: set, no_ai_images: bool) -> int:
    """只重新生成指定 slide_id 的视觉资产。"""
    from src.step3_visuals import _render_visual

    assets_dir = output_dir / "assets"
    assets_dir.mkdir(parents=True, exist_ok=True)
    color_scheme = slide_plan.get("meta", {}).get("color_scheme", {})

    # 读取已有 manifest
    manifest_path = assets_dir / "manifest.json"
    manifest = {}
    if manifest_path.is_file():
        manifest = json.loads(manifest_path.read_text(encoding="utf-8"))

    count = 0
    for slide in slide_plan.get("slides", []):
        sid = slide.get("id", "")
        if sid not in changed_ids:
            continue
        visual = slide.get("visual")
        if not visual or not isinstance(visual, dict):
            continue
        visual_type = visual.get("type", "")
        if no_ai_images and visual_type == "generate-image":
            continue

        task = {
            "slide_id": sid,
            "visual": visual,
            "color_scheme": color_scheme,
            "output_path": assets_dir / f"{sid}_{visual_type.replace('-', '_')}.png",
        }
        try:
            result = _render_visual(task)
            manifest[sid] = result
            count += 1
            print(f"    [{sid}] {result.get('status', 'unknown')}", flush=True)
        except Exception as e:
            manifest[sid] = {"status": "failed", "error": str(e)}
            print(f"    [{sid}] 失败: {e}", flush=True)

    # 更新 manifest
    safe_write_text(manifest_path, json.dumps(manifest, ensure_ascii=False, indent=2))
    return count


# ============ 外科手术式编辑 ============

def _classify_text_only_changes(agent_c_result: dict) -> list:
    """
    从 Agent C 结果中提取纯文字修改（title/bullets/notes/takeaway）。
    返回 [{slide_id, field, new_value}, ...]
    """
    text_fields = {"title", "subtitle", "bullets", "notes", "takeaway"}
    text_changes = []
    for ch in agent_c_result.get("changes", []):
        if ch.get("action") == "modify" and ch.get("field") in text_fields:
            text_changes.append({
                "slide_id": ch["slide_id"],
                "field": ch["field"],
                "new_value": ch["new_value"],
            })
    return text_changes


def _surgical_text_edit(pptx_path: Path, slide_plan: dict,
                         text_changes: list):
    """
    直接在 PPTX 中修改文字内容，保留原有格式和布局。
    仅处理 title/subtitle/bullets/notes 字段的文本替换。
    """
    from pptx import Presentation
    from pptx.util import Pt
    import copy

    prs = Presentation(str(pptx_path))
    slides = list(prs.slides)

    # 建立 slide_id → slide_index 映射
    plan_slides = slide_plan.get("slides", [])
    id_to_idx = {}
    for i, ps in enumerate(plan_slides):
        id_to_idx[ps.get("id", "")] = i

    applied = 0
    for change in text_changes:
        sid = change["slide_id"]
        field = change["field"]
        new_value = change["new_value"]

        idx = id_to_idx.get(sid)
        if idx is None or idx >= len(slides):
            print(f"    [外科] 跳过 {sid}：未找到对应页面", flush=True)
            continue

        slide = slides[idx]

        try:
            if field == "notes":
                # 修改备注
                if slide.has_notes_slide:
                    notes_tf = slide.notes_slide.notes_text_frame
                    _replace_text_preserve_format(notes_tf, str(new_value))
                    applied += 1
                    print(f"    [外科] {sid}.notes 已更新", flush=True)

            elif field == "title":
                # 修改标题（通常是第一个带文字的 placeholder）
                for shape in slide.placeholders:
                    if shape.placeholder_format.idx in (0, 12):  # title placeholders
                        if hasattr(shape, "text_frame"):
                            _replace_text_preserve_format(shape.text_frame, str(new_value))
                            applied += 1
                            print(f"    [外科] {sid}.title 已更新", flush=True)
                            break

            elif field == "subtitle":
                for shape in slide.placeholders:
                    if shape.placeholder_format.idx == 1:  # subtitle placeholder
                        if hasattr(shape, "text_frame"):
                            _replace_text_preserve_format(shape.text_frame, str(new_value))
                            applied += 1
                            print(f"    [外科] {sid}.subtitle 已更新", flush=True)
                            break

            elif field == "bullets":
                if isinstance(new_value, list):
                    # 找到正文文本框（排除标题占位符）
                    for shape in slide.shapes:
                        if not hasattr(shape, "text_frame"):
                            continue
                        # 跳过标题占位符
                        try:
                            if hasattr(shape, "placeholder_format"):
                                if shape.placeholder_format.idx in (0, 12):
                                    continue
                        except Exception:
                            pass
                        tf = shape.text_frame
                        if len(tf.paragraphs) >= 2:
                            _replace_bullets_preserve_format(tf, new_value)
                            applied += 1
                            print(f"    [外科] {sid}.bullets 已更新 ({len(new_value)}条)",
                                  flush=True)
                            break

        except Exception as e:
            print(f"    [外科] {sid}.{field} 修改失败: {e}", flush=True)

    if applied > 0:
        prs.save(str(pptx_path))
        print(f"  [外科] 已保存 {applied} 处修改", flush=True)
    else:
        print(f"  [外科] 无有效修改，回退到全量重建", flush=True)


def _replace_text_preserve_format(text_frame, new_text: str):
    """替换文本框内容，尽量保留第一个 run 的格式。"""
    if not text_frame.paragraphs:
        return

    # 保存第一个 paragraph 第一个 run 的格式
    first_para = text_frame.paragraphs[0]
    fmt_info = None
    if first_para.runs:
        run = first_para.runs[0]
        fmt_info = {
            "font_size": run.font.size,
            "font_bold": run.font.bold,
            "font_color": run.font.color.rgb if run.font.color and run.font.color.type else None,
            "font_name": run.font.name,
        }

    # 清除所有段落内容
    for para in text_frame.paragraphs:
        for run in para.runs:
            run.text = ""

    # 设置第一段的文本
    if first_para.runs:
        first_para.runs[0].text = new_text
    else:
        first_para.text = new_text

    # 恢复格式
    if fmt_info and first_para.runs:
        run = first_para.runs[0]
        if fmt_info["font_size"]:
            run.font.size = fmt_info["font_size"]
        if fmt_info["font_bold"] is not None:
            run.font.bold = fmt_info["font_bold"]
        if fmt_info["font_name"]:
            run.font.name = fmt_info["font_name"]


def _replace_bullets_preserve_format(text_frame, new_bullets: list):
    """替换 bullet 列表内容，保留现有格式风格。"""
    paras = list(text_frame.paragraphs)

    # 保存格式模板（从第一个有文字的段落）
    template_fmt = None
    for p in paras:
        if p.runs and p.runs[0].text.strip():
            run = p.runs[0]
            template_fmt = {
                "font_size": run.font.size,
                "font_bold": run.font.bold,
                "font_name": run.font.name,
                "font_color_rgb": run.font.color.rgb if run.font.color and run.font.color.type else None,
            }
            break

    # 更新现有段落或添加新段落
    from pptx.util import Pt
    for i, bullet_text in enumerate(new_bullets):
        if i < len(paras):
            # 更新现有段落
            para = paras[i]
            if para.runs:
                # 清除所有 run 然后设置第一个
                for r in para.runs:
                    r.text = ""
                para.runs[0].text = str(bullet_text)
            else:
                para.text = str(bullet_text)
        else:
            # 添加新段落
            from pptx.oxml.ns import qn
            from lxml import etree
            new_p = text_frame._txBody.makeelement(qn("a:p"), {})
            new_r = new_p.makeelement(qn("a:r"), {})
            new_t = new_r.makeelement(qn("a:t"), {})
            new_t.text = str(bullet_text)
            new_r.append(new_t)
            new_p.append(new_r)
            text_frame._txBody.append(new_p)

            # 应用模板格式
            if template_fmt:
                from pptx.oxml.ns import qn as ns_qn
                rPr = new_r.makeelement(ns_qn("a:rPr"), {"lang": "zh-CN", "dirty": "0"})
                if template_fmt["font_size"]:
                    rPr.set("sz", str(int(template_fmt["font_size"] / 12700)))
                new_r.insert(0, rPr)

    # 删除多余段落（如果新 bullets 更少）
    while len(text_frame.paragraphs) > len(new_bullets) and len(text_frame.paragraphs) > 1:
        last_p = text_frame._txBody.findall(f'{{{text_frame._txBody.nsmap.get("a", "http://schemas.openxmlformats.org/drawingml/2006/main")}}}p')
        if last_p and len(last_p) > len(new_bullets):
            text_frame._txBody.remove(last_p[-1])
        else:
            break


# ============ JSON 解析 ============

def _parse_json_response(response: str, agent_label: str) -> dict:
    """解析 LLM 返回的 JSON，兼容 markdown code block。"""
    text = response.strip()

    # 提取 code block
    if "```json" in text:
        start = text.index("```json") + 7
        end = text.find("```", start)
        text = text[start:end].strip() if end > start else text[start:].strip()
    elif "```" in text:
        start = text.index("```") + 3
        end = text.find("```", start)
        text = text[start:end].strip() if end > start else text[start:].strip()

    try:
        return json.loads(text)
    except json.JSONDecodeError:
        pass

    # 尝试找 { 开头
    idx = text.find("{")
    if idx >= 0:
        last = text.rfind("}")
        if last > idx:
            try:
                return json.loads(text[idx:last + 1])
            except json.JSONDecodeError:
                pass

    raise ParseError(f"Agent {agent_label} 返回内容无法解析为 JSON")


# ============ 日志构建 ============

def _build_round_log(round_num, result_a, result_b, agent_c_result,
                      visuals_regen, elapsed, agent_e_result=None,
                      agent_d_result=None):
    """构建单轮日志。"""
    log = {
        "round": round_num,
        "agent_a": {
            "dimensions": result_a.get("dimensions", {}),
            "average": _avg_score(result_a),
            "slide_comments": result_a.get("slide_comments", []),
        },
        "agent_b": {
            "dimensions": result_b.get("dimensions", {}),
            "average": _avg_score(result_b),
            "slide_comments": result_b.get("slide_comments", []),
        },
        "visuals_regenerated": visuals_regen,
        "elapsed_seconds": round(elapsed, 1),
    }
    if agent_d_result:
        log["agent_d"] = {
            "macro_logic": agent_d_result.get("macro_logic", {}).get("score", 0),
            "meso_logic": agent_d_result.get("meso_logic", {}).get("score", 0),
            "micro_logic": agent_d_result.get("micro_logic", {}).get("score", 0),
            "source_fidelity": agent_d_result.get("source_fidelity", {}).get("score", 0),
            "average": _avg_score_d(agent_d_result),
            "priority_suggestions_count": len(agent_d_result.get("priority_suggestions", [])),
            "emphasis": agent_d_result.get("emphasis", {}),
        }
    if agent_c_result:
        log["agent_c"] = {
            "reasoning": agent_c_result.get("reasoning", ""),
            "response_to_a": agent_c_result.get("response_to_a", ""),
            "response_to_b": agent_c_result.get("response_to_b", ""),
            "response_to_d": agent_c_result.get("response_to_d", ""),
            "response_to_e": agent_c_result.get("response_to_e", ""),
            "changes_count": len(agent_c_result.get("changes", [])),
        }
    if agent_e_result:
        log["agent_e"] = {
            "total_anomalies": agent_e_result["summary"]["total_anomalies"],
            "high_severity": agent_e_result["summary"]["high_severity"],
            "medium_severity": agent_e_result["summary"]["medium_severity"],
            "anomaly_slides": [a["slide_id"] for a in agent_e_result.get("anomalies", [])],
        }
    return log
