# -*- coding: utf-8 -*-
"""
Agent E: PPTX 视觉排版检测器。
对已生成的 PPTX 文件进行逐页程序化分析，检测6类排版异常。
"""
import math
from pathlib import Path
from typing import Optional

import src  # noqa: F401


def inspect_pptx(pptx_path: Path, slide_plan: dict = None) -> dict:
    """
    对 PPTX 文件进行逐页视觉排版检测。

    Returns:
        {
            "slide_width": float,  # inches
            "slide_height": float,
            "total_slides": int,
            "anomalies": [
                {
                    "slide_index": 1,        # 1-based
                    "slide_id": "s01",        # from slide_plan
                    "layout": "title_content", # from slide_plan
                    "issues": [
                        {
                            "type": "excessive_whitespace",
                            "severity": "high",     # high/medium/low
                            "detail": "页面留白占比62%，仅有文字无配图，疑似图片丢失",
                            "metrics": {"whitespace_pct": 62, "has_image": false}
                        }
                    ]
                }
            ],
            "summary": {
                "total_anomalies": 12,
                "high_severity": 3,
                "medium_severity": 5,
                "low_severity": 4,
                "most_common_type": "excessive_whitespace"
            }
        }
    """
    from pptx import Presentation
    from pptx.util import Inches, Emu

    prs = Presentation(str(pptx_path))
    slide_width = prs.slide_width.inches
    slide_height = prs.slide_height.inches
    slide_area = slide_width * slide_height

    # 从 slide_plan 获取 slide_id 和 layout 映射
    plan_slides = (slide_plan or {}).get("slides", [])
    plan_map = {}
    for i, ps in enumerate(plan_slides):
        plan_map[i] = {
            "slide_id": ps.get("id", f"s{i+1:02d}"),
            "layout": ps.get("layout", "unknown"),
        }

    anomalies = []

    for idx, slide in enumerate(prs.slides):
        slide_num = idx + 1
        plan_info = plan_map.get(idx, {"slide_id": f"s{slide_num:02d}", "layout": "unknown"})
        slide_id = plan_info["slide_id"]
        layout = plan_info["layout"]

        # 跳过封面和结束页（设计自由度高）
        if layout in ("cover", "end"):
            continue

        issues = []

        # 收集该页所有形状的信息
        shapes_info = _analyze_shapes(slide, slide_width, slide_height)

        # ── 检测 1: 大面积留白 + 仅文字 ──
        _check_excessive_whitespace(shapes_info, slide_area, issues)

        # ── 检测 2: 少量文字无配图 ──
        _check_sparse_content(shapes_info, layout, slide_area, issues)

        # ── 检测 3: 信息图不饱满/字体偏小 ──
        _check_infographic_quality(shapes_info, slide_area, issues)

        # ── 检测 4: 图文混排图片偏小 ──
        _check_image_too_small(shapes_info, slide_area, layout, issues)

        # ── 检测 5: 图片拉伸变形 ──
        _check_image_distortion(slide, shapes_info, issues)

        # ── 检测 6: 文字颜色与背景区分不清 ──
        _check_color_contrast(slide, shapes_info, issues)

        # ── 检测 7: 内容偏向一侧，另一侧大片空白 ──
        _check_unbalanced_layout(shapes_info, slide_width, slide_height, layout, issues)

        # ── 检测 8: 视觉单调性（借鉴PPTEval Style评分标准）──
        _check_visual_monotony(shapes_info, layout, issues)

        if issues:
            anomalies.append({
                "slide_index": slide_num,
                "slide_id": slide_id,
                "layout": layout,
                "issues": issues,
            })

    # 汇总
    all_issues = [iss for a in anomalies for iss in a["issues"]]
    type_counts = {}
    for iss in all_issues:
        type_counts[iss["type"]] = type_counts.get(iss["type"], 0) + 1

    # 视觉单调性统计
    monotony_count = sum(
        1 for iss in all_issues if iss["type"] == "visual_monotony"
    )

    summary = {
        "total_anomalies": len(all_issues),
        "high_severity": sum(1 for i in all_issues if i["severity"] == "high"),
        "medium_severity": sum(1 for i in all_issues if i["severity"] == "medium"),
        "low_severity": sum(1 for i in all_issues if i["severity"] == "low"),
        "most_common_type": max(type_counts, key=type_counts.get) if type_counts else "none",
        "visual_monotony_slides": monotony_count,
    }

    return {
        "slide_width": slide_width,
        "slide_height": slide_height,
        "total_slides": len(prs.slides),
        "anomalies": anomalies,
        "summary": summary,
    }


def _analyze_shapes(slide, slide_width: float, slide_height: float) -> dict:
    """分析一个 slide 的所有形状，返回汇总信息。"""
    from pptx.util import Emu
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    info = {
        "total_shapes": 0,
        "text_shapes": [],       # 含文字的形状
        "image_shapes": [],      # 图片形状
        "chart_shapes": [],      # 图表形状
        "other_shapes": [],      # 装饰/其他
        "decoration_shapes": [], # 几何装饰元素（矩形、圆形、线条等）
        "total_text_chars": 0,
        "has_image": False,
        "has_chart": False,
        "covered_area": 0.0,     # 所有形状覆盖的总面积（inches²）
        "image_area": 0.0,       # 图片总面积
        "text_area": 0.0,        # 文字框总面积
        "min_font_size": 999,
        "font_sizes": [],
        "text_colors": set(),    # 文字使用的颜色集合（RGB 元组）
        "fill_colors": set(),    # 形状填充色集合
        "has_background_fill": False,  # 是否有非白色背景
    }

    for shape in slide.shapes:
        info["total_shapes"] += 1

        # 位置和尺寸
        try:
            left = shape.left / 914400.0 if shape.left else 0   # EMU to inches
            top = shape.top / 914400.0 if shape.top else 0
            width = shape.width / 914400.0 if shape.width else 0
            height = shape.height / 914400.0 if shape.height else 0
        except Exception:
            width = height = 0

        shape_area = width * height

        # 判断形状类型
        is_image = False
        is_text = False

        try:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                is_image = True
            elif shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
                # Placeholder 可能含图片或文字
                if hasattr(shape, "image"):
                    is_image = True
        except Exception:
            pass

        if is_image:
            info["image_shapes"].append({
                "left": left, "top": top, "width": width, "height": height,
                "area": shape_area,
            })
            info["has_image"] = True
            info["image_area"] += shape_area

        # 文字检测
        if hasattr(shape, "text_frame"):
            try:
                text = shape.text_frame.text.strip()
                if text:
                    is_text = True
                    info["text_shapes"].append({
                        "left": left, "top": top, "width": width, "height": height,
                        "area": shape_area, "text_len": len(text), "text": text[:100],
                    })
                    info["total_text_chars"] += len(text)
                    info["text_area"] += shape_area

                    # 字号和颜色检测
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if run.font.size:
                                pt = run.font.size.pt
                                info["font_sizes"].append(pt)
                                if pt < info["min_font_size"]:
                                    info["min_font_size"] = pt
                            # 收集文字颜色
                            try:
                                if run.font.color and run.font.color.type is not None:
                                    if hasattr(run.font.color, 'rgb') and run.font.color.rgb:
                                        rgb_str = str(run.font.color.rgb)
                                        info["text_colors"].add(rgb_str)
                            except Exception:
                                pass
            except Exception:
                pass

        # 形状填充色检测（几何装饰元素）
        try:
            if hasattr(shape, "fill"):
                fill = shape.fill
                if fill.type is not None:
                    try:
                        fc = fill.fore_color
                        if hasattr(fc, 'rgb') and fc.rgb:
                            info["fill_colors"].add(str(fc.rgb))
                    except Exception:
                        pass
                    # 非文字非图片的有填充形状 = 装饰元素
                    if not is_image and not is_text and shape_area > 0.1:
                        info["decoration_shapes"].append({
                            "left": left, "top": top, "width": width, "height": height,
                            "area": shape_area,
                        })
        except Exception:
            pass

        # 图表检测
        if hasattr(shape, "chart"):
            info["has_chart"] = True
            info["chart_shapes"].append({
                "left": left, "top": top, "width": width, "height": height,
                "area": shape_area,
            })

        # 只有包含实际内容的形状（文字/图片/图表）或小型装饰元素才计入覆盖面积
        # 空的大面积占位框不应虚增覆盖率
        if is_image or is_text or info["has_chart"]:
            info["covered_area"] += shape_area
        elif shape_area < 3.0:
            # 小型装饰元素（线条、小色块等）正常计入
            info["covered_area"] += shape_area
        else:
            # 大面积空形状（可能是空占位框），记入 other 但不计入 covered_area
            info["other_shapes"].append({
                "left": left, "top": top, "width": width, "height": height,
                "area": shape_area,
            })

    if info["min_font_size"] == 999:
        info["min_font_size"] = None

    # 转换 set → list（JSON 序列化兼容）
    info["text_colors"] = list(info["text_colors"])
    info["fill_colors"] = list(info["fill_colors"])

    return info


def _check_excessive_whitespace(shapes_info: dict, slide_area: float, issues: list):
    """检测1: 页面大量留白，分两种情况检测。"""
    if slide_area <= 0:
        return

    # 计算有效内容覆盖率
    content_area = shapes_info["covered_area"]
    coverage_pct = (content_area / slide_area) * 100
    whitespace_pct = max(0, 100 - coverage_pct)

    has_visual = shapes_info["has_image"] or shapes_info["has_chart"]

    # ── 情况A：纯文字页面，留白>40% ──
    if (whitespace_pct > 40
            and not has_visual
            and shapes_info["total_text_chars"] > 10):
        severity = "high" if whitespace_pct > 60 else "medium"
        issues.append({
            "type": "excessive_whitespace",
            "severity": severity,
            "detail": f"页面留白占比约{whitespace_pct:.0f}%，仅有文字({shapes_info['total_text_chars']}字)无配图，"
                      f"可能是配图丢失或布局不合理",
            "metrics": {
                "whitespace_pct": round(whitespace_pct, 1),
                "has_image": False,
                "text_chars": shapes_info["total_text_chars"],
            },
        })

    # ── 情况B：有图/表但内容覆盖仍然很低（图片太小+文字太少，大片留白） ──
    elif (has_visual
          and whitespace_pct > 50
          and shapes_info["total_text_chars"] > 10):
        # 文字+图片加起来仍覆盖不到50%，说明布局不饱满
        img_pct = (shapes_info["image_area"] / slide_area) * 100
        txt_pct = (shapes_info["text_area"] / slide_area) * 100
        severity = "high" if whitespace_pct > 65 else "medium"
        issues.append({
            "type": "excessive_whitespace",
            "severity": severity,
            "detail": f"页面留白占比约{whitespace_pct:.0f}%，图片占{img_pct:.0f}%+文字占{txt_pct:.0f}%，"
                      f"内容覆盖不足，布局不饱满（建议扩大图表或补充内容）",
            "metrics": {
                "whitespace_pct": round(whitespace_pct, 1),
                "has_image": shapes_info["has_image"],
                "image_area_pct": round(img_pct, 1),
                "text_area_pct": round(txt_pct, 1),
                "text_chars": shapes_info["total_text_chars"],
            },
        })


def _check_sparse_content(shapes_info: dict, layout: str, slide_area: float, issues: list):
    """检测2: 内容稀疏——文字过少或整体内容覆盖率低。"""
    # 跳过本身就是轻内容的布局
    if layout in ("section_break", "quote", "image_full"):
        return

    text_chars = shapes_info["total_text_chars"]
    has_visual = shapes_info["has_image"] or shapes_info["has_chart"]

    # ── 情况A：文字极少且无视觉元素 ──
    if text_chars < 50 and not has_visual and shapes_info["total_shapes"] < 5:
        issues.append({
            "type": "sparse_content",
            "severity": "high",
            "detail": f"页面仅{text_chars}字文字，无图片/图表/信息图，内容过于稀疏",
            "metrics": {
                "text_chars": text_chars,
                "shape_count": shapes_info["total_shapes"],
            },
        })

    # ── 情况B：文字较少（<100字），即使有图但整体内容覆盖率低 ──
    elif text_chars < 100 and slide_area > 0:
        coverage_pct = (shapes_info["covered_area"] / slide_area) * 100
        if coverage_pct < 40:
            issues.append({
                "type": "sparse_content",
                "severity": "medium",
                "detail": f"页面仅{text_chars}字文字，内容覆盖率{coverage_pct:.0f}%，"
                          f"页面显得空旷（建议补充内容或合并到相邻页面）",
                "metrics": {
                    "text_chars": text_chars,
                    "coverage_pct": round(coverage_pct, 1),
                },
            })


def _check_infographic_quality(shapes_info: dict, slide_area: float, issues: list):
    """检测3: 信息图不饱满（尺寸小于页面30%）或字体偏小。"""
    if not shapes_info["has_image"]:
        return

    # 检查最大图片是否够大
    if shapes_info["image_shapes"]:
        max_img = max(shapes_info["image_shapes"], key=lambda x: x["area"])
        img_pct = (max_img["area"] / slide_area) * 100 if slide_area > 0 else 0

        if img_pct < 15:
            issues.append({
                "type": "infographic_too_small",
                "severity": "medium",
                "detail": f"最大图片/信息图仅占页面{img_pct:.0f}%，不够饱满，"
                          f"信息图尺寸偏小导致周围留白过多",
                "metrics": {
                    "max_image_pct": round(img_pct, 1),
                    "image_width": round(max_img["width"], 1),
                    "image_height": round(max_img["height"], 1),
                },
            })

    # 检查字号偏小
    if shapes_info["min_font_size"] and shapes_info["min_font_size"] < 12:
        issues.append({
            "type": "font_too_small",
            "severity": "medium",
            "detail": f"页面最小字号仅{shapes_info['min_font_size']:.0f}pt，"
                      f"低于12pt阅读舒适度下限",
            "metrics": {
                "min_font_pt": shapes_info["min_font_size"],
            },
        })


def _check_image_too_small(shapes_info: dict, slide_area: float, layout: str, issues: list):
    """检测4: 图文混排页面，图片占页面比例偏小。"""
    # 所有应该有图文搭配的布局类型
    mixed_layouts = (
        "title_content", "infographic", "key_insight", "two_column",
        "data_chart", "comparison", "methodology", "architecture",
    )
    if layout not in mixed_layouts:
        return

    if not shapes_info["has_image"]:
        return

    img_pct = (shapes_info["image_area"] / slide_area) * 100 if slide_area > 0 else 0

    # data_chart 布局要求图表至少占25%（图表是核心内容）
    if layout == "data_chart":
        threshold = 25
    else:
        threshold = 20

    if img_pct < threshold and shapes_info["total_text_chars"] > 50:
        issues.append({
            "type": "image_too_small_in_layout",
            "severity": "high" if img_pct < 10 else "medium",
            "detail": f"图文混排页面({layout})图片仅占{img_pct:.0f}%（建议≥{threshold}%），"
                      f"图片偏小导致页面留白过多",
            "metrics": {
                "image_area_pct": round(img_pct, 1),
                "text_chars": shapes_info["total_text_chars"],
                "layout": layout,
            },
        })


def _check_image_distortion(slide, shapes_info: dict, issues: list):
    """检测5: 图片拉伸变形。"""
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    for shape in slide.shapes:
        try:
            if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
                continue
        except Exception:
            continue

        try:
            img = shape.image
            # 原始图片宽高
            blob = img.blob
            from PIL import Image
            import io
            pil_img = Image.open(io.BytesIO(blob))
            orig_w, orig_h = pil_img.size
            if orig_w == 0 or orig_h == 0:
                continue

            orig_ratio = orig_w / orig_h

            # 显示宽高
            disp_w = shape.width / 914400.0  # EMU to inches
            disp_h = shape.height / 914400.0
            if disp_h == 0:
                continue

            disp_ratio = disp_w / disp_h

            # 变形比（>20% 差异视为变形）
            ratio_diff = abs(orig_ratio - disp_ratio) / orig_ratio
            if ratio_diff > 0.20:
                issues.append({
                    "type": "image_distortion",
                    "severity": "high",
                    "detail": f"图片宽高比变形{ratio_diff*100:.0f}%"
                              f"（原始{orig_w}x{orig_h}，显示比例{disp_ratio:.2f}），"
                              f"视觉拉伸明显",
                    "metrics": {
                        "original_ratio": round(orig_ratio, 3),
                        "display_ratio": round(disp_ratio, 3),
                        "distortion_pct": round(ratio_diff * 100, 1),
                    },
                })
        except Exception:
            continue


def _check_color_contrast(slide, shapes_info: dict, issues: list):
    """检测6: 文字颜色与背景色区分不清。"""
    # 简化检测：检查是否有白色/浅色文字在浅色背景上
    # 或深色文字在深色背景上

    bg_is_dark = _slide_has_dark_background(slide)

    for ts in shapes_info["text_shapes"]:
        # 检查文字中是否有与背景混淆的颜色
        pass  # 完整颜色对比需要遍历每个run的font.color

    # 简化版：遍历形状检查
    for shape in slide.shapes:
        if not hasattr(shape, "text_frame"):
            continue
        try:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if not run.text.strip():
                        continue
                    color = _get_run_color(run)
                    if color is None:
                        continue

                    r, g, b = color
                    text_luminance = 0.299 * r + 0.587 * g + 0.114 * b

                    if bg_is_dark:
                        # 深色背景上的深色文字
                        if text_luminance < 80:
                            issues.append({
                                "type": "poor_color_contrast",
                                "severity": "high",
                                "detail": f"深色背景上出现深色文字(RGB:{r},{g},{b})，"
                                          f"文字几乎不可读",
                                "metrics": {
                                    "text_rgb": [r, g, b],
                                    "text_luminance": round(text_luminance, 1),
                                    "bg_is_dark": True,
                                },
                            })
                            return  # 只报一次
                    else:
                        # 浅色背景上的浅色文字
                        if text_luminance > 220 and len(run.text.strip()) > 3:
                            issues.append({
                                "type": "poor_color_contrast",
                                "severity": "high",
                                "detail": f"浅色背景上出现浅色文字(RGB:{r},{g},{b})，"
                                          f"文字几乎不可读",
                                "metrics": {
                                    "text_rgb": [r, g, b],
                                    "text_luminance": round(text_luminance, 1),
                                    "bg_is_dark": False,
                                },
                            })
                            return  # 只报一次
        except Exception:
            continue


def _check_unbalanced_layout(shapes_info: dict, slide_width: float,
                              slide_height: float, layout: str, issues: list):
    """检测7: 内容集中在页面一侧，另一侧大片空白。"""
    if layout in ("cover", "end", "section_break", "image_full"):
        return

    # 收集所有有意义的形状（文字+图片+图表）的位置
    meaningful = (shapes_info["text_shapes"] +
                  shapes_info["image_shapes"] +
                  shapes_info["chart_shapes"])
    if len(meaningful) < 1 or slide_width <= 0:
        return

    mid_x = slide_width / 2.0

    # 过滤掉跨全宽的形状（标题、底栏等），它们不参与左右平衡判断
    # 定义：宽度超过页面宽度70%的形状视为全宽
    content_shapes = [s for s in meaningful if s["width"] < slide_width * 0.7]
    if not content_shapes:
        return

    # 计算左半区和右半区的内容面积
    left_area = 0.0
    right_area = 0.0
    for s in content_shapes:
        center_x = s["left"] + s["width"] / 2.0
        if center_x < mid_x:
            left_area += s["area"]
        else:
            right_area += s["area"]

    total_content = left_area + right_area
    if total_content <= 0:
        return

    # 如果一侧内容面积 < 总内容的10%，视为不平衡
    left_ratio = left_area / total_content
    right_ratio = right_area / total_content

    empty_side = None
    if left_ratio < 0.10 and right_area > 0:
        empty_side = "左侧"
    elif right_ratio < 0.10 and left_area > 0:
        empty_side = "右侧"

    if empty_side:
        issues.append({
            "type": "unbalanced_layout",
            "severity": "high",
            "detail": f"页面{empty_side}几乎无内容（仅占总内容{min(left_ratio, right_ratio)*100:.0f}%），"
                      f"形成半页空白，建议补充图表/信息图或调整为全宽布局",
            "metrics": {
                "left_content_pct": round(left_ratio * 100, 1),
                "right_content_pct": round(right_ratio * 100, 1),
                "empty_side": empty_side,
            },
        })


def _check_visual_monotony(shapes_info: dict, layout: str, issues: list):
    """检测8: 视觉单调性——页面缺乏视觉元素，纯黑白文字无装饰。
    借鉴 PPTEval Style 评分标准：
    - 1分：风格冲突影响阅读
    - 2分：纯黑白单调，可读但无吸引力
    - 3分：有基础配色但缺视觉元素（图标、背景、几何装饰）
    - 4分：配色和谐+部分视觉元素，小瑕疵
    - 5分：视觉元素丰富，增强整体吸引力
    """
    # 跳过天然轻内容布局
    if layout in ("section_break", "quote", "image_full", "cover", "end"):
        return

    has_image = shapes_info["has_image"]
    has_chart = shapes_info["has_chart"]
    n_decorations = len(shapes_info.get("decoration_shapes", []))
    n_fill_colors = len(shapes_info.get("fill_colors", []))
    n_text_colors = len(shapes_info.get("text_colors", []))
    total_text = shapes_info["total_text_chars"]

    # 计算视觉元素计数
    visual_elements = 0
    if has_image:
        visual_elements += 2  # 图片权重高
    if has_chart:
        visual_elements += 2
    visual_elements += min(n_decorations, 3)  # 装饰元素最多贡献3分
    if n_fill_colors >= 2:
        visual_elements += 1  # 有多色填充
    if n_text_colors >= 2:
        visual_elements += 1  # 文字不只是黑色

    # 只有有一定文字量的内容页才检测单调性
    if total_text < 20:
        return

    # 判断是否单调
    is_monotone = (
        visual_elements <= 1
        and not has_image
        and not has_chart
        and n_decorations == 0
    )

    # 轻度单调：有文字无任何视觉辅助
    is_plain = (
        visual_elements <= 2
        and not has_image
        and not has_chart
    )

    if is_monotone:
        # 纯文字零装饰（PPTEval Style 2分级别）
        color_desc = "纯黑白" if n_text_colors <= 1 and n_fill_colors == 0 else "少量颜色"
        issues.append({
            "type": "visual_monotony",
            "severity": "medium",
            "detail": f"页面视觉单调（{color_desc}，{total_text}字纯文字，"
                      f"无图表/信息图/几何装饰），建议补充视觉元素提升吸引力",
            "metrics": {
                "visual_element_count": visual_elements,
                "has_image": has_image,
                "has_chart": has_chart,
                "decoration_count": n_decorations,
                "fill_color_count": n_fill_colors,
                "text_color_count": n_text_colors,
                "ppteval_style_level": 2,
            },
        })
    elif is_plain and total_text > 60:
        # 有配色但缺视觉元素（PPTEval Style 3分级别）
        issues.append({
            "type": "visual_monotony",
            "severity": "low",
            "detail": f"页面视觉元素不足（文字{total_text}字，仅{n_decorations}个装饰元素，"
                      f"无图片/图表），建议添加图标、背景纹理或几何形状辅助",
            "metrics": {
                "visual_element_count": visual_elements,
                "has_image": has_image,
                "has_chart": has_chart,
                "decoration_count": n_decorations,
                "fill_color_count": n_fill_colors,
                "text_color_count": n_text_colors,
                "ppteval_style_level": 3,
            },
        })


def _slide_has_dark_background(slide) -> bool:
    """判断幻灯片是否使用深色背景。"""
    try:
        bg = slide.background
        fill = bg.fill
        if fill.type is not None:
            try:
                color = fill.fore_color
                if hasattr(color, 'rgb') and color.rgb:
                    r = int(str(color.rgb)[0:2], 16)
                    g = int(str(color.rgb)[2:4], 16)
                    b = int(str(color.rgb)[4:6], 16)
                    luminance = 0.299 * r + 0.587 * g + 0.114 * b
                    return luminance < 128
            except Exception:
                pass
    except Exception:
        pass
    return False  # 默认白色背景


def _get_run_color(run) -> Optional[tuple]:
    """获取文字 run 的 RGB 颜色值。"""
    try:
        font_color = run.font.color
        if font_color and font_color.type is not None:
            if hasattr(font_color, 'rgb') and font_color.rgb:
                rgb_str = str(font_color.rgb)
                r = int(rgb_str[0:2], 16)
                g = int(rgb_str[2:4], 16)
                b = int(rgb_str[4:6], 16)
                return (r, g, b)
    except Exception:
        pass
    return None


def generate_thumbnail_grid(pptx_path: Path, output_path: Path = None,
                            cols: int = 5, thumb_width: int = 300) -> Path:
    """
    从 PPTX 生成缩略图网格（纯 Python，无需 soffice/pdftoppm）。
    使用 python-pptx 提取形状信息 + PIL 绘制简化的页面缩略图。

    Args:
        pptx_path: PPTX 文件路径
        output_path: 输出图片路径（默认在同目录生成 _thumbnails.jpg）
        cols: 列数
        thumb_width: 每个缩略图宽度（像素）

    Returns:
        生成的缩略图网格文件路径
    """
    from PIL import Image, ImageDraw, ImageFont
    from pptx import Presentation
    from pptx.util import Emu
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    prs = Presentation(str(pptx_path))
    slide_w_inch = prs.slide_width / 914400.0
    slide_h_inch = prs.slide_height / 914400.0
    aspect = slide_h_inch / slide_w_inch
    thumb_height = int(thumb_width * aspect)

    n_slides = len(prs.slides)
    if n_slides == 0:
        raise ValueError("PPTX 无幻灯片")

    rows = (n_slides + cols - 1) // cols
    padding = 15
    label_h = 28
    grid_w = cols * thumb_width + (cols + 1) * padding
    grid_h = rows * (thumb_height + label_h) + (rows + 1) * padding

    grid = Image.new("RGB", (grid_w, grid_h), "white")
    draw = ImageDraw.Draw(grid)

    try:
        font = ImageFont.load_default(size=18)
    except Exception:
        font = ImageFont.load_default()

    # 颜色映射
    bg_color = (245, 245, 250)
    text_color = (60, 60, 80)
    image_color = (200, 220, 240)
    chart_color = (220, 240, 220)
    shape_color = (240, 220, 200)
    border_color = (180, 180, 190)

    for idx, slide in enumerate(prs.slides):
        row, col = idx // cols, idx % cols
        x = col * thumb_width + (col + 1) * padding
        y = row * (thumb_height + label_h) + (row + 1) * padding

        # 标签
        label = f"S{idx + 1}"
        bbox = draw.textbbox((0, 0), label, font=font)
        tw = bbox[2] - bbox[0]
        draw.text((x + (thumb_width - tw) // 2, y + 4), label,
                  fill="black", font=font)

        ty = y + label_h
        # 背景
        draw.rectangle([(x, ty), (x + thumb_width, ty + thumb_height)],
                       fill=bg_color, outline=border_color, width=1)

        # 缩放系数
        sx = thumb_width / slide_w_inch
        sy = thumb_height / slide_h_inch

        # 绘制每个形状的简化表示
        for shape in slide.shapes:
            try:
                sl = (shape.left or 0) / 914400.0
                st = (shape.top or 0) / 914400.0
                sw = (shape.width or 0) / 914400.0
                sh = (shape.height or 0) / 914400.0
            except Exception:
                continue

            px = x + int(sl * sx)
            py = ty + int(st * sy)
            pw = max(2, int(sw * sx))
            ph = max(2, int(sh * sy))

            # 裁剪到缩略图范围内
            px = max(x, min(px, x + thumb_width - 2))
            py = max(ty, min(py, ty + thumb_height - 2))
            pw = min(pw, x + thumb_width - px)
            ph = min(ph, ty + thumb_height - py)

            is_image = False
            try:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    is_image = True
                elif hasattr(shape, "image"):
                    is_image = True
            except Exception:
                pass

            has_chart = hasattr(shape, "chart")
            has_text = False
            text_content = ""
            try:
                if hasattr(shape, "text_frame") and shape.text_frame.text.strip():
                    has_text = True
                    text_content = shape.text_frame.text.strip()[:20]
            except Exception:
                pass

            if is_image:
                # 蓝色块表示图片
                draw.rectangle([(px, py), (px + pw, py + ph)],
                               fill=image_color, outline=(100, 150, 200))
                draw.line([(px, py), (px + pw, py + ph)],
                          fill=(100, 150, 200), width=1)
                draw.line([(px + pw, py), (px, py + ph)],
                          fill=(100, 150, 200), width=1)
            elif has_chart:
                draw.rectangle([(px, py), (px + pw, py + ph)],
                               fill=chart_color, outline=(100, 180, 100))
            elif has_text:
                # 文字区域
                fill = shape_color if ph < thumb_height * 0.3 else (250, 250, 255)
                draw.rectangle([(px, py), (px + pw, py + ph)],
                               fill=fill, outline=(180, 180, 200))
                if pw > 30 and ph > 12:
                    # 绘制文字行模拟
                    line_y = py + 3
                    line_count = min(4, max(1, ph // 8))
                    for _ in range(line_count):
                        if line_y + 3 >= py + ph:
                            break
                        line_w = min(pw - 6, max(10, pw * 2 // 3))
                        draw.line([(px + 3, line_y), (px + 3 + line_w, line_y)],
                                  fill=text_color, width=1)
                        line_y += 7
            else:
                if pw > 5 and ph > 5:
                    draw.rectangle([(px, py), (px + pw, py + ph)],
                                   outline=(200, 200, 210))

    if output_path is None:
        output_path = pptx_path.parent / f"{pptx_path.stem}_thumbnails.jpg"

    grid.save(str(output_path), quality=92)
    return output_path


def format_agent_e_report(inspection: dict) -> str:
    """将检测结果格式化为供 Agent C 使用的结构化文本。"""
    import json

    lines = []
    lines.append(f"## Agent E 视觉排版检测报告")
    lines.append(f"检测页数: {inspection['total_slides']}")

    summary = inspection["summary"]
    lines.append(f"异常总数: {summary['total_anomalies']} "
                 f"(严重{summary['high_severity']} / "
                 f"中等{summary['medium_severity']} / "
                 f"轻微{summary['low_severity']})")

    if summary["total_anomalies"] == 0:
        lines.append("未检测到排版异常。")
        return "\n".join(lines)

    lines.append(f"最常见问题: {summary['most_common_type']}")
    monotony = summary.get("visual_monotony_slides", 0)
    if monotony > 0:
        lines.append(f"视觉单调页面: {monotony}页（缺乏图表/信息图/装饰元素，PPTEval Style≤3分）")
    lines.append("")

    for anomaly in inspection["anomalies"]:
        sid = anomaly["slide_id"]
        layout = anomaly["layout"]
        lines.append(f"### {sid} [{layout}]")
        for iss in anomaly["issues"]:
            severity_icon = {"high": "🔴", "medium": "🟡", "low": "🟢"}.get(iss["severity"], "⚪")
            lines.append(f"  {severity_icon} [{iss['type']}] {iss['detail']}")
        lines.append("")

    return "\n".join(lines)
