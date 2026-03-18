#!/usr/bin/env python3
"""读取PPT内容概览"""
from pptx import Presentation
from pathlib import Path

pptx_path = "output/japan_defense/japan_defense_slides.pptx"

if not Path(pptx_path).exists():
    print(f"文件不存在: {pptx_path}")
    exit(1)

prs = Presentation(pptx_path)

print("="*70)
print("PPT 内容概览: 日本军工企业深度研究")
print("="*70)
print(f"总页数: {len(prs.slides)}")
print(f"幻灯片尺寸: {prs.slide_width.inches:.2f}\" x {prs.slide_height.inches:.2f}\" (16:9)")
print()

for i, slide in enumerate(prs.slides, 1):
    print(f"\n--- 第 {i:2d} 页 ---")
    
    # 获取所有文本
    texts = []
    for shape in slide.shapes:
        if hasattr(shape, 'text') and shape.text.strip():
            text = shape.text.strip().replace('\n', ' ')
            texts.append(text)
    
    if texts:
        # 显示标题（通常第一个文本框）
        title = texts[0][:60] if texts[0] else "[无标题]"
        print(f"  标题: {title}")
        
        # 显示要点（其他文本框）
        for t in texts[1:3]:
            if len(t) > 10:
                print(f"  · {t[:80]}...")
    else:
        print("  [视觉元素/图片页]")

print("\n" + "="*70)
print("PPT 读取完成")
print(f"文件位置: {pptx_path}")
print("="*70)
