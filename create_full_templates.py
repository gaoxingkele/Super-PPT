#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
创建完整的主题模板（16 种标准 layout）
"""
import sys
sys.path.insert(0, '.')

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pathlib import Path
import json


# 4 个主题配置
TEMPLATES = {
    'yili_power': {
        'name': '亿力科技-电力生产安全',
        'color_scheme': {
            'primary': '#024177',
            'secondary': '#005D7F',
            'accent': '#E8612D',
            'background': '#FFFFFF',
            'text': '#333333',
            'text_light': '#666666'
        },
        'fonts': {'title': '微软雅黑', 'body': '微软雅黑'}
    },
    'xmu_graph': {
        'name': '厦门大学-图理论研究',
        'color_scheme': {
            'primary': '#178F95',
            'secondary': '#4472C4',
            'accent': '#E8612D',
            'background': '#FFFFFF',
            'text': '#333333',
            'text_light': '#666666'
        },
        'fonts': {'title': '微软雅黑', 'body': '微软雅黑'}
    },
    'epri_nature': {
        'name': '电科院-自然基金',
        'color_scheme': {
            'primary': '#0070C0',
            'secondary': '#4C9857',
            'accent': '#E8612D',
            'background': '#FFFFFF',
            'text': '#333333',
            'text_light': '#666666'
        },
        'fonts': {'title': '微软雅黑', 'body': '微软雅黑'}
    },
    'zhinang_qa': {
        'name': '输配作业智囊',
        'color_scheme': {
            'primary': '#00479D',
            'secondary': '#5B9BD5',
            'accent': '#F2F2F2',
            'background': '#FFFFFF',
            'text': '#333333',
            'text_light': '#666666'
        },
        'fonts': {'title': '微软雅黑', 'body': '微软雅黑'}
    }
}


def hex_to_rgb(hex_color):
    hex_color = hex_color.lstrip('#')
    return RGBColor(int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))


class TemplateGenerator:
    def __init__(self, theme_key, theme_config):
        self.theme = theme_config
        self.colors = theme_config['color_scheme']
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)
        self.slide_count = 0
        
    def add_title(self, slide, text, top=Inches(0.3), left=Inches(0.5), 
                  width=Inches(12), height=Inches(1), size=Pt(28), bold=True):
        """添加标题文本框"""
        box = slide.shapes.add_textbox(left, top, width, height)
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = size
        p.font.bold = bold
        p.font.color.rgb = hex_to_rgb(self.colors['primary'])
        return box
    
    def add_accent_line(self, slide, top=Inches(1.25), left=Inches(0.5)):
        """添加装饰线"""
        line = slide.shapes.add_shape(
            1, left, top, Inches(2), Pt(4)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = hex_to_rgb(self.colors['accent'])
        line.line.fill.background()
        
    def add_bottom_bar(self, slide):
        """添加底部色条"""
        bar = slide.shapes.add_shape(
            1, Inches(0), Inches(7.35), self.prs.slide_width, Pt(8)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = hex_to_rgb(self.colors['primary'])
        bar.line.fill.background()

    # ========== 16 种 Layout 生成器 ==========
    
    def create_cover(self):
        """1. 封面"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # 标题
        self.add_title(slide, "标题", top=Inches(2.5), size=Pt(44))
        
        # 副标题
        box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12), Inches(1))
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = "副标题"
        p.font.size = Pt(24)
        p.font.color.rgb = hex_to_rgb(self.colors['text'])
        
        self.add_accent_line(slide, top=Inches(4.0))
        self.slide_count += 1
        return slide
    
    def create_agenda(self):
        """2. 目录页"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.add_title(slide, "目录", size=Pt(32))
        self.add_accent_line(slide)
        self.add_bottom_bar(slide)
        
        # 目录项占位区
        for i in range(4):
            y = Inches(2.0 + i * 1.0)
            # 编号圆圈
            circle = slide.shapes.add_shape(9, Inches(1.6), y, Inches(0.45), Inches(0.45))
            circle.fill.solid()
            circle.fill.fore_color.rgb = hex_to_rgb(self.colors['primary'])
            
            # 文字
            box = slide.shapes.add_textbox(Inches(2.3), y, Inches(9), Inches(0.5))
            tf = box.text_frame
            p = tf.paragraphs[0]
            p.text = f"章节 {i+1}"
            p.font.size = Pt(20)
            
        self.slide_count += 1
        return slide
    
    def create_section_break(self):
        """3. 章节过渡页"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # 背景色块
        bg = slide.shapes.add_shape(
            1, Inches(0), Inches(0), self.prs.slide_width, self.prs.slide_height
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = hex_to_rgb(self.colors['primary'])
        bg.line.fill.background()
        
        # 编号徽章
        badge = slide.shapes.add_shape(9, Inches(1.5), Inches(1.8), Inches(1.6), Inches(1.6))
        badge.fill.solid()
        badge.fill.fore_color.rgb = hex_to_rgb(self.colors['accent'])
        
        box = badge.text_frame
        p = box.paragraphs[0]
        p.text = "01"
        p.font.size = Pt(56)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
        
        # 标题
        box = slide.shapes.add_textbox(Inches(1.5), Inches(3.8), Inches(10), Inches(1.2))
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = "章节标题"
        p.font.size = Pt(40)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        
        self.slide_count += 1
        return slide
    
    def create_title_content(self):
        """4. 标题+内容页"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.add_title(slide, "标题")
        self.add_accent_line(slide)
        self.add_bottom_bar(slide)
        
        # 左侧内容区
        box = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(6.5), Inches(5))
        tf = box.text_frame
        tf.word_wrap = True
        for i in range(3):
            p = tf.add_paragraph()
            p.text = f"要点 {i+1}"
            p.font.size = Pt(16)
            
        # 右侧图片占位
        placeholder = slide.shapes.add_shape(
            1, Inches(7.2), Inches(1.6), Inches(5.3), Inches(5)
        )
        placeholder.fill.solid()
        placeholder.fill.fore_color.rgb = hex_to_rgb(self.colors['accent'])
        
        self.slide_count += 1
        return slide
    
    def create_data_chart(self):
        """5. 数据图表页"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.add_title(slide, "数据标题")
        self.add_accent_line(slide)
        self.add_bottom_bar(slide)
        
        # 图表区域占位
        chart_area = slide.shapes.add_shape(
            1, Inches(1), Inches(1.5), Inches(11), Inches(5)
        )
        chart_area.fill.solid()
        chart_area.fill.fore_color.rgb = RGBColor(245, 245, 245)
        
        self.slide_count += 1
        return slide
    
    def create_infographic(self):
        """6. 信息图页"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.add_title(slide, "信息图标题")
        self.add_accent_line(slide)
        self.add_bottom_bar(slide)
        
        # 信息图占位
        placeholder = slide.shapes.add_shape(
            1, Inches(1), Inches(1.6), Inches(11), Inches(5.2)
        )
        placeholder.fill.solid()
        placeholder.fill.fore_color.rgb = hex_to_rgb(self.colors['secondary'])
        
        self.slide_count += 1
        return slide
    
    def create_two_column(self):
        """7. 双栏对比页"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.add_title(slide, "对比标题")
        self.add_accent_line(slide)
        self.add_bottom_bar(slide)
        
        # 中间分隔线
        divider = slide.shapes.add_shape(
            1, Inches(6.5), Inches(1.8), Pt(2), Inches(4.8)
        )
        divider.fill.solid()
        divider.fill.fore_color.rgb = RGBColor(221, 221, 221)
        
        # 左栏标题
        box = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(5.8), Inches(0.6))
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = "方案 A"
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = hex_to_rgb(self.colors['primary'])
        
        # 右栏标题
        box = slide.shapes.add_textbox(Inches(7), Inches(1.8), Inches(5.8), Inches(0.6))
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = "方案 B"
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = hex_to_rgb(self.colors['accent'])
        
        self.slide_count += 1
        return slide
    
    def create_key_insight(self):
        """8. 核心发现页"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.add_title(slide, "核心发现")
        self.add_accent_line(slide)
        self.add_bottom_bar(slide)
        
        # 顶部强调条
        bar = slide.shapes.add_shape(
            1, Inches(0), Inches(0), self.prs.slide_width, Inches(0.12)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = hex_to_rgb(self.colors['accent'])
        
        # 左侧要点区
        box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(6.5), Inches(4.5))
        tf = box.text_frame
        tf.word_wrap = True
        for i in range(3):
            p = tf.add_paragraph()
            p.text = f"关键要点 {i+1}"
            p.font.size = Pt(16)
        
        # 右侧 KPI 区占位
        placeholder = slide.shapes.add_shape(
            1, Inches(7.8), Inches(1.8), Inches(5), Inches(4.5)
        )
        placeholder.fill.solid()
        placeholder.fill.fore_color.rgb = hex_to_rgb(self.colors['secondary'])
        
        self.slide_count += 1
        return slide
    
    def create_table(self):
        """9. 表格页"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.add_title(slide, "数据表格")
        self.add_accent_line(slide)
        self.add_bottom_bar(slide)
        
        # 表格占位区
        table_area = slide.shapes.add_shape(
            1, Inches(0.8), Inches(1.8), Inches(11.5), Inches(4.8)
        )
        table_area.fill.solid()
        table_area.fill.fore_color.rgb = RGBColor(250, 250, 250)
        table_area.line.color.rgb = hex_to_rgb(self.colors['primary'])
        
        self.slide_count += 1
        return slide
    
    def create_image_full(self):
        """10. 全屏图片页"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # 全屏图片占位
        placeholder = slide.shapes.add_shape(
            1, Inches(0), Inches(0), self.prs.slide_width, self.prs.slide_height
        )
        placeholder.fill.solid()
        placeholder.fill.fore_color.rgb = hex_to_rgb(self.colors['secondary'])
        
        self.slide_count += 1
        return slide
    
    def create_quote(self):
        """11. 引用页"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # 浅色背景
        bg = slide.shapes.add_shape(
            1, Inches(0), Inches(0), self.prs.slide_width, self.prs.slide_height
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = RGBColor(248, 250, 252)
        bg.line.fill.background()
        
        # 引号装饰
        box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(2), Inches(1.5))
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = "\u201C"
        p.font.size = Pt(120)
        p.font.color.rgb = hex_to_rgb(self.colors['accent'])
        
        # 引用文字
        box = slide.shapes.add_textbox(Inches(2), Inches(2.5), Inches(9), Inches(2.5))
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = "引用内容"
        p.font.size = Pt(32)
        p.font.italic = True
        p.font.color.rgb = hex_to_rgb(self.colors['primary'])
        p.alignment = PP_ALIGN.CENTER
        
        self.slide_count += 1
        return slide
    
    def create_timeline(self):
        """12. 时间线页"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.add_title(slide, "发展历程")
        self.add_accent_line(slide)
        self.add_bottom_bar(slide)
        
        # 时间线占位
        timeline = slide.shapes.add_shape(
            1, Inches(0.5), Inches(2.5), Inches(12.3), Inches(3)
        )
        timeline.fill.solid()
        timeline.fill.fore_color.rgb = hex_to_rgb(self.colors['secondary'])
        
        self.slide_count += 1
        return slide
    
    def create_summary(self):
        """13. 总结页"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.add_title(slide, "总结与展望")
        self.add_accent_line(slide)
        self.add_bottom_bar(slide)
        
        # 内容卡片
        card = slide.shapes.add_shape(
            1, Inches(0.5), Inches(1.6), Inches(12), Inches(3.2)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(245, 247, 250)
        card.line.fill.background()
        
        # CTA 按钮
        cta = slide.shapes.add_shape(
            1, Inches(2), Inches(5.0), Inches(9), Inches(1)
        )
        cta.fill.solid()
        cta.fill.fore_color.rgb = hex_to_rgb(self.colors['accent'])
        cta.line.fill.background()
        
        self.slide_count += 1
        return slide
    
    def create_methodology(self):
        """14. 研究方法/技术路线页"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.add_title(slide, "技术路线")
        self.add_accent_line(slide)
        self.add_bottom_bar(slide)
        
        # 步骤展示
        for i in range(4):
            x = Inches(1.5 + i * 2.8)
            # 编号徽章
            badge = slide.shapes.add_shape(9, x, Inches(1.8), Inches(0.5), Inches(0.5))
            badge.fill.solid()
            badge.fill.fore_color.rgb = hex_to_rgb(self.colors['secondary'])
            
            tx = badge.text_frame
            p = tx.paragraphs[0]
            p.text = str(i+1)
            p.font.size = Pt(18)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.alignment = PP_ALIGN.CENTER
        
        # 流程图占位
        flow = slide.shapes.add_shape(
            1, Inches(1), Inches(3.0), Inches(11), Inches(3)
        )
        flow.fill.solid()
        flow.fill.fore_color.rgb = hex_to_rgb(self.colors['secondary'])
        
        self.slide_count += 1
        return slide
    
    def create_architecture(self):
        """15. 系统架构页"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.add_title(slide, "系统架构")
        self.add_accent_line(slide)
        self.add_bottom_bar(slide)
        
        # 架构图占位
        arch = slide.shapes.add_shape(
            1, Inches(0.8), Inches(1.6), Inches(11.5), Inches(4)
        )
        arch.fill.solid()
        arch.fill.fore_color.rgb = hex_to_rgb(self.colors['secondary'])
        
        # 底部要点卡片
        for i in range(3):
            x = Inches(1 + i * 4)
            card = slide.shapes.add_shape(
                1, x, Inches(5.8), Inches(3.5), Inches(1.2)
            )
            card.fill.solid()
            card.fill.fore_color.rgb = RGBColor(240, 244, 250)
        
        self.slide_count += 1
        return slide
    
    def create_end(self):
        """16. 结束页"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # 背景
        bg = slide.shapes.add_shape(
            1, Inches(0), Inches(0), self.prs.slide_width, self.prs.slide_height
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = hex_to_rgb(self.colors['primary'])
        bg.line.fill.background()
        
        # 装饰线
        line = slide.shapes.add_shape(
            1, Inches(5.5), Inches(2.5), Inches(2.3), Pt(5)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = hex_to_rgb(self.colors['accent'])
        line.line.fill.background()
        
        # 感谢文字
        box = slide.shapes.add_textbox(Inches(1), Inches(2.8), Inches(11), Inches(2))
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = "感谢聆听"
        p.font.size = Pt(56)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
        
        self.slide_count += 1
        return slide
    
    def generate(self, output_path):
        """生成完整模板"""
        # 按顺序创建所有 16 种 layout
        layouts = [
            ('cover', self.create_cover),
            ('agenda', self.create_agenda),
            ('section_break', self.create_section_break),
            ('title_content', self.create_title_content),
            ('data_chart', self.create_data_chart),
            ('infographic', self.create_infographic),
            ('two_column', self.create_two_column),
            ('key_insight', self.create_key_insight),
            ('table', self.create_table),
            ('image_full', self.create_image_full),
            ('quote', self.create_quote),
            ('timeline', self.create_timeline),
            ('summary', self.create_summary),
            ('methodology', self.create_methodology),
            ('architecture', self.create_architecture),
            ('end', self.create_end),
        ]
        
        for name, creator in layouts:
            creator()
            print(f"  [OK] Created layout: {name}")
        
        # 保存
        self.prs.save(output_path)
        print(f"  [OK] Saved: {output_path}")
        
        return self.slide_count


def main():
    print("="*60)
    print("Creating Full Templates (16 layouts)")
    print("="*60)
    
    Path('themes').mkdir(exist_ok=True)
    
    for key, config in TEMPLATES.items():
        print(f"\nGenerating: {config['name']}")
        gen = TemplateGenerator(key, config)
        output_path = f"themes/{key}.pptx"
        count = gen.generate(output_path)
        
        # 更新 profile
        profile = {
            "theme_name": key,
            "display_name": config['name'],
            "color_scheme": config['color_scheme'],
            "typography": {
                "title_font": config['fonts']['title'],
                "body_font": config['fonts']['body'],
                "title_size_pt": 32,
                "body_size_pt": 18
            },
            "available_layouts": 16,
            "layouts": [
                "cover", "agenda", "section_break", "title_content",
                "data_chart", "infographic", "two_column", "key_insight",
                "table", "image_full", "quote", "timeline",
                "summary", "methodology", "architecture", "end"
            ]
        }
        
        profile_path = f"themes/{key}_profile.json"
        with open(profile_path, 'w', encoding='utf-8') as f:
            json.dump(profile, f, ensure_ascii=False, indent=2)
        print(f"  [OK] Profile: {profile_path}")
    
    print("\n" + "="*60)
    print("All templates generated successfully!")
    print("="*60)


if __name__ == '__main__':
    main()
