#!/usr/bin/env python3
"""修复批次2: 从raw_response重新解析slide_plan, 重跑Step3+Step4, 然后合并所有批次"""
import sys
sys.path.insert(0, '.')

import json
import re
import shutil
from pathlib import Path

OUTPUT_DIR = Path("output/japan_defense_50pages")
BATCH2_DIR = OUTPUT_DIR / "batch_2"

# ============================================================
# Step 1: 修复 slide_plan.json
# ============================================================
print("="*70)
print("[修复] 从 raw_response 重新解析 slide_plan")
print("="*70)

plan_path = BATCH2_DIR / "slide_plan.json"
plan = json.loads(plan_path.read_text(encoding='utf-8'))

if plan.get('parse_error') and plan.get('raw_response'):
    raw = plan['raw_response']
    # 去除 ```json ... ``` 包裹
    raw = re.sub(r'^```(?:json)?\s*\n?', '', raw)
    raw = re.sub(r'\n?```\s*$', '', raw)

    try:
        parsed = json.loads(raw)
        slides = parsed.get('slides', [])
        meta = parsed.get('meta', {})
        print(f"  成功解析: {len(slides)} 张幻灯片")

        # 覆盖 slide_plan.json
        fixed_plan = {
            "meta": meta,
            "slides": slides
        }
        plan_path.write_text(
            json.dumps(fixed_plan, ensure_ascii=False, indent=2),
            encoding='utf-8'
        )
        print(f"  已保存修复后的 slide_plan.json")
    except json.JSONDecodeError as e:
        print(f"  解析失败: {e}")
        print("  尝试截取到最后一个完整的 slide 对象...")
        # 尝试找到最后一个完整的 } 并截断
        # 找 "slides": [ ... ] 的结构
        # 先试试能不能找到slides数组
        match = re.search(r'"slides"\s*:\s*\[', raw)
        if match:
            # 尝试逐步截断找到合法JSON
            for i in range(len(raw)-1, match.end(), -1):
                if raw[i] == '}':
                    try:
                        # 尝试闭合
                        attempt = raw[:i+1]
                        # 计算未闭合的括号
                        open_sq = attempt.count('[') - attempt.count(']')
                        open_cr = attempt.count('{') - attempt.count('}')
                        attempt += ']' * open_sq + '}' * open_cr
                        parsed = json.loads(attempt)
                        slides = parsed.get('slides', [])
                        meta = parsed.get('meta', {})
                        print(f"  截断修复成功: {len(slides)} 张幻灯片")
                        fixed_plan = {"meta": meta, "slides": slides}
                        plan_path.write_text(
                            json.dumps(fixed_plan, ensure_ascii=False, indent=2),
                            encoding='utf-8'
                        )
                        print(f"  已保存修复后的 slide_plan.json")
                        break
                    except:
                        continue
        else:
            print("  无法修复，需要重新运行 Step 2")
            sys.exit(1)
else:
    print("  slide_plan 无需修复")

# ============================================================
# Step 2: 重跑 Step 3 (视觉资产) + Step 4 (生成PPT)
# ============================================================
from src.step3_visuals import run_visuals
from src.step4_build import run_build

batch_base = "japan_defense_50pages_batch2"

print(f"\n[批次 2] Step 3: 视觉资产...")
result3 = run_visuals(batch_base, BATCH2_DIR, no_ai_images=True)
print(f"  完成: {len(result3.get('manifest', {}))} 个")

print(f"\n[批次 2] Step 4: 生成PPT...")
result4 = run_build(batch_base, BATCH2_DIR, theme="tech")

# 复制到汇总目录
src_pptx = result4['pptx_path']
dst_pptx = OUTPUT_DIR / "batch_2_电子与潜艇.pptx"
shutil.copy(src_pptx, dst_pptx)

# 读取页数
plan_fixed = json.loads(plan_path.read_text(encoding='utf-8'))
pages = len(plan_fixed.get('slides', []))
print(f"  批次2修复完成: {pages} 页 -> {dst_pptx.name}")

# ============================================================
# Step 3: 合并所有批次
# ============================================================
print("\n" + "="*70)
print("合并所有批次")
print("="*70)

from pptx import Presentation
from copy import deepcopy
import copy

batch_files = sorted(OUTPUT_DIR.glob("batch_*_*.pptx"))
print(f"找到 {len(batch_files)} 个批次文件:")
for f in batch_files:
    prs = Presentation(str(f))
    print(f"  {f.name}: {len(prs.slides)} 页")

# 合并
merged = Presentation(str(batch_files[0]))
print(f"\n基础: {batch_files[0].name} ({len(merged.slides)} 页)")

for bf in batch_files[1:]:
    src_prs = Presentation(str(bf))
    if len(src_prs.slides) == 0:
        print(f"  跳过空文件: {bf.name}")
        continue

    for slide in src_prs.slides:
        # 复制 slide layout
        # 使用空白布局
        blank_layout = merged.slide_layouts[6]  # blank
        new_slide = merged.slides.add_slide(blank_layout)

        # 复制slide尺寸和背景
        # 复制所有shapes
        for shape in slide.shapes:
            el = copy.deepcopy(shape._element)
            new_slide.shapes._spTree.append(el)

        # 复制背景
        if slide.background.fill.type is not None:
            bg_el = copy.deepcopy(slide.background._element)
            new_slide.background._element.getparent().replace(
                new_slide.background._element, bg_el
            )

    print(f"  合并: {bf.name} ({len(src_prs.slides)} 页)")

merged_path = OUTPUT_DIR / "日本军工企业深度研究_完整版.pptx"
merged.save(str(merged_path))
total_pages = len(merged.slides)
print(f"\n合并完成: {merged_path.name} ({total_pages} 页)")
print(f"路径: {merged_path}")
