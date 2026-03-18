#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
分析 source-doc 目录下的 PPTX 文件，提取设计元素
修复编码和错误处理问题
"""
import sys
sys.path.insert(0, '.')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pathlib import Path
import json
import os


def safe_get_text(shape):
    """安全获取文本"""
    try:
        if hasattr(shape, 'text'):
            return shape.text.strip()
    except:
        pass
    return None


def analyze_shape(shape, indent=0):
    """递归分析形状"""
    info = {
        'shape_type': str(shape.shape_type),
        'has_text_frame': hasattr(shape, 'text_frame'),
        'has_table': shape.shape_type == MSO_SHAPE_TYPE.TABLE,
        'has_chart': shape.shape_type == MSO_SHAPE_TYPE.CHART,
        'has_picture': shape.shape_type == MSO_SHAPE_TYPE.PICTURE,
    }
    
    # 尝试获取文本
    text = safe_get_text(shape)
    if text:
        info['text_preview'] = text[:50] if len(text) > 50 else text
    
    # 尝试获取位置
    try:
        info['position'] = {
            'left': str(shape.left),
            'top': str(shape.top),
            'width': str(shape.width),
            'height': str(shape.height)
        }
    except:
        pass
    
    # 尝试获取填充颜色
    try:
        if hasattr(shape, 'fill'):
            fill = shape.fill
            if fill.type is not None:
                info['fill_type'] = str(fill.type)
                try:
                    if hasattr(fill, 'fore_color'):
                        rgb = fill.fore_color.rgb
                        if rgb:
                            info['fill_color'] = f'#{rgb}'
                except:
                    pass
    except:
        pass
    
    return info


def analyze_slide(slide, slide_idx):
    """分析单页幻灯片"""
    shapes_info = []
    has_chart = False
    has_table = False
    has_image = False
    title_candidates = []
    
    for shape in slide.shapes:
        info = analyze_shape(shape)
        shapes_info.append(info)
        
        if info['has_chart']:
            has_chart = True
        if info['has_table']:
            has_table = True
        if info['has_picture']:
            has_image = True
        if info.get('text_preview'):
            # 可能是标题（文字较短且位置靠上）
            text = info['text_preview']
            if len(text) < 60 and len(text) > 5:
                try:
                    top = shape.top
                    title_candidates.append((top, text))
                except:
                    title_candidates.append((0, text))
    
    # 分类布局类型
    layout_type = 'unknown'
    if slide_idx == 0:
        layout_type = 'cover'
    elif has_chart:
        layout_type = 'data_chart'
    elif has_table:
        layout_type = 'table'
    elif has_image and len(shapes_info) <= 3:
        layout_type = 'image_full'
    elif len([s for s in shapes_info if s.get('text_preview')]) >= 2:
        layout_type = 'title_content'
    
    # 找标题
    title = None
    if title_candidates:
        title_candidates.sort(key=lambda x: x[0])  # 按位置排序
        title = title_candidates[0][1][:60]
    
    return {
        'slide_idx': slide_idx,
        'shape_count': len(shapes_info),
        'has_chart': has_chart,
        'has_table': has_table,
        'has_image': has_image,
        'layout_type': layout_type,
        'title': title,
        'shapes': shapes_info[:5]  # 只存前5个形状的详情
    }


def extract_design_features(prs):
    """提取整体设计特征"""
    features = {
        'slide_count': len(prs.slides),
        'slide_width': str(prs.slide_width),
        'slide_height': str(prs.slide_height),
        'masters_count': len(prs.slide_masters),
        'color_samples': [],
        'font_samples': []
    }
    
    # 收集颜色和字体样本
    slide_list = list(prs.slides)
    for slide in slide_list[:5]:  # 只看前5页
        for shape in slide.shapes:
            # 颜色
            try:
                if hasattr(shape, 'fill') and shape.fill.type is not None:
                    try:
                        rgb = shape.fill.fore_color.rgb
                        if rgb:
                            features['color_samples'].append(f'#{rgb}')
                    except:
                        pass
            except:
                pass
            
            # 字体
            try:
                if hasattr(shape, 'text_frame'):
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if run.font.name:
                                features['font_samples'].append(run.font.name)
            except:
                pass
    
    # 统计最常见的
    from collections import Counter
    if features['color_samples']:
        features['top_colors'] = Counter(features['color_samples']).most_common(5)
    if features['font_samples']:
        features['top_fonts'] = Counter(features['font_samples']).most_common(3)
    
    return features


def analyze_pptx(filepath):
    """分析单个 PPTX 文件"""
    filename = os.path.basename(filepath)
    
    print(f"\n{'='*70}")
    print(f"文件: {filename}")
    print('='*70)
    
    try:
        prs = Presentation(filepath)
    except Exception as e:
        print(f"  无法打开: {e}")
        return None
    
    # 提取整体特征
    features = extract_design_features(prs)
    print(f"  幻灯片总数: {features['slide_count']}")
    print(f"  尺寸: {features['slide_width']} x {features['slide_height']}")
    
    if features.get('top_colors'):
        print(f"\n  主要颜色:")
        for color, count in features['top_colors']:
            print(f"    {color} ({count}次)")
    
    if features.get('top_fonts'):
        print(f"\n  主要字体:")
        for font, count in features['top_fonts']:
            print(f"    {font} ({count}次)")
    
    # 分析每页
    print(f"\n  幻灯片布局分析:")
    layouts = []
    for idx, slide in enumerate(prs.slides):
        layout = analyze_slide(slide, idx)
        layouts.append(layout)
        
        title_str = layout.get('title') or '无标题'
        print(f"    [{idx+1:2d}] {layout['layout_type']:15s} - {title_str[:40]}")
    
    return {
        'filename': filename,
        'features': features,
        'layouts': layouts
    }


def main():
    source_dir = Path('source-doc')
    
    # 目标文件关键词
    target_keywords = ['多模态', '图理论', '电网', '智囊']
    
    # 找到匹配的文件
    pptx_files = []
    for f in source_dir.glob('*.pptx'):
        fstr = f.name
        for kw in target_keywords:
            if kw in fstr:
                pptx_files.append(f)
                break
    
    print(f"找到 {len(pptx_files)} 个目标文件")
    
    # 分析
    results = []
    for filepath in pptx_files:
        result = analyze_pptx(filepath)
        if result:
            results.append(result)
    
    # 保存
    output_file = 'pptx_analysis_report.json'
    with open(output_file, 'w', encoding='utf-8') as f:
        json.dump(results, f, ensure_ascii=False, indent=2)
    
    print(f"\n{'='*70}")
    print(f"报告已保存: {output_file}")
    print('='*70)


if __name__ == '__main__':
    main()
