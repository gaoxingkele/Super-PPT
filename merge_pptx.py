#!/usr/bin/env python3
"""
merge_pptx.py — Relationship-aware PPTX batch merger.

Merges multiple batch_*_*.pptx files into a single presentation.
- Batch 1: keep ALL slides
- Batch 2+: skip slides with layout types "cover", "agenda", "end"

Uses proper relationship handling so images and charts survive the copy.

Usage:
    python merge_pptx.py [--input-dir DIR] [--output NAME]

Defaults:
    --input-dir  output/japan_defense_50pages
    --output     日本军工企业深度研究_完整版.pptx
"""

import sys
import json
import copy
import argparse
import tempfile
import io
from pathlib import Path
from typing import Optional

from pptx import Presentation
from pptx.util import Inches, Emu
from pptx.oxml.ns import qn
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.opc.package import Part
from pptx.opc.packuri import PackURI
import lxml.etree as etree


# ── Slide dimension constants (16:9 widescreen) ──────────────────────
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# Layout types to skip from batches 2+
SKIP_LAYOUTS = {"cover", "agenda", "end"}


# =====================================================================
# Slide plan helpers
# =====================================================================

def load_slide_plan(batch_dir: Path) -> Optional[dict]:
    """Load slide_plan.json from a batch directory."""
    plan_file = batch_dir / "slide_plan.json"
    if plan_file.exists():
        return json.loads(plan_file.read_text(encoding="utf-8"))
    return None


def get_skip_indices(batch_dir: Path) -> set:
    """
    Determine which slide indices (0-based) to skip based on slide_plan.json.
    Returns empty set if plan not found (conservative: keep all).
    """
    plan = load_slide_plan(batch_dir)
    if not plan:
        return set()

    skip = set()
    slides = plan.get("slides", [])
    for i, slide_info in enumerate(slides):
        layout = slide_info.get("layout", "").lower().strip()
        if layout in SKIP_LAYOUTS:
            skip.add(i)
    return skip


def detect_skip_by_content(slide, slide_idx: int, total_slides: int) -> bool:
    """
    Fallback heuristic: detect cover/agenda/end slides by content patterns.
    Used when slide_plan.json is not available.
    """
    # Collect all text from the slide
    all_text = ""
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                all_text += para.text + " "

    text = all_text.strip()

    # End slide patterns
    end_patterns = ["谢谢聆听", "感谢聆听", "Thank You", "Q&A", "Questions"]
    for pat in end_patterns:
        if pat in text:
            return True

    # Agenda patterns
    agenda_patterns = ["报告目录", "报告议程", "研究框架", "叙事逻辑"]
    for pat in agenda_patterns:
        if pat in text:
            return True

    # Cover: first slide often has very few shapes and the main title
    if slide_idx == 0:
        return True

    # Last slide is often "end"
    if slide_idx == total_slides - 1:
        # Check if it has very few content shapes
        content_shapes = sum(1 for s in slide.shapes if s.has_text_frame)
        if content_shapes <= 5:
            return True

    return False


# =====================================================================
# Relationship-aware shape copying
# =====================================================================

def _copy_image_shape(src_shape, dst_slide):
    """
    Copy a picture shape by extracting the image blob and re-adding it
    to the destination slide via add_picture().
    """
    # Extract image blob from source
    image = src_shape.image
    image_blob = image.blob
    content_type = image.content_type

    # Determine file extension from content type
    ext_map = {
        "image/png": ".png",
        "image/jpeg": ".jpg",
        "image/gif": ".gif",
        "image/bmp": ".bmp",
        "image/tiff": ".tiff",
        "image/svg+xml": ".svg",
        "image/webp": ".webp",
    }
    ext = ext_map.get(content_type, ".png")

    # Get position and size from source shape
    left = src_shape.left
    top = src_shape.top
    width = src_shape.width
    height = src_shape.height

    # Write blob to a BytesIO stream and add to destination
    img_stream = io.BytesIO(image_blob)
    pic = dst_slide.shapes.add_picture(img_stream, left, top, width, height)

    # Copy crop settings if any
    try:
        src_crop = src_shape._element.find(qn("p:spPr") + "/" + qn("a:blipFill") + "/" + qn("a:srcRect"))
        if src_crop is None:
            # Try in the pic:blipFill path
            src_blipFill = src_shape._element.find(".//" + qn("a:blipFill"))
            if src_blipFill is not None:
                src_crop = src_blipFill.find(qn("a:srcRect"))
        if src_crop is not None:
            dst_blipFill = pic._element.find(".//" + qn("a:blipFill"))
            if dst_blipFill is not None:
                existing = dst_blipFill.find(qn("a:srcRect"))
                if existing is not None:
                    dst_blipFill.remove(existing)
                dst_blipFill.append(copy.deepcopy(src_crop))
    except Exception:
        pass

    return pic


def _copy_chart_shape(src_shape, src_slide, dst_slide, dst_prs):
    """
    Copy a chart shape by duplicating the chart part (XML + embedded Excel)
    and re-linking it in the destination slide.
    """
    # Find the chart relationship in source slide
    chart_rId = None
    graphic_frame_el = src_shape._element
    # The rId is in the graphicFrame's graphic > graphicData > chart element
    chart_ref = graphic_frame_el.find(".//" + qn("c:chart"))
    if chart_ref is not None:
        chart_rId = chart_ref.get(qn("r:id"))

    if chart_rId is None:
        # Fallback: deepcopy the element (chart won't work but at least placeholder)
        el = copy.deepcopy(graphic_frame_el)
        dst_slide.shapes._spTree.append(el)
        return

    # Get the chart part from source
    src_chart_rel = src_slide.part.rels[chart_rId]
    src_chart_part = src_chart_rel.target_part

    # Deep copy the graphicFrame element
    new_el = copy.deepcopy(graphic_frame_el)

    # Add the chart part to the destination slide's package
    # We need to create a new relationship and copy the chart part

    # Create a unique partname for the chart in destination
    existing_chart_names = set()
    for part in dst_prs.part.package.iter_parts():
        if "/charts/" in str(part.partname):
            existing_chart_names.add(str(part.partname))

    chart_num = 1
    while f"/ppt/charts/chart{chart_num}.xml" in existing_chart_names:
        chart_num += 1
    new_chart_partname = f"/ppt/charts/chart{chart_num}.xml"

    # Create the chart part as a copy (note: Part(partname, content_type, package, blob))
    new_chart_part = Part(
        PackURI(new_chart_partname),
        src_chart_part.content_type,
        dst_prs.part.package,
        src_chart_part.blob,
    )

    # Copy chart's own relationships (embedded Excel, etc.)
    for rel in src_chart_part.rels.values():
        target_part = rel.target_part
        # Create unique partname for embedded parts
        if "embeddings" in str(target_part.partname):
            existing_embed_names = set()
            for p in dst_prs.part.package.iter_parts():
                if "/embeddings/" in str(p.partname):
                    existing_embed_names.add(str(p.partname))
            embed_num = 1
            orig_ext = Path(str(target_part.partname)).suffix
            while f"/ppt/embeddings/Microsoft_Excel_Sheet{embed_num}{orig_ext}" in existing_embed_names:
                embed_num += 1
            new_embed_partname = f"/ppt/embeddings/Microsoft_Excel_Sheet{embed_num}{orig_ext}"

            new_embed_part = Part(
                PackURI(new_embed_partname),
                target_part.content_type,
                dst_prs.part.package,
                target_part.blob,
            )
            # Use _add_relationship to preserve the original rId
            new_chart_part.rels._add_relationship(rel.reltype, new_embed_part)
        else:
            # Other rel types - try to copy
            try:
                new_sub_part = Part(
                    target_part.partname,
                    target_part.content_type,
                    dst_prs.part.package,
                    target_part.blob,
                )
                new_chart_part.rels._add_relationship(rel.reltype, new_sub_part)
            except Exception:
                pass

    # Add relationship from destination slide to new chart part
    new_rId = dst_slide.part.relate_to(new_chart_part, src_chart_rel.reltype)

    # Update the rId in the copied element
    chart_ref_new = new_el.find(".//" + qn("c:chart"))
    if chart_ref_new is not None:
        chart_ref_new.set(qn("r:id"), new_rId)

    # Append the element to the destination slide
    dst_slide.shapes._spTree.append(new_el)


def _copy_simple_shape(src_shape, dst_slide):
    """
    Copy a non-image, non-chart shape via deepcopy of its XML element.
    Works for text boxes, auto shapes, rectangles, etc.
    """
    el = copy.deepcopy(src_shape._element)
    dst_slide.shapes._spTree.append(el)


def copy_slide(src_slide, src_prs, dst_prs):
    """
    Copy a single slide from src_prs to dst_prs with proper relationship handling.
    Returns the new slide object.
    """
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    # Add a blank slide
    blank_layout = dst_prs.slide_layouts[6]  # blank layout
    new_slide = dst_prs.slides.add_slide(blank_layout)

    # Copy each shape with appropriate method
    for shape in src_slide.shapes:
        shape_type = shape.shape_type

        if shape_type == MSO_SHAPE_TYPE.PICTURE:
            try:
                _copy_image_shape(shape, new_slide)
            except Exception as e:
                print(f"    WARNING: Failed to copy picture '{shape.name}': {e}")
                _copy_simple_shape(shape, new_slide)

        elif hasattr(shape, "has_chart") and shape.has_chart:
            try:
                _copy_chart_shape(shape, src_slide, new_slide, dst_prs)
            except Exception as e:
                print(f"    WARNING: Failed to copy chart '{shape.name}': {e}")
                _copy_simple_shape(shape, new_slide)

        elif shape_type == MSO_SHAPE_TYPE.GROUP:
            # Group shapes may contain pictures; for now do deepcopy
            # (group images are rarer in this project)
            try:
                _copy_group_shape(shape, src_slide, new_slide, dst_prs)
            except Exception as e:
                print(f"    WARNING: Failed to copy group '{shape.name}': {e}")
                _copy_simple_shape(shape, new_slide)

        else:
            _copy_simple_shape(shape, new_slide)

    # Copy slide-level background if it has image fill
    _copy_slide_background(src_slide, new_slide, dst_prs)

    # Copy speaker notes
    _copy_notes(src_slide, new_slide)

    return new_slide


def _copy_group_shape(src_shape, src_slide, dst_slide, dst_prs):
    """
    Copy a group shape. For groups containing pictures, we need to handle
    image relationships. For simple groups, deepcopy works.
    """
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    group_el = src_shape._element
    # Check if group contains any pictures (blip references)
    blips = group_el.findall(".//" + qn("a:blip"))

    if not blips:
        # No images in group, safe to deepcopy
        _copy_simple_shape(src_shape, dst_slide)
        return

    # Group has images - need to copy the image parts and update rIds
    new_el = copy.deepcopy(group_el)

    for blip in new_el.findall(".//" + qn("a:blip")):
        old_rId = blip.get(qn("r:embed"))
        if old_rId and old_rId in src_slide.part.rels:
            rel = src_slide.part.rels[old_rId]
            if "image" in rel.reltype:
                # Copy the image part
                img_part = rel.target_part
                new_rId = dst_slide.part.relate_to(img_part, rel.reltype)
                blip.set(qn("r:embed"), new_rId)

    dst_slide.shapes._spTree.append(new_el)


def _copy_slide_background(src_slide, dst_slide, dst_prs):
    """
    Copy the slide background, handling image fills if present.
    """
    src_bg = src_slide._element.find(qn("p:bg"))
    if src_bg is None:
        return  # No explicit background

    new_bg = copy.deepcopy(src_bg)

    # Check for image references in background
    for blip in new_bg.findall(".//" + qn("a:blip")):
        old_rId = blip.get(qn("r:embed"))
        if old_rId and old_rId in src_slide.part.rels:
            rel = src_slide.part.rels[old_rId]
            if "image" in rel.reltype:
                img_part = rel.target_part
                new_rId = dst_slide.part.relate_to(img_part, rel.reltype)
                blip.set(qn("r:embed"), new_rId)

    # Replace or add the background in destination
    dst_bg = dst_slide._element.find(qn("p:bg"))
    cSld = dst_slide._element
    if dst_bg is not None:
        cSld.remove(dst_bg)

    # Insert bg before spTree
    spTree = cSld.find(qn("p:spTree"))
    if spTree is not None:
        cSld.insert(list(cSld).index(spTree), new_bg)
    else:
        cSld.insert(0, new_bg)


def _copy_notes(src_slide, dst_slide):
    """Copy speaker notes from source to destination slide."""
    try:
        if not src_slide.has_notes_slide:
            return
        src_notes = src_slide.notes_slide
        src_text = src_notes.notes_text_frame.text
        if not src_text.strip():
            return
        # Access or create notes slide on destination
        dst_notes = dst_slide.notes_slide
        dst_notes.notes_text_frame.text = src_text
    except Exception:
        pass


# =====================================================================
# Main merge logic
# =====================================================================

def merge_batches(input_dir: Path, output_name: str = "日本军工企业深度研究_完整版.pptx"):
    """
    Merge all batch_*_*.pptx files in input_dir into a single presentation.

    - Batch 1: all slides kept
    - Batch 2+: slides with layout "cover", "agenda", "end" are skipped
    """
    batch_files = sorted(input_dir.glob("batch_*_*.pptx"))
    if not batch_files:
        print("ERROR: No batch_*_*.pptx files found in", input_dir)
        sys.exit(1)

    print(f"Found {len(batch_files)} batch files:")
    for f in batch_files:
        prs = Presentation(str(f))
        print(f"  {f.name}: {len(prs.slides)} slides")

    # Create a fresh presentation with correct dimensions
    merged = Presentation()
    merged.slide_width = SLIDE_WIDTH
    merged.slide_height = SLIDE_HEIGHT

    total_kept = 0
    total_skipped = 0

    for batch_idx, bf in enumerate(batch_files):
        batch_num = batch_idx + 1
        src_prs = Presentation(str(bf))
        n_slides = len(src_prs.slides)

        if n_slides == 0:
            print(f"\n  Skipping empty file: {bf.name}")
            continue

        # Determine batch directory for slide plan
        # batch files are named like batch_1_宏观与巨头.pptx
        # batch dirs are named like batch_1
        batch_dir_name = bf.stem.split("_")[0] + "_" + bf.stem.split("_")[1]
        batch_dir = input_dir / batch_dir_name

        # Determine which slides to skip.
        # Always use content-based detection because step4_build may add
        # extra slides (e.g. "要点解读" detail pages) that cause slide_plan
        # indices to be misaligned with actual PPTX slide indices.
        #
        # For batch 1: only skip end slides (keep cover + agenda)
        # For batch 2+: skip cover, agenda, and end slides
        skip_indices = set()
        is_last_batch = (batch_idx == len(batch_files) - 1)
        for i, slide in enumerate(src_prs.slides):
            if batch_num == 1:
                # Only skip end-type slides from batch 1 (they'd appear mid-deck)
                all_text = " ".join(
                    shape.text_frame.text for shape in slide.shapes if shape.has_text_frame
                ).strip()
                end_patterns = ["谢谢聆听", "感谢聆听", "Thank You", "Q&A", "Questions"]
                is_end = any(pat in all_text for pat in end_patterns) and len(all_text) < 60
                if is_end:
                    skip_indices.add(i)
            else:
                if detect_skip_by_content(slide, i, n_slides):
                    # For the last batch, keep the final end slide
                    if is_last_batch and i == n_slides - 1:
                        continue
                    skip_indices.add(i)

        # Copy slides
        kept = 0
        skipped_list = []
        for i, slide in enumerate(src_prs.slides):
            if i in skip_indices:
                # Get some text for logging
                title_text = ""
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        title_text = shape.text_frame.text[:40]
                        break
                skipped_list.append(f"#{i+1}({title_text})")
                total_skipped += 1
                continue

            copy_slide(slide, src_prs, merged)
            kept += 1
            total_kept += 1

        print(f"\n  Batch {batch_num} [{bf.name}]: {kept}/{n_slides} slides kept")
        if skipped_list:
            print(f"    Skipped: {', '.join(skipped_list)}")

    # Save (handle file-in-use gracefully)
    output_path = input_dir / output_name
    try:
        merged.save(str(output_path))
    except PermissionError:
        # File might be open in PowerPoint; save with suffix
        stem = Path(output_name).stem
        ext = Path(output_name).suffix
        alt_path = input_dir / f"{stem}_new{ext}"
        print(f"\n  WARNING: Cannot write to {output_name} (file is locked).")
        print(f"  Saving to: {alt_path.name}")
        merged.save(str(alt_path))
        output_path = alt_path

    print(f"\n{'='*70}")
    print(f"Merge complete!")
    print(f"  Total slides kept:    {total_kept}")
    print(f"  Total slides skipped: {total_skipped}")
    print(f"  Output: {output_path}")
    print(f"{'='*70}")

    # Verify the result
    verify = Presentation(str(output_path))
    print(f"\nVerification: {len(verify.slides)} slides in output file")

    return output_path


# =====================================================================
# CLI
# =====================================================================

def main():
    parser = argparse.ArgumentParser(description="Merge batch PPTX files into one")
    parser.add_argument(
        "--input-dir",
        default="output/japan_defense_50pages",
        help="Directory containing batch_*_*.pptx files",
    )
    parser.add_argument(
        "--output",
        default="日本军工企业深度研究_完整版.pptx",
        help="Output filename (saved in input-dir)",
    )
    args = parser.parse_args()

    input_dir = Path(args.input_dir)
    if not input_dir.exists():
        print(f"ERROR: Input directory not found: {input_dir}")
        sys.exit(1)

    merge_batches(input_dir, args.output)


if __name__ == "__main__":
    main()
