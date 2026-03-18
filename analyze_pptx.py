#!/usr/bin/env python3
"""
分析 source-doc 目录下的 PPTX 文件，提取设计元素用于模板生成
"""
import sys
sys.path.insert(0, '.')

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor
from pathlib import Path
import json
import os


def extract_color_scheme(prs):
    """提取配色方案"""
    colors = {}
    try:
        # 尝试从 theme 提取
        slide_master = prs.slide_masters[0]
        theme = slide_master.theme
        
        # 从 slide layouts 的背景和形状中提取实际使用的颜色
        color_samples = []
        for slide in prs.slides[:5]:  # 看前5页
            for shape in slide.shapes:
                if hasattr(shape, 'fill'):
                    try:
                        if shape.fill.type == 1:  # SOLID
                            rgb = shape.fill.fore_color.rgb
                            if rgb:
                                color_samples.append(f'#{rgb}')
                    except:
                        pass
        
        # 统计最常见的颜色
        from collections import Counter
        if color_samples:
            top_colors = Counter(color_samples).most_common(6)
            color_names = ['primary', 'secondary', 'accent', 'background', 'text', 'text_light']
            for i, (color, count) in enumerate(top_colors):
                if i < len(color_names):
                    colors[color_names[i]] = color
                    
    except Exception as e:
        print(f"  提取颜色失败: {e}")
    
    return colors


def extract_fonts(prs):
    """提取字体信息"""
    fonts = {'title_font': None, 'body_font': None}
    
    try:
        # 从文本框中提取字体
        font_samples = {'title': [], 'body': []}
        
        for slide in prs.slides[:3]:
            for shape in slide.shapes:
                if hasattr(shape, 'text_frame'):
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if run.font.name:
                                # 根据字号判断是标题还是正文
                                if run.font.size and run.font.size > Pt(24):
                                    font_samples['title'].append(run.font.name)
                                else:
                                    font_samples['body'].append(run.font.name)
        
        from collections import Counter
        if font_samples['title']:
            fonts['title_font'] = Counter(font_samples['title']).most_common(1)[0][0]
        if font_samples['body']:
            fonts['body_font'] = Counter(font_samples['body']).most_common(1)[0][0]
            
        # 如果没有区分出来，默认一样
        if not fonts['body_font'] and fonts['title_font']:
            fonts['body_font'] = fonts['title_font']
            
    except Exception as e:
        print(f"  提取字体失败: {e}")
    
    return fonts


def analyze_layout(slide, slide_idx):
    """分析单页的布局类型"""
    layout_info = {
        'slide_idx': slide_idx,
        'shape_count': len(slide.shapes),
        'has_chart': False,
        'has_table': False,
        'has_image': False,
        'text_boxes': 0,
        'title': None,
        'layout_type': 'unknown'
    }
    
    for shape in slide.shapes:
        # 检测图表
        if shape.shape_type == MSO_SHAPE_TYPE.CHART:
            layout_info['has_chart'] = True
        # 检测表格
        elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            layout_info['has_table'] = True
        # 检测图片
        elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            layout_info['has_image'] = True
        # 检测文本
        elif hasattr(shape, 'text'):
            if shape.text.strip():
                layout_info['text_boxes'] += 1
                # 尝试找标题（通常是最上面的文本框）
                if not layout_info['title'] and len(shape.text) < 100:
                    layout_info['title'] = shape.text.strip()[:50]
    
    # 简单分类逻辑
    if slide_idx == 0:
        layout_info['layout_type'] = 'cover'
    elif layout_info['has_chart']:
        layout_info['layout_type'] = 'data_chart'
    elif layout_info['has_table']:
        layout_info['layout_type'] = 'table'
    elif layout_info['has_image'] and layout_info['text_boxes'] <= 2:
        layout_info['layout_type'] = 'image_full'
    elif layout_info['text_boxes'] >= 3:
        layout_info['layout_type'] = 'title_content'
    
    return layout_info


def analyze_pptx(filepath):
    """分析单个 PPTX 文件"""
    print(f"\n{'='*60}")
    print(f"分析文件: {os.path.basename(filepath)}")
    print('='*60)
    
    try:
        prs = Presentation(filepath)
    except Exception as e:
        print(f"  无法打开文件: {e}")
        return None
    
    result = {
        'filename': os.path.basename(filepath),
        'slide_count': len(prs.slides),
        'slide_width': str(prs.slide_width),
        'slide_height': str(prs.slide_height),
        'color_scheme': extract_color_scheme(prs),
        'fonts': extract_fonts(prs),
        'layouts': []
    }
    
    # 分析每页布局
    for idx, slide in enumerate(prs.slides):
        layout = analyze_layout(slide, idx)
        result['layouts'].append(layout)
        print(f"  幻灯片 {idx+1}: {layout['layout_type']} - {layout.get('title', '无标题')[:30]}...")
    
    # 打印配色和字体摘要
    print(f"\n  配色方案:")
    for k, v in result['color_scheme'].items():
        print(f"    {k}: {v}")
    
    print(f"\n  字体:")
    for k, v in result['fonts'].items():
        print(f"    {k}: {v}")
    
    return result


def main():
    source_dir = Path('source-doc')
    
    # 用户指定的4个文件关键词
    target_keywords = [
        '多模态',      # 国网信通亿力科技
        '图理论',      # 厦门大学
        '电网',        # 中国电力科学研究院
        '智囊',        # 输配作业智囊
    ]
    
    # 找到匹配的pptx文件
    pptx_files = []
    for f in source_dir.glob('*.pptx'):
        for keyword in target_keywords:
            if keyword in f.name:
                pptx_files.append(f)
                break
    
    print(f"找到 {len(pptx_files)} 个目标文件")
    
    # 分析每个文件
    all_results = []
    for filepath in pptx_files:
        result = analyze_pptx(filepath)
        if result:
            all_results.append(result)
    
    # 保存分析结果
    output_file = 'pptx_analysis_report.json'
    with open(output_file, 'w', encoding='utf-8') as f:
        json.dump(all_results, f, ensure_ascii=False, indent=2)
    
    print(f"\n{'='*60}")
    print(f"分析报告已保存到: {output_file}")
    print('='*60)


if __name__ == '__main__':
    main()
