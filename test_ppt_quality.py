#!/usr/bin/env python3
"""PPT 质量检测脚本 - 全面检查生成的PPTX文件"""
import sys
sys.path.insert(0, '.')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pathlib import Path
import json

# ============================================================
# 配置
# ============================================================
PPTX_PATH = Path("output/japan_defense_50pages_v2/japan_defense_50pages_v2_slides.pptx")

if not PPTX_PATH.exists():
    # fallback
    candidates = list(Path("output/japan_defense_50pages").glob("*完整版*.pptx"))
    if candidates:
        PPTX_PATH = max(candidates, key=lambda p: p.stat().st_mtime)
    else:
        print("找不到合并后的PPTX文件")
        sys.exit(1)

print(f"检测文件: {PPTX_PATH}")
print(f"文件大小: {PPTX_PATH.stat().st_size / 1024 / 1024:.1f} MB")
print("=" * 80)

prs = Presentation(str(PPTX_PATH))
slides = list(prs.slides)
total = len(slides)
print(f"总页数: {total}")

# ============================================================
# 检测项目
# ============================================================
issues = []       # (severity, slide_num, description)
warnings = []     # 警告
stats = {
    'total': total,
    'with_images': 0,
    'with_charts': 0,
    'with_tables': 0,
    'with_text': 0,
    'blank_or_near_blank': 0,
    'duplicate_titles': [],
    'image_count': 0,
    'chart_count': 0,
}

slide_info = []

for i, slide in enumerate(slides):
    snum = i + 1
    shapes = list(slide.shapes)

    # 收集信息
    info = {
        'num': snum,
        'title': '',
        'shape_count': len(shapes),
        'has_image': False,
        'has_chart': False,
        'has_table': False,
        'text_length': 0,
        'image_count': 0,
        'shape_types': [],
        'issues': [],
    }

    title_text = ''
    all_text = ''
    image_count = 0
    chart_count = 0
    table_count = 0
    textbox_count = 0
    rect_count = 0

    for shape in shapes:
        info['shape_types'].append(shape.shape_type)

        # 文本
        if shape.has_text_frame:
            txt = shape.text_frame.text.strip()
            all_text += txt + ' '
            if not title_text and len(txt) > 2:
                title_text = txt[:60]
            textbox_count += 1

        # 图片
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            image_count += 1
            # 检查图片是否有效
            try:
                blob = shape.image.blob
                if len(blob) < 100:
                    info['issues'].append("图片数据异常小")
            except Exception as e:
                info['issues'].append(f"图片读取失败: {e}")

        # 图表
        if shape.has_chart:
            chart_count += 1

        # 表格
        if shape.has_table:
            table_count += 1
            tbl = shape.table
            if tbl.rows and len(tbl.rows) < 2:
                info['issues'].append("表格只有标题行，无数据")

        # 自动形状 (矩形等装饰)
        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            rect_count += 1

    info['title'] = title_text
    info['text_length'] = len(all_text.strip())
    info['has_image'] = image_count > 0
    info['has_chart'] = chart_count > 0
    info['has_table'] = table_count > 0
    info['image_count'] = image_count

    # ============================================================
    # 质量检查规则
    # ============================================================

    # 1. 空白页检测 (只有装饰矩形和短标题)
    # 但 cover(P1), section_break, end 页面内容少是设计如此，不算空白
    content_shapes = textbox_count + image_count + chart_count + table_count
    meaningful_text = len(all_text.strip()) > 20
    text_stripped = all_text.strip()

    # 检测结构性页面 (不算空白)
    is_structural = False
    structural_patterns = ["PART ", "谢谢聆听", "感谢聆听", "Thank You", "Q&A",
                           "Questions", "交流讨论"]
    if any(pat in text_stripped for pat in structural_patterns):
        is_structural = True
    if snum == 1:  # 封面
        is_structural = True

    if not meaningful_text and image_count == 0 and chart_count == 0 and table_count == 0:
        if is_structural:
            # 结构性页面，不算严重问题
            pass
        else:
            issues.append(('CRITICAL', snum, f'空白页 - 无实质内容 (shapes={len(shapes)}, text_len={len(text_stripped)})'))
            info['issues'].append('空白页')
            stats['blank_or_near_blank'] += 1
    elif len(shapes) <= 3 and not info['has_image'] and not info['has_chart'] and not info['has_table'] and len(text_stripped) < 50:
        if not is_structural:
            warnings.append(('WARN', snum, f'内容稀薄 - 仅{len(shapes)}个形状, {len(text_stripped)}字'))
            stats['blank_or_near_blank'] += 1

    # 2. 图片完整性
    if image_count > 0:
        stats['with_images'] += 1
        stats['image_count'] += image_count

    if chart_count > 0:
        stats['with_charts'] += 1
        stats['chart_count'] += chart_count

    if table_count > 0:
        stats['with_tables'] += 1

    if meaningful_text:
        stats['with_text'] += 1

    # 3. 重复封面检测
    title_lower = title_text.lower().strip()

    # 4. 文字溢出检测 (粗略: 单页文字过多)
    if len(all_text.strip()) > 800:
        warnings.append(('WARN', snum, f'文字可能过多 ({len(all_text.strip())}字) - 检查是否溢出'))

    for iss in info['issues']:
        issues.append(('ERROR', snum, iss))

    slide_info.append(info)

# 5. 重复标题检测
titles = [(s['num'], s['title']) for s in slide_info if s['title']]
seen_titles = {}
for num, title in titles:
    # 只检查完全相同的标题
    if title in seen_titles:
        issues.append(('WARN', num, f'标题重复 - "{title}" (与第{seen_titles[title]}页相同)'))
    else:
        seen_titles[title] = num

# 6. 连续结构页检测 (不应有连续的 section_break 或类似)
# (通过检查连续的低内容页)
for i in range(len(slide_info) - 1):
    s1 = slide_info[i]
    s2 = slide_info[i+1]
    if s1['text_length'] < 30 and s2['text_length'] < 30 and not s1['has_image'] and not s2['has_image']:
        warnings.append(('WARN', s1['num'], f'连续稀薄页 (P{s1["num"]}-P{s2["num"]})'))

# ============================================================
# 报告输出
# ============================================================
print("\n" + "=" * 80)
print("[STATS] 总体统计")
print("=" * 80)
print(f"  总页数:        {stats['total']}")
print(f"  有图片的页:    {stats['with_images']} ({stats['with_images']*100//stats['total']}%)")
print(f"  有图表的页:    {stats['with_charts']}")
print(f"  有表格的页:    {stats['with_tables']}")
print(f"  有文字内容的:  {stats['with_text']} ({stats['with_text']*100//stats['total']}%)")
print(f"  空白/稀薄页:   {stats['blank_or_near_blank']}")
print(f"  图片总数:      {stats['image_count']}")
print(f"  图表总数:      {stats['chart_count']}")

# 详细页面清单
print("\n" + "=" * 80)
print("[DETAIL] 逐页详情")
print("=" * 80)
print(f"{'页':>3} {'形状':>4} {'图':>2} {'表':>2} {'图表':>2} {'文字':>5}  标题")
print("-" * 80)
for s in slide_info:
    img_mark = f"{s['image_count']}" if s['has_image'] else '-'
    tbl_mark = 'Y' if s['has_table'] else '-'
    cht_mark = 'Y' if s['has_chart'] else '-'
    flag = ' !!' if s['issues'] else ''
    print(f"P{s['num']:>2} {s['shape_count']:>4} {img_mark:>3} {tbl_mark:>3} {cht_mark:>4} {s['text_length']:>5}  {s['title'][:50]}{flag}")

# 问题汇总
if issues:
    print("\n" + "=" * 80)
    print(f"[CRITICAL] 严重问题 ({len(issues)})")
    print("=" * 80)
    for severity, snum, desc in sorted(issues, key=lambda x: (0 if x[0]=='CRITICAL' else 1, x[1])):
        print(f"  [{severity}] P{snum}: {desc}")

if warnings:
    print("\n" + "=" * 80)
    print(f"⚠ 警告 ({len(warnings)})")
    print("=" * 80)
    for severity, snum, desc in sorted(warnings, key=lambda x: x[1]):
        print(f"  [{severity}] P{snum}: {desc}")

# 最终评分
print("\n" + "=" * 80)
print("[SCORE] 质量评分")
print("=" * 80)
critical_count = sum(1 for s, _, _ in issues if s == 'CRITICAL')
error_count = sum(1 for s, _, _ in issues if s == 'ERROR')
warn_count = len(warnings)

score = 100
score -= critical_count * 15   # 空白页严重扣分
score -= error_count * 5       # 错误扣5分
score -= warn_count * 2        # 警告扣2分

# 图片覆盖率奖惩
img_ratio = stats['with_images'] / stats['total'] if stats['total'] > 0 else 0
if img_ratio < 0.3:
    score -= 15
    print(f"  图片覆盖率过低: {img_ratio*100:.0f}% (目标>50%) → -15分")
elif img_ratio < 0.5:
    score -= 8
    print(f"  图片覆盖率偏低: {img_ratio*100:.0f}% (目标>50%) → -8分")
else:
    print(f"  图片覆盖率: {img_ratio*100:.0f}% → OK")

# 空白页
if stats['blank_or_near_blank'] > 0:
    print(f"  空白/稀薄页: {stats['blank_or_near_blank']} → -{stats['blank_or_near_blank']*10}分")

score = max(0, min(100, score))
print(f"\n  最终得分: {score}/100", end='')
if score >= 80:
    print(" [PASS] 良好")
elif score >= 60:
    print(" !! 合格但需改进")
else:
    print(" [CRITICAL] 需要修复")
