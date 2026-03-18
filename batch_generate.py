#!/usr/bin/env python3
"""
分批生成 PPT：将长文档拆分为多份，每份生成 12-15 页，最后合并
"""
import sys
sys.path.insert(0, '.')

import json
import shutil
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt

# 配置
SOURCE_PDF = "source-doc/日本军工企业深度研究 (2).pdf"
TARGET_TOTAL_PAGES = 50
BATCH_SIZE = 12  # 每批生成页数，略低于15以留余量
NUM_BATCHES = (TARGET_TOTAL_PAGES + BATCH_SIZE - 1) // BATCH_SIZE  # 向上取整

print("="*70)
print("分批生成 PPT 方案")
print("="*70)
print(f"目标总页数: {TARGET_TOTAL_PAGES}")
print(f"每批页数: {BATCH_SIZE}")
print(f"分批数: {NUM_BATCHES}")
print(f"源文件: {SOURCE_PDF}")
print("="*70)

# Step 0: 先完整提取内容
print("\n[Step 0] 提取原始文档内容...")
from src.step0_ingest import run_ingest
result = run_ingest(SOURCE_PDF, base="japan_defense_batch")
base = result['base']
output_dir = result['output_dir']

print(f"[Step 0] 完成: {output_dir}")
print(f"  - 字符数: {len((output_dir / 'raw_content.md').read_text(encoding='utf-8'))}")

# 读取原始内容
raw_content = (output_dir / 'raw_content.md').read_text(encoding='utf-8')
raw_meta = json.loads((output_dir / 'raw_meta.json').read_text(encoding='utf-8'))
raw_tables = json.loads((output_dir / 'raw_tables.json').read_text(encoding='utf-8'))

# Step 1: 完整分析（只做一次）
print("\n[Step 1] 结构化分析...")
from src.step1_analyze import run_analyze
result1 = run_analyze(base, output_dir)
analysis = result1['analysis']

print(f"[Step 1] 完成: {len(analysis.get('chapters', []))} 章节")

# Step 2: 按章节拆分生成多批大纲
print("\n[Step 2] 生成多批大纲...")

chapters = analysis.get('chapters', [])
print(f"  总章节数: {len(chapters)}")

# 将章节分配给各批次
batches = []
chapters_per_batch = len(chapters) // NUM_BATCHES
remainder = len(chapters) % NUM_BATCHES

start_idx = 0
for i in range(NUM_BATCHES):
    # 分配章节（前面的批次多分一个余数）
    count = chapters_per_batch + (1 if i < remainder else 0)
    end_idx = min(start_idx + count, len(chapters))
    
    batch_chapters = chapters[start_idx:end_idx]
    batches.append({
        'batch_id': i + 1,
        'start_chapter': start_idx,
        'end_chapter': end_idx,
        'chapters': batch_chapters,
        'target_pages': BATCH_SIZE
    })
    
    print(f"  批次 {i+1}: 章节 {start_idx+1}-{end_idx} ({len(batch_chapters)} 章) → {BATCH_SIZE} 页")
    start_idx = end_idx

# 为每批生成大纲
all_slide_plans = []

for batch in batches:
    batch_id = batch['batch_id']
    print(f"\n[Step 2-Batch {batch_id}] 生成大纲...")
    
    # 创建批次特定的分析结果
    batch_analysis = {
        **analysis,
        'chapters': batch['chapters'],
        'batch_info': {
            'batch_id': batch_id,
            'total_batches': NUM_BATCHES,
            'target_pages': batch['target_pages']
        }
    }
    
    # 保存批次分析
    batch_analysis_path = output_dir / f"analysis_batch_{batch_id}.json"
    batch_analysis_path.write_text(
        json.dumps(batch_analysis, ensure_ascii=False, indent=2),
        encoding='utf-8'
    )
    
    # 生成大纲
    from src.step2_outline import run_outline
    
    # 临时替换分析文件
    original_analysis = output_dir / "analysis.json"
    backup_analysis = output_dir / "analysis_backup.json"
    shutil.copy(original_analysis, backup_analysis)
    shutil.copy(batch_analysis_path, original_analysis)
    
    try:
        result2 = run_outline(
            base=base,
            output_dir=output_dir,
            style_profile=None,
            slide_range=(batch['target_pages'] - 2, batch['target_pages'] + 2),  # 允许±2页浮动
            two_phase=False  # 单次模式更快
        )
        
        slide_plan = result2['slide_plan']
        batch['slide_plan'] = slide_plan
        batch['actual_pages'] = len(slide_plan.get('slides', []))
        all_slide_plans.append(slide_plan)
        
        print(f"  [Batch {batch_id}] 生成完成: {batch['actual_pages']} 页")
        
    finally:
        # 恢复原分析文件
        shutil.copy(backup_analysis, original_analysis)
        backup_analysis.unlink()

print(f"\n[Step 2] 所有批次大纲生成完成")
print(f"  总页数: {sum(b['actual_pages'] for b in batches)}")

# Step 3: 为每批生成视觉资产
print("\n[Step 3] 生成视觉资产...")
from src.step3_visuals import run_visuals

for batch in batches:
    batch_id = batch['batch_id']
    print(f"\n[Step 3-Batch {batch_id}] 生成视觉资产...")
    
    # 临时替换 slide_plan
    original_plan = output_dir / "slide_plan.json"
    backup_plan = output_dir / "slide_plan_backup.json"
    
    if original_plan.exists():
        shutil.copy(original_plan, backup_plan)
    
    # 写入批次大纲
    original_plan.write_text(
        json.dumps(batch['slide_plan'], ensure_ascii=False, indent=2),
        encoding='utf-8'
    )
    
    try:
        result3 = run_visuals(
            base=base,
            output_dir=output_dir,
            no_ai_images=True,
            max_workers=4
        )
        
        # 复制视觉资产到批次目录
        assets_dir = output_dir / "assets"
        batch_assets_dir = output_dir / f"assets_batch_{batch_id}"
        if assets_dir.exists():
            shutil.copytree(assets_dir, batch_assets_dir, dirs_exist_ok=True)
            print(f"  [Batch {batch_id}] 视觉资产: {len(list(batch_assets_dir.glob('*.png')))} 个")
        
    finally:
        # 恢复
        if backup_plan.exists():
            shutil.copy(backup_plan, original_plan)
            backup_plan.unlink()

# Step 4: 为每批生成 PPT
print("\n[Step 4] 生成各批次 PPT...")
from src.step4_build import run_build

batch_pptx_files = []

for batch in batches:
    batch_id = batch['batch_id']
    print(f"\n[Step 4-Batch {batch_id}] 生成 PPT...")
    
    # 临时替换 slide_plan 和 assets
    original_plan = output_dir / "slide_plan.json"
    backup_plan = output_dir / "slide_plan_backup.json"
    
    if original_plan.exists():
        shutil.copy(original_plan, backup_plan)
    
    original_plan.write_text(
        json.dumps(batch['slide_plan'], ensure_ascii=False, indent=2),
        encoding='utf-8'
    )
    
    # 恢复对应批次的视觉资产
    assets_dir = output_dir / "assets"
    batch_assets_dir = output_dir / f"assets_batch_{batch_id}"
    
    if batch_assets_dir.exists():
        if assets_dir.exists():
            shutil.rmtree(assets_dir)
        shutil.copytree(batch_assets_dir, assets_dir)
    
    try:
        # 生成 PPT
        result4 = run_build(
            base=base,
            output_dir=output_dir,
            theme="tech"
        )
        
        # 重命名批次 PPT
        original_pptx = result4['pptx_path']
        batch_pptx = output_dir / f"{base}_batch_{batch_id}.pptx"
        shutil.copy(original_pptx, batch_pptx)
        batch_pptx_files.append(batch_pptx)
        
        print(f"  [Batch {batch_id}] PPT: {batch_pptx.name} ({batch['actual_pages']} 页)")
        
    finally:
        # 恢复
        if backup_plan.exists():
            shutil.copy(backup_plan, original_plan)
            backup_plan.unlink()

# Step 5: 合并所有 PPT
print("\n[Step 5] 合并所有批次 PPT...")

def merge_pptx_files(pptx_files, output_path):
    """合并多个 PPTX 文件"""
    merged = Presentation()
    merged.slide_width = Inches(13.333)
    merged.slide_height = Inches(7.5)
    
    total_slides = 0
    
    for i, pptx_file in enumerate(pptx_files, 1):
        print(f"  合并批次 {i}/{len(pptx_files)}: {pptx_file.name}")
        prs = Presentation(str(pptx_file))
        
        # 添加批次分隔页（除第一批外）
        if i > 1:
            blank_layout = merged.slide_layouts[6]
            separator = merged.slides.add_slide(blank_layout)
            
            # 添加批次标题
            from pptx.util import Pt
            from pptx.dml.color import RGBColor
            
            title_box = separator.shapes.add_textbox(
                Inches(1), Inches(3), Inches(11), Inches(1.5)
            )
            tf = title_box.text_frame
            p = tf.paragraphs[0]
            p.text = f"Part {i}"
            p.font.size = Pt(44)
            p.font.bold = True
            p.font.color.rgb = RGBColor(0, 32, 96)
            
            total_slides += 1
        
        # 复制所有幻灯片
        for slide in prs.slides:
            # 获取空白布局
            blank_layout = merged.slide_layouts[6]
            new_slide = merged.slides.add_slide(blank_layout)
            
            # 复制形状（简化版）
            for shape in slide.shapes:
                if shape.has_text_frame:
                    # 复制文本框
                    new_shape = new_slide.shapes.add_textbox(
                        shape.left, shape.top, shape.width, shape.height
                    )
                    new_shape.text_frame.text = shape.text_frame.text
                    
                    # 复制字体样式
                    for para_idx, para in enumerate(shape.text_frame.paragraphs):
                        if para_idx < len(new_shape.text_frame.paragraphs):
                            new_para = new_shape.text_frame.paragraphs[para_idx]
                            for run_idx, run in enumerate(para.runs):
                                if run_idx < len(new_para.runs):
                                    new_run = new_para.runs[run_idx]
                                    new_run.font.size = run.font.size
                                    new_run.font.bold = run.font.bold
                                    new_run.font.color.rgb = run.font.color.rgb
            
            total_slides += 1
    
    # 保存
    merged.save(str(output_path))
    return total_slides

final_pptx = output_dir / f"{base}_merged_{TARGET_TOTAL_PAGES}pages.pptx"
total = merge_pptx_files(batch_pptx_files, final_pptx)

print(f"\n[Step 5] 合并完成!")
print(f"  最终文件: {final_pptx}")
print(f"  总页数: {total}")

# 生成报告
print("\n" + "="*70)
print("分批生成报告")
print("="*70)
for batch in batches:
    print(f"批次 {batch['batch_id']}: {batch['actual_pages']} 页 (章节 {batch['start_chapter']+1}-{batch['end_chapter']})")
print("-"*70)
print(f"总计: {total} 页")
print(f"目标: {TARGET_TOTAL_PAGES} 页")
print(f"差异: {total - TARGET_TOTAL_PAGES} 页")
print("="*70)
