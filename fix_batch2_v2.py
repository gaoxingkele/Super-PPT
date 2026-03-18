#!/usr/bin/env python3
"""重新运行批次2的完整流程 (Step2+Step3+Step4)，然后合并所有批次"""
import sys
sys.path.insert(0, '.')

import json
import shutil
import copy
from pathlib import Path

OUTPUT_DIR = Path("output/japan_defense_50pages")
BATCH2_DIR = OUTPUT_DIR / "batch_2"

# 读取分析结果
analysis = json.loads((OUTPUT_DIR / "analysis.json").read_text(encoding='utf-8'))
chapters = analysis.get('chapters', [])

# 批次2的章节: ch03, ch04
target_chapters = [ch for ch in chapters if ch.get('id') in ('ch03', 'ch04')]
print(f"批次2 章节: {[ch['id'] + ' ' + ch['title'] for ch in target_chapters]}")

# 重写 batch_2 的 analysis.json
batch_analysis = {**analysis, 'chapters': target_chapters}
(BATCH2_DIR / "analysis.json").write_text(
    json.dumps(batch_analysis, ensure_ascii=False, indent=2),
    encoding='utf-8'
)

# 删除旧的 slide_plan 以强制重新生成
old_plan = BATCH2_DIR / "slide_plan.json"
if old_plan.exists():
    old_plan.unlink()
    print("已删除旧的 slide_plan.json")

# 重新运行 Step2 + Step3 + Step4
from src.step2_outline import run_outline
from src.step3_visuals import run_visuals
from src.step4_build import run_build

batch_base = "japan_defense_50pages_batch2"

print("\n" + "="*70)
print("[批次 2] Step 2: 重新生成大纲...")
print("="*70)
result2 = run_outline(
    base=batch_base,
    output_dir=BATCH2_DIR,
    slide_range=(12, 14),
    two_phase=False
)
pages = len(result2['slide_plan'].get('slides', []))
print(f"  生成: {pages} 页")

if pages == 0:
    print("错误: 大纲生成仍然为0页!")
    sys.exit(1)

print(f"\n[批次 2] Step 3: 视觉资产...")
result3 = run_visuals(batch_base, BATCH2_DIR, no_ai_images=True)
print(f"  完成: {len(result3.get('manifest', {}))} 个")

print(f"\n[批次 2] Step 4: 生成PPT...")
result4 = run_build(batch_base, BATCH2_DIR, theme="tech")

# 复制到汇总目录
src_pptx = result4['pptx_path']
dst_pptx = OUTPUT_DIR / "batch_2_电子与潜艇.pptx"
shutil.copy(src_pptx, dst_pptx)
print(f"  批次2完成: {pages} 页 -> {dst_pptx.name}")

# ============================================================
# 合并所有批次
# ============================================================
print("\n" + "="*70)
print("合并所有批次")
print("="*70)

from pptx import Presentation

batch_files = sorted(OUTPUT_DIR.glob("batch_*_*.pptx"))
print(f"找到 {len(batch_files)} 个批次文件:")
total = 0
for f in batch_files:
    prs = Presentation(str(f))
    n = len(prs.slides)
    total += n
    print(f"  {f.name}: {n} 页")
print(f"预计总页数: {total}")

# 合并: 以第一个文件为基础，追加其他文件的slides
merged = Presentation(str(batch_files[0]))
print(f"\n基础: {batch_files[0].name} ({len(merged.slides)} 页)")

for bf in batch_files[1:]:
    src_prs = Presentation(str(bf))
    if len(src_prs.slides) == 0:
        print(f"  跳过空文件: {bf.name}")
        continue

    for slide in src_prs.slides:
        blank_layout = merged.slide_layouts[6]  # blank
        new_slide = merged.slides.add_slide(blank_layout)

        for shape in slide.shapes:
            el = copy.deepcopy(shape._element)
            new_slide.shapes._spTree.append(el)

        # 复制背景
        try:
            if slide.background.fill.type is not None:
                bg_el = copy.deepcopy(slide.background._element)
                new_slide.background._element.getparent().replace(
                    new_slide.background._element, bg_el
                )
        except:
            pass

    print(f"  合并: {bf.name} ({len(src_prs.slides)} 页)")

merged_path = OUTPUT_DIR / "日本军工企业深度研究_完整版.pptx"
merged.save(str(merged_path))
final_pages = len(merged.slides)
print(f"\n{'='*70}")
print(f"最终结果: {merged_path.name} ({final_pages} 页)")
print(f"路径: {merged_path}")
print(f"{'='*70}")
